import copy
import datetime
from functools import wraps
from glob import glob
import math
from matplotlib.figure import Figure
from matplotlib.backends.backend_tkagg import (
    FigureCanvasTkAgg, NavigationToolbar2Tk)
import os
import pandas as pd
import pickle
import pprint
import re
import shutil
import threading
import tkinter as tk
from tkinter import filedialog
from tkinter import messagebox
from tkinter import scrolledtext
from tkinter import ttk
import traceback
import statistics
import time


import openpyxl
from openpyxl.styles import PatternFill, Font
from openpyxl.utils.dataframe import dataframe_to_rows
import pptx
from pptx.util import Inches


from Data_Preprocessor import DataPreprocessor
import ReAnalyzer
from Spectrum_analyzer import SingleAnalyzer
from Spectrum_analyzer import BatchAnalyzer

IconPath = 'image/OAD_Analyzer_icon_ver2.png'

class MSRIDD_GUI(tk.Frame):
    def __init__(self, master):
        super().__init__(master)
        self.master.attributes('-topmost', False)
        self.main_window_geometry(self.master)
        self.master.title("MS-RIDD")
        self.master.resizable(width=True, height=True)
        self.master.minsize(width=1200, height=600)
        # self.iconfile = './favicon.ico'
        # self.master.iconbitmap(default=self.iconfile)
        self.master.iconphoto(False, tk.PhotoImage(file=IconPath))
        self.create_menu()
        style = ttk.Style()  
        style.map('Treeview', foreground=self.fixed_map('foreground', style), 
            background=self.fixed_map('background', style))
        self.create_main_widgets()
        self.disable_widgets()

    def fixed_map(self, option, style):
        """ Fix for setting text colour for Tkinter 8.6.9
            From: https://core.tcl.tk/tk/info/509cafafae
            Returns the style map for 'option' with any styles starting with
            ('!disabled', '!selected', ...) filtered out.
            style.map() returns an empty list for missing options, so this
            should be future-safe. """
        return [elm for elm in style.map('Treeview', query_opt=option) 
            if elm[:2] != ('!disabled', '!selected')]

    def create_menu(self):
        self.menubar = tk.Menu(self.master)
        self.filemenu = tk.Menu(self.menubar, tearoff=0)
        self.newprjmenu = tk.Menu(self.filemenu, tearoff=0)
        # self.filemenu.add_command(
        #     label='New project', command=self.create_setting_window)
        self.newprjmenu.add_command(
            label='Batch analysis', command=self.create_batch_setwin)
        self.newprjmenu.add_command(
            label='Single analysis', command=self.create_setting_window)
        self.filemenu.add_cascade(label='New project', menu=self.newprjmenu)
        self.filemenu.add_command(
            label='Open project', command=self.import_project_files)
        self.filemenu.add_command(
            label='Save', command=self.start_saving)
        self.menubar.add_cascade(label='File', menu=self.filemenu)
        self.datamenu = tk.Menu(self.menubar, tearoff=0)
        self.datamenu.add_command(
            label='Merge Neg&Pos CID-MS/MS data', 
            command=self.merge_bipolarity_cid)
        self.datamenu.add_command(
            label='Merge CID&OAD MS/MS data', command=self.merge_cid_and_oad)
        self.menubar.add_cascade(label='Data preparation', menu=self.datamenu)
        # self.datamenu.entryconfigure(
        #     'Merge Neg&Pos CID-MS/MS data', state='disabled')
        self.exportmenu = tk.Menu(self.menubar, tearoff=0)
        self.exportmenu.add_command(
            label='Result export', command=self.create_data_export_window)
        self.menubar.add_cascade(label='Export', menu=self.exportmenu)
        self.master.config(menu=self.menubar)

    def main_window_geometry(self, win):
        width = int(win.winfo_screenwidth()*0.95)
        height = int(win.winfo_screenheight()*0.9)
        x = (win.winfo_screenwidth() // 2) - (width // 2)
        win.geometry('{}x{}+{}+{}'.format(width, height, x, 0))

    def center_position(self, win, width, height):
        x = (win.winfo_screenwidth() // 2) - (width // 2)
        y = (win.winfo_screenheight() // 2) - (height // 2)
        win.geometry('{}x{}+{}+{}'.format(width, height, x, y))

    def close_window(self, win):
        win.destroy()

    #region Main widgets
    def create_main_widgets(self):
        # Result table
        self.left_wrap_lebel = tk.StringVar()
        self.left_wrap_lebel.set("Result table")
        self.left_wrap_frame = ttk.LabelFrame(
            self.master, text=self.left_wrap_lebel.get())
        self.result_table_pwin = ttk.PanedWindow(
            self.left_wrap_frame, orient="vertical") 
        self.result_table_pwin.place(
            relx=0, rely=0, relwidth=0.97, relheight=0.98)
        self.left_wrap_frame.place(
            relx=0.005, rely=0.005, relwidth=0.36, relheight=0.44)
        # Show only solved molecules CheckBox
        self.show_only_solved_mol = tk.BooleanVar()
        self.show_only_solved_mol.set(False)
        self.show_only_solved_molcheckbtn = ttk.Checkbutton(self.master, 
            text="Show only C=C positions solved molecules", 
            variable=self.show_only_solved_mol, 
            command=self.switch_result_table_tree)
        self.show_only_solved_molcheckbtn.place(relx=0.235, rely=0.445)
        # Comment input Button
        self.comment_btn = ttk.Button(self.master, text="Input comment", 
            command=self.enter_comment_on_result_table)
        self.comment_btn.place(
            relx=0.285, rely=0.47, relwidth=0.08, relheight=0.025)
        self.comment_str = tk.StringVar()
        self.comment_entry = ttk.Entry(
            self.master, textvariable=self.comment_str)
        self.comment_entry.place(
            relx=0.005, rely=0.47, relwidth=0.275, relheight=0.025)
        # C=C candidates table
        self.db_candidates_frame = ttk.LabelFrame(self.master, 
            text="C=C positional candidates")
        self.db_1_frame = tk.Frame(self.db_candidates_frame)
        self.db_2_frame = tk.Frame(self.db_candidates_frame)
        self.db_3_frame = tk.Frame(self.db_candidates_frame)
        self.db_1_pwin = ttk.PanedWindow(self.db_1_frame, orient="vertical") 
        self.db_2_pwin = ttk.PanedWindow(self.db_2_frame, orient="vertical") 
        self.db_3_pwin = ttk.PanedWindow(self.db_3_frame, orient="vertical") 
        self.db_1_pwin.place(relx=0, rely=0, relwidth=0.96, relheight=1)
        self.db_2_pwin.place(relx=0, rely=0, relwidth=0.96, relheight=1)
        self.db_3_pwin.place(relx=0, rely=0, relwidth=0.96, relheight=1)
        self.db_1_frame.place(relx=0, rely=0, relwidth=0.33, relheight=1)
        self.db_2_frame.place(relx=0.34, rely=0, relwidth=0.33, relheight=1)
        self.db_3_frame.place(relx=0.67, rely=0, relwidth=0.33, relheight=1)
        self.db_candidates_frame.place(
            relx=0.37, rely=0.005, relwidth=0.625, relheight=0.12)
        # Re-analysis Button
        self.re_analysis_btn = ttk.Button(self.master, 
            text='Re-analize C=C position', 
            command=self.set_re_analysis_condition)
        self.re_analysis_btn.place(
            relx=0.8, rely=0.13, relwidth=0.08, relheight=0.025)
        # Modifying C=C Button
        self.modifying_db_btn = ttk.Button(
            self.master, text="Modify C=C position", command=self.modify_db_pos)
        self.modifying_db_btn.place(
            relx=0.9, rely=0.13, relwidth=0.08, relheight=0.025)
        # Show only matched fragment ions CheckBox
        self.show_only_matched_ions = tk.BooleanVar()
        self.show_only_matched_ions.set(False)
        self.show_only_matched_ions_checkbtn = ttk.Checkbutton(self.master, 
            text="Show only matched fragment ions", 
            variable=self.show_only_matched_ions, 
            command=self.switch_matched_fragment_ions_tree)
        self.show_only_matched_ions_checkbtn.place(relx=0.37, rely=0.138)
        # Matched OAD fragment ions
        self.matched_ions_label = tk.StringVar()
        self.matched_ions_label.set("Mass table")
        self.matched_fragments_frame = ttk.LabelFrame(
            self.master, text=self.matched_ions_label.get())
        self.matched_fragments_pwin = ttk.PanedWindow(
            self.matched_fragments_frame, orient="vertical") 
        self.matched_fragments_pwin.place(
            relx=0, rely=0, relwidth=0.99, relheight=1)
        self.matched_fragments_frame.place(
            relx=0.37, rely=0.155, relwidth=0.625, relheight=0.34)
        # Centroid OAD-MS/MS
        self.centroid_msms_frame = ttk.LabelFrame(
            self.master, text="Centroid OAD-MS/MS")
        self.centroid_msms_frame.place(
            relx=0.005, rely=0.515, relwidth=0.36, relheight=0.48)
        # Measured vs Reference OAD-MS/MS
        self.measured_vs_ref_msms_frame = ttk.LabelFrame(
            self.master, text="Measurement vs Reference OAD-MS/MS")
        self.measured_vs_ref_msms_frame.place(
            relx=0.37, rely=0.515, relwidth=0.625, relheight=0.48)

    def recreate_main_widgets(self):
        # Result table
        self.result_table_pwin = ttk.PanedWindow(self.left_wrap_frame, orient="vertical") 
        self.result_table_pwin.place(relx=0, rely=0, relwidth=0.97, relheight=0.98)
        # C=C candidates table
        self.db_1_frame = tk.Frame(self.db_candidates_frame)
        self.db_2_frame = tk.Frame(self.db_candidates_frame)
        self.db_3_frame = tk.Frame(self.db_candidates_frame)
        self.db_1_pwin = ttk.PanedWindow(self.db_1_frame, orient="vertical") 
        self.db_2_pwin = ttk.PanedWindow(self.db_2_frame, orient="vertical") 
        self.db_3_pwin = ttk.PanedWindow(self.db_3_frame, orient="vertical") 
        self.db_1_pwin.place(relx=0, rely=0, relwidth=0.96, relheight=1)
        self.db_2_pwin.place(relx=0, rely=0, relwidth=0.96, relheight=1)
        self.db_3_pwin.place(relx=0, rely=0, relwidth=0.96, relheight=1)
        self.db_1_frame.place(relx=0, rely=0, relwidth=0.33, relheight=1)
        self.db_2_frame.place(relx=0.34, rely=0, relwidth=0.33, relheight=1)
        self.db_3_frame.place(relx=0.67, rely=0, relwidth=0.33, relheight=1)
        # Matched OAD fragment ions
        self.matched_fragments_pwin = ttk.PanedWindow(self.matched_fragments_frame, orient="vertical") 
        self.matched_fragments_pwin.place(relx=0, rely=0, relwidth=0.99, relheight=1)

    def create_result_table_tree(self, parent):
        self.result_table_tree = ttk.Treeview(parent)
        win_w = parent.winfo_width()
        w_unit = int(win_w/33)
        ion_v_col = 'pmol/mg tissue' if self.normalized_data else 'Height'
        self.result_table_tree_column = [
            'ID', 'RT', 'm/z', 'm/z type', 'Solved','Result name', 
            'Data from', ion_v_col, 'Ontology', 'Comment']
        self.result_table_tree_column_width = [
            w_unit*2, w_unit*2, w_unit*3, w_unit*3, w_unit*3, w_unit*15, 
            w_unit*3, w_unit*3, w_unit*3, w_unit*8]
        self.result_table_tree["column"] = self.result_table_tree_column
        self.result_table_tree["show"] = "headings"
        for i, (colname, width) in enumerate(zip(self.result_table_tree_column, self.result_table_tree_column_width)):
            self.result_table_tree.heading(i, text=colname)
            self.result_table_tree.column(i, width=width)
        h_scrollbar = tk.Scrollbar(self.left_wrap_frame, orient=tk.HORIZONTAL,
            command=self.result_table_tree.xview)
        v_scrollbar = tk.Scrollbar(self.left_wrap_frame, orient=tk.VERTICAL,
            command=self.result_table_tree.yview)
        self.result_table_tree.configure(xscrollcommand=h_scrollbar.set, 
            yscrollcommand=v_scrollbar.set)
        v_scrollbar.pack(side=tk.RIGHT, fill='y')
        h_scrollbar.pack(side=tk.BOTTOM, fill='x')
        self.result_table_tree.focus_set()

    def create_db_1_table_tree(self, parent):
        self.db_1_tree = ttk.Treeview(parent)
        win_w = parent.winfo_width()
        w_unit = int(win_w/6)
        self.db_1_tree_column = ['Rank', 'C=C moiety-1', 'Score', 'Presence', 
                                 'Rel. Int']
        self.db_1_tree_column_width = [w_unit, w_unit*2, w_unit, w_unit, w_unit]
        self.db_1_tree["column"] = self.db_1_tree_column
        self.db_1_tree["show"] = "headings"
        for i, (colname, width) in enumerate(zip(
            self.db_1_tree_column, self.db_1_tree_column_width)):
            self.db_1_tree.heading(i, text=colname)
            self.db_1_tree.column(i, width=width)
        v_scrollbar = tk.Scrollbar(self.db_1_frame, orient=tk.VERTICAL,
            command=self.db_1_tree.yview)
        self.db_1_tree.configure(yscrollcommand=v_scrollbar.set)
        v_scrollbar.pack(side=tk.RIGHT, fill='y')
        self.db_1_tree.focus_set()

    def create_db_2_table_tree(self, parent):
        self.db_2_tree = ttk.Treeview(parent)
        win_w = parent.winfo_width()
        w_unit = int(win_w/6)
        self.db_2_tree_column = ['Rank', 'C=C moiety-2', 'Score', 'Presence', 
                                 'Rel. Int']
        self.db_2_tree_column_width = [w_unit, w_unit*2, w_unit, w_unit, w_unit]
        self.db_2_tree["column"] = self.db_2_tree_column
        self.db_2_tree["show"] = "headings"
        for i, (colname, width) in enumerate(zip(
            self.db_2_tree_column, self.db_2_tree_column_width)):
            self.db_2_tree.heading(i, text=colname)
            self.db_2_tree.column(i, width=width)
        v_scrollbar = tk.Scrollbar(self.db_2_frame, orient=tk.VERTICAL,
            command=self.db_2_tree.yview)
        self.db_2_tree.configure(yscrollcommand=v_scrollbar.set)
        v_scrollbar.pack(side=tk.RIGHT, fill='y')
        self.db_2_tree.focus_set()

    def create_db_3_table_tree(self, parent):
        self.db_3_tree = ttk.Treeview(parent)
        win_w = parent.winfo_width()
        w_unit = int(win_w/6)
        self.db_3_tree_column = ['Rank', 'C=C moiety-3', 'Score', 'Presence', 
                                 'Rel. Int']
        self.db_3_tree_column_width = [w_unit, w_unit*2, w_unit, w_unit, w_unit]
        self.db_3_tree["column"] = self.db_3_tree_column
        self.db_3_tree["show"] = "headings"
        for i, (colname, width) in enumerate(zip(
            self.db_3_tree_column, self.db_3_tree_column_width)):
            self.db_3_tree.heading(i, text=colname)
            self.db_3_tree.column(i, width=width)
        v_scrollbar = tk.Scrollbar(self.db_3_frame, orient=tk.VERTICAL,
            command=self.db_3_tree.yview)
        self.db_3_tree.configure(yscrollcommand=v_scrollbar.set)
        v_scrollbar.pack(side=tk.RIGHT, fill='y')
        self.db_3_tree.focus_set()

    def create_matched_fragments_table_tree(self, parent):
        self.matched_fragments_tree = ttk.Treeview(parent)
        self.matched_fragments_tree_column = ['m/z', 'Intensity', 'Ratio(%)', 'Delta', 
                                              'Ref m/z1', 'ppm1', 'Type1',
                                              'Ref m/z2', 'ppm2', 'Type2',
                                              'Ref m/z3', 'ppm3', 'Type3']
        self.matched_fragments_tree_column_width = [10, 10, 10, 10,
                                                    10, 4, 100, 10, 4, 100, 10, 4, 100]
        self.matched_fragments_tree["column"] = self.matched_fragments_tree_column
        self.matched_fragments_tree["show"] = "headings"
        for i, (colname, width) in enumerate(zip(
            self.matched_fragments_tree_column, 
            self.matched_fragments_tree_column_width)):
            self.matched_fragments_tree.heading(i, text=colname)
            self.matched_fragments_tree.column(i, width=width)
        v_scrollbar = tk.Scrollbar(self.matched_fragments_frame, orient=tk.VERTICAL,
            command=self.matched_fragments_tree.yview)
        self.matched_fragments_tree.configure(yscrollcommand=v_scrollbar.set)
        v_scrollbar.pack(side=tk.RIGHT, fill='y')
        self.matched_fragments_tree.focus_set()
    
    def disable_widgets(self):
        filenmenu = ['Save']
        for menu in filenmenu:
            self.filemenu.entryconfigure(menu, state='disabled')
        exportmenu = ['Result export']
        for menu in exportmenu:
            self.exportmenu.entryconfigure(menu, state='disabled')
        self.re_analysis_btn['state'] = tk.DISABLED
        self.modifying_db_btn['state'] = tk.DISABLED
        self.show_only_matched_ions_checkbtn['state'] = tk.DISABLED
        self.comment_entry['state'] = tk.DISABLED
        self.comment_btn['state'] = tk.DISABLED
        self.show_only_solved_molcheckbtn['state'] = tk.DISABLED

    def activate_widgets(self):
        filenmenu = ['Save']
        for menu in filenmenu:
            self.filemenu.entryconfigure(menu, state='normal')
        exportmenu = ['Result export']
        for menu in exportmenu:
            self.exportmenu.entryconfigure(menu, state='normal')
        self.re_analysis_btn['state'] = tk.NORMAL
        self.modifying_db_btn['state'] = tk.NORMAL
        self.show_only_matched_ions_checkbtn['state'] = tk.NORMAL
        self.comment_entry['state'] = tk.NORMAL
        self.comment_btn['state'] = tk.NORMAL
        self.show_only_solved_molcheckbtn['state'] = tk.NORMAL
    #endregion

    #region Batch analysis
    def create_batch_setwin(self):
        self.batch_setwin = BatchSettingWindow(master=self.master)
        self.batch_setwin.finish_btn["command"] = self.check_batch_paths

    def check_batch_paths(self):
        prj_path = self.batch_setwin.prj_path_entry.get().rsplit('/', 1)[0]
        alignment_path = self.batch_setwin.alignment_entry.get()
        peakid_path = self.batch_setwin.peakid_entry.get()
        peaklists_dict = {name: path for name, path in zip(
            self.batch_setwin.entries['Files'], 
            self.batch_setwin.entries['Paths'])}
        if prj_path == "":
            messagebox.showwarning(
                "Empty project folder", "Please select project folder")
        elif alignment_path == "":
            messagebox.showwarning(
                "Empty Alignment data", "Please select Alignment data")
        elif peakid_path == "":
            messagebox.showwarning(
                "Empty PeakID data", "Please select PeakID data")
        elif not peaklists_dict:
            messagebox.showwarning(
                "Empty PeakLists data", "Please select PeakLists data")
        else:
            self.delete_contents()
            self.recreate_main_widgets()
            self.disable_widgets()
            self.analysis_start = False
            thread_1 = threading.Thread(target=self.create_progress_bar_window)
            thread_1.start()
            thread_2 = threading.Thread(target=self.start_batch_analysis)
            thread_2.start()

    def start_batch_analysis(self):
        start = time.time()
        user_def_path = self.batch_setwin.prj_path_entry.get()
        postfix = '_analysis_table.pkl'
        self.prj_path = user_def_path.rsplit('/', 1)[0]
        self.user_prefix = user_def_path.rsplit('/', 1)[1].replace(postfix, '')
        alignment_path = self.batch_setwin.alignment_entry.get()
        peakid_path = self.batch_setwin.peakid_entry.get()
        normalized_data = self.batch_setwin.normalized_data.get()
        # merged_check = self.batch_setwin.merged_data.get()
        self. normalized_data = normalized_data
        peaklists_dict = {name.get(): path.get() for name, path in zip(
            self.batch_setwin.entries['Files'], 
            self.batch_setwin.entries['Paths'])}
        self.ms_tolerance_ppm = 15
        self.must_nl_cut_off_dict = {'diagnostic_1': ['OAD03', 0.01], 
                                     'diagnostic_2': ['OAD16', 0.01], 
                                     'sphingobase': 0.005}
        cut_off_ratio = 0
        while not self.analysis_start:
            time.sleep(0.1)
        self.batch_setwin.close_window()
        try:
            spectrum_analyzer = BatchAnalyzer(
                directry=self.prj_path, prefix=self.user_prefix,
                alignment_path=alignment_path, 
                # merged=merged_check,
                peakid_path = peakid_path, peaklists_dict=peaklists_dict, 
                ms_tolerance_ppm=self.ms_tolerance_ppm, 
                must_nl_cut_off_dict=self.must_nl_cut_off_dict, 
                cut_off_ratio=cut_off_ratio, normalized=normalized_data,
                sec_rep=self.section_report, sec_bar=self.section_bar,
                each_rep=self.each_report, each_bar=self.each_prgbar, 
                timer=self.small_timer)
            # result setting
            self.target_table = spectrum_analyzer.target_table
            self.msms_dict = spectrum_analyzer.msms_dict
            self.lipid_structural_info_dict \
            = spectrum_analyzer.lipid_structural_info_dict
            self.cid_result_dict = spectrum_analyzer.cid_result_dict
            self.oad_result_dict = spectrum_analyzer.oad_result_dict
            self.graph_dict = spectrum_analyzer.graph_dict
            # updating main window
            self.master.title(f"MS-RIDD | {user_def_path}")
            self.set_analysis_result_into_table()
            self.data_setting(event=None, is_first=True)
            self.activate_widgets()
            self.data_analysis_running = False
        except Exception as e:
            self.data_analysis_running = False
            self.section_report.set("--- Error occurred ---")
            messagebox.showerror(
                "Error message", 
                """Error has occurred while the analysis, 
                   please confirm input data.""")
            messagebox.showwarning("Error details", "{}".format(e))
            traceback.print_exc()
        required_time = time.time() - start
        print(f'Required time for all process: {required_time} [sec]')
    #endregion

    #region Single analysis
    def create_setting_window(self):
        self.setting_window = tk.Toplevel(self.master)
        self.setting_window.attributes('-topmost', False)
        self.center_position(self.setting_window, width=690, height=190)
        self.setting_window.resizable(width=False, height=False)
        self.setting_window.title("Start up a project")
        self.setting_window.grab_set()
        self.setting_window.iconphoto(False, tk.PhotoImage(file=IconPath))
        self.create_setting_widgets()
        self.setting_window.protocol(
            'WM_DELETE_WINDOW', lambda : self.close_window(self.setting_window))

    def create_setting_widgets(self):
        # Project path
        prj_entry_label = ttk.Label(self.setting_window, text="Project file path :")
        prj_entry_label.place(x=30, y=30)
        self.prj_path = tk.StringVar()
        self.prj_path_entry = ttk.Entry(self.setting_window, textvariable=self.prj_path, width=68)
        self.prj_path_entry.place(x=150, y=30)
        self.prj_path_btn = ttk.Button(self.setting_window, text="Browse", command=self.get_filepath)
        self.prj_path_btn.place(x=580, y=28)
        # Input data path
        inputdata_label = ttk.Label(self.setting_window, text="Input data path :")
        inputdata_label.place(x=30, y=80)
        self.inputdata_path = tk.StringVar()
        self.inputdata_entry = ttk.Entry(self.setting_window, textvariable=self.inputdata_path, width=68)
        self.inputdata_entry.place(x=150, y=80)
        self.inputdata_btn = ttk.Button(self.setting_window, text="Browse", command=self.get_inputdata_path)
        self.inputdata_btn.place(x=580, y=78)
        # Data format
        dataformat_label = ttk.Label(self.setting_window, text="Data format :")
        dataformat_label.place(x=30, y=130)
        self.radio_btn_value = tk.StringVar()
        self.radio_btn_value.set("")
        self.radio_btn1 = ttk.Radiobutton(self.setting_window, text="PeakList", value="PeakList", variable=self.radio_btn_value)
        self.radio_btn2 = ttk.Radiobutton(self.setting_window, text="Alignment", value="Alignment", variable=self.radio_btn_value)
        # self.radio_btn3 = ttk.Radiobutton(self.setting_window, text="Merged text", value="Merged text", variable=self.radio_btn_value)
        self.radio_btn1.place(x=150, y=130)
        self.radio_btn2.place(x=250, y=130)
        # self.radio_btn3.place(x=350, y=130)
        # Clear button
        self.clear_btn = ttk.Button(self.setting_window, text="Clear", command=self.clear_button)
        self.clear_btn.place(x=489, y=128)
        # Finish button
        self.finish_btn = ttk.Button(self.setting_window, text="Start", command=self.check_parameters)
        self.finish_btn.place(x=580, y=128) 

    def get_filepath(self):
        old_path = self.prj_path_entry.get()
        new_path = filedialog.askdirectory(title="Select project folder")
        if new_path == "":
            self.prj_path.set(old_path)
        else:
            timestamp = get_timestamp()
            default_path = f'{new_path}/{timestamp}_analysis_table.pkl'
            self.prj_path.set(default_path)

    def get_inputdata_path(self):
        old_path = self.inputdata_entry.get()
        filetype_list = [("text file", "*.txt")]
        new_path = filedialog.askopenfilename(
            filetypes=filetype_list, title="Select file")
        if new_path == "":
            self.inputdata_path.set(old_path)
        else:
            self.inputdata_path.set(new_path)

    def clear_button(self):
        self.prj_path.set("")
        self.inputdata_path.set("")
        self.radio_btn_value.set("")
        # self.sample_name.set("")

    def check_parameters(self):
        prj_path = self.prj_path_entry.get().rsplit('/', 1)[0]
        inputdata_path = self.inputdata_entry.get()
        dataformat = self.radio_btn_value.get()
        # sample_name = self.samplename_entry.get()
        if prj_path == "":
            messagebox.showwarning(
                "Empty project folder", "Please select project folder")
        elif inputdata_path == "":
            messagebox.showwarning(
                "Empty input data", "Please select input data")
        elif dataformat == "":
            messagebox.showwarning(
                "Select data format", "Please select data format")
        # elif sample_name == "":
        #     messagebox.showwarning("Empty sample name", "Please fill in sample name")
        else:
            self.delete_contents()
            self.recreate_main_widgets()
            self.disable_widgets()
            self.analysis_start = False
            thread_1 = threading.Thread(target=self.create_progress_bar_window)
            thread_1.start()
            thread_2 = threading.Thread(target=self.start_single_analysis)
            thread_2.start()

    def start_single_analysis(self):
        start = time.time()
        dataformat = self.radio_btn_value.get()
        user_def_path = self.prj_path_entry.get()
        postfix = '_analysis_table.pkl'
        self.prj_path = user_def_path.rsplit('/', 1)[0]
        self.user_prefix = user_def_path.rsplit('/', 1)[1].replace(postfix, '')
        inputdata_path = self.inputdata_entry.get()
        self. normalized_data = False
        self.ms_tolerance_ppm = 15
        self.must_nl_cut_off_dict = {'diagnostic_1': ['OAD03', 0.01], 
                                     'diagnostic_2': ['OAD16', 0.01], 
                                     'sphingobase': 0.005}
        cut_off_ratio = 0
        while not self.analysis_start: time.sleep(0.1)
        self.setting_window.destroy()
        try:
            spectrum_analyzer = SingleAnalyzer(
                tabel_format=dataformat, directry=self.prj_path, 
                prefix=self.user_prefix, input_data=inputdata_path, 
                ms_tolerance_ppm=self.ms_tolerance_ppm, 
                must_nl_cut_off_dict=self.must_nl_cut_off_dict, 
                cut_off_ratio=cut_off_ratio, file_name='',
                sec_rep=self.section_report, sec_bar=self.section_bar,
                each_rep=self.each_report, each_bar=self.each_prgbar, 
                timer=self.small_timer)
            # result setting
            self.target_table = spectrum_analyzer.target_table
            self.msms_dict = spectrum_analyzer.msms_dict
            self.lipid_structural_info_dict \
            = spectrum_analyzer.lipid_structural_info_dict
            self.cid_result_dict = spectrum_analyzer.cid_result_dict
            self.oad_result_dict = spectrum_analyzer.oad_result_dict
            self.graph_dict = spectrum_analyzer.graph_dict
            # updating main window
            self.master.title(f"MS-RIDD | {user_def_path}")
            self.set_analysis_result_into_table()
            self.data_setting(event=None, is_first=True)
            self.activate_widgets()
            self.data_analysis_running = False
        except Exception as e:
            self.data_analysis_running = False
            self.section_report.set("--- Error occurred ---")
            messagebox.showerror(
                "Error message", 
                """Error has occurred while the analysis, 
                   please confirm input data.""")
            messagebox.showwarning("Error details", "{}".format(e))
            traceback.print_exc()
        print(f'Process time: {time.time() - start:.3f} [sec]')
    
    def delete_contents(self):
        frame_list = [self.left_wrap_frame, self.db_candidates_frame, 
                      self.matched_fragments_frame, self.centroid_msms_frame, 
                      self.measured_vs_ref_msms_frame]
        for frame in frame_list:
            for widget in frame.winfo_children():
                widget.destroy()
    #endregion

    #region Open project
    def import_project_files(self):
        filetype_list = [("pkl file", "*.pkl")]
        self.imported_data = filedialog.askopenfilename(
            filetypes=filetype_list, title="Select file")
        if self.imported_data != '':
            self.delete_contents()
            self.recreate_main_widgets()
            self.disable_widgets()
            postfix = '_analysis_table.pkl'
            self.master.title(f"MS-RIDD | {self.imported_data}")
            self.prj_path = self.imported_data.rsplit('/', 1)[0]
            self.user_prefix = self.imported_data.rsplit('/', 1)[1].replace(
                postfix, '')
            thread_1 = threading.Thread(
                target=self.reflect_imported_data_on_main_widgets)
            thread_2 = threading.Thread(
                target=self.create_importing_prgbar_window)
            thread_1.start()
            thread_2.start()
    
    def reflect_imported_data_on_main_widgets(self):
        def load_picklefile(path):
            with open(path, 'rb') as web:
                loadfile = pickle.load(web)
                return loadfile
        self.target_table = pd.read_pickle(self.imported_data)
        pickle_files = glob(f'{self.prj_path}/*.pickle')
        for path in pickle_files:
            if f'{self.user_prefix}_extracted_msms.pickle' in path:
                self.msms_dict = load_picklefile(path)
            elif f'{self.user_prefix}_cid_result.pickle' in path:
                self.cid_result_dict = load_picklefile(path)
            elif f'{self.user_prefix}_oad_result.pickle' in path:
                self.oad_result_dict = load_picklefile(path)
            elif f'{self.user_prefix}_structure_info.pickle' in path:
                self.lipid_structural_info_dict = load_picklefile(path)
            elif f'{self.user_prefix}_graph_info.pickle' in path:
                self.graph_dict = load_picklefile(path)
        cols = self.target_table.columns.tolist()
        if 'pmol/mg tissue' in cols: self.normalized_data = True
        else: self.normalized_data = False
        self.ms_tolerance_ppm = 15
        self.must_nl_cut_off_dict = {'diagnostic_1': ['OAD03', 0.01], 
                                     'diagnostic_2': ['OAD16', 0.01], 
                                     'sphingobase': 0.005}
        self.set_analysis_result_into_table()
        self.data_setting(event=None, is_first=True)
        self.activate_widgets()
        self.temp_prgbar.close_window()

    def create_importing_prgbar_window(self):
        self.temp_prgbar = ProgressBar(self.master,
            title="Data import in progress", detail="--- Data importing ---")
    #endregion

    #region Save data
    def start_saving(self):
        thread_1 = threading.Thread(target=self.popup_for_saving)
        thread_2 = threading.Thread(target=self.save_analysis_files)
        thread_1.start()
        thread_2.start()

    def save_analysis_files(self):
        dataframe_path = f'{self.prj_path}/{self.user_prefix}_analysis_table.pkl'
        msms_path = f'{self.prj_path}/{self.user_prefix}_extracted_msms.pickle'
        cid_reslut_path = f'{self.prj_path}/{self.user_prefix}_cid_result.pickle'
        oad_result_path = f'{self.prj_path}/{self.user_prefix}_oad_result.pickle'
        str_info_path = f'{self.prj_path}/{self.user_prefix}_structure_info.pickle'
        graph_info_path = f'{self.prj_path}/{self.user_prefix}_graph_info.pickle'
        self.target_table.to_pickle(dataframe_path)
        with open(msms_path, 'wb') as output_msms_file:
            pickle.dump(self.msms_dict, output_msms_file)
        with open(cid_reslut_path, 'wb') as output_cid_file:
            pickle.dump(self.cid_result_dict, output_cid_file)
        with open(oad_result_path, 'wb') as output_oad_file:
            pickle.dump(self.oad_result_dict, output_oad_file)
        with open(str_info_path, 'wb') as output_dict_file:
            pickle.dump(self.lipid_structural_info_dict, output_dict_file)
        with open(graph_info_path, 'wb') as output_graph_file:
            pickle.dump(self.graph_dict, output_graph_file)
        time.sleep(2)
        self.temp_popup.close_window()

    def popup_for_saving(self):
        self.temp_popup = PopUpWindow(master=self.master,
            title='Saving files', message='Saving files ...')
    #endregion

    #region Data preparation
    def merge_bipolarity_cid(self):
        self.disable_widgets()
        self.twopath_importer = TwoPathImporter(master=self.master,
            title="Merge Neg&Pos CID-MS/MS data",
            type1="Neg(-) ion mode", type2="Pos(+) ion mode", 
            fmt1="Neg", fmt2="Pos")
        self.activate_widgets()

    def merge_cid_and_oad(self):
        self.disable_widgets()
        self.twopath_importer = TwoPathImporter(master=self.master,
            title="Merge CID&OAD MS/MS data",
            type1="CID-MS/MS", type2="OAD-MS/MS", fmt1="CID", fmt2="OAD")
        self.activate_widgets()
    #endregion

    #region Export window
    def create_data_export_window(self):
        self.export_window = tk.Toplevel(self.master)
        self.export_window.attributes('-topmost', False)
        self.center_position(self.export_window, width=690, height=120)
        self.export_window.resizable(width=False, height=False)
        self.export_window.title("Result data export")
        self.export_window.grab_set()
        self.export_window.iconphoto(False, tk.PhotoImage(file=IconPath))
        self.create_data_export_widgets()
        self.export_window.protocol('WM_DELETE_WINDOW', lambda : self.close_window(self.export_window))

    def create_data_export_widgets(self):
        # Export path
        exporter_label = ttk.Label(self.export_window, text="Directory :")
        exporter_label.place(x=30, y=20)
        self.export_path = tk.StringVar()
        self.export_path_entry = ttk.Entry(
            self.export_window, textvariable=self.export_path, width=75)
        self.export_path_entry.place(x=100, y=20)
        self.export_path_btn = ttk.Button(
            self.export_window, text="Select", command=self.get_export_path)
        self.export_path_btn.place(x=580, y=18)
        # txt file check
        self.txt_export_bool = tk.BooleanVar()
        self.txt_export_bool.set(False)
        self.txt_export_check = ttk.Checkbutton(
            self.export_window, text="Result sheet (.txt)", 
            variable=self.txt_export_bool)
        self.txt_export_check.place(x=30, y=68)
        # xlsx file check
        self.excel_export_bool = tk.BooleanVar()
        self.excel_export_bool.set(True)
        self.excel_export_check = ttk.Checkbutton(
            self.export_window, text="Analysis report (.xlsx)", 
            variable=self.excel_export_bool)
        self.excel_export_check.place(x=190, y=68)
        # pptx file check
        self.pptx_export_bool = tk.BooleanVar()
        self.pptx_export_bool.set(False)
        self.pptx_export_check = ttk.Checkbutton(
            self.export_window, text="MS/MS figures (.pptx)", 
            variable=self.pptx_export_bool)
        self.pptx_export_check.place(x=370, y=68)
        # NL HeatMap check
        # self.nl_heatmap_export_bool = tk.BooleanVar()
        # self.nl_heatmap_export_bool.set(False)
        # self.nl_heatmap_export_check = ttk.Checkbutton(
        #     self.export_window, text="HeatMap of NL (.xlsx)", 
        #     variable=self.nl_heatmap_export_bool)
        # self.nl_heatmap_export_check.place(x=30, y=98)
        # Export button
        self.export_btn = ttk.Button(self.export_window, text="Export", command=self.check_export_parameter)
        self.export_btn.place(x=580, y=68)

    def get_export_path(self):
        old_path = self.export_path_entry.get()
        new_path = filedialog.askdirectory(title="Select export directory")
        if new_path == "":
            self.export_path.set(old_path)
        else:
            self.export_path.set(new_path)

    def check_export_parameter(self):
        export_path = self.export_path_entry.get()
        txt_bool = self.txt_export_bool.get()
        excel_bool = self.excel_export_bool.get()
        pptx_bool = self.pptx_export_bool.get()
        # heatmap_bool = self.nl_heatmap_export_bool.get()
        if export_path == "":
            messagebox.showwarning("Select export directory", "Please select export directory")
        elif not any([txt_bool, excel_bool, pptx_bool]):
            messagebox.showwarning("Select export data format", "Please select export data format")
        else:
            self.start_export()    

    def start_export(self):
        self.temp_vert_prgbar = VerticalProgressBars(
            master=self.master, title="Result export in progress",
            start_text="Start exporting")
        thread_1 = threading.Thread(
            target=self.start_vertical_prgbar_window)
        thread_2 = threading.Thread(target=self.export_analysis_data)
        thread_1.start()
        thread_2.start()

    def export_analysis_data(self):
        start = time.time()
        export_path = self.export_path_entry.get()
        txt_bool = self.txt_export_bool.get()
        excel_bool = self.excel_export_bool.get()
        pptx_bool = self.pptx_export_bool.get()
        self.export_window.destroy()
        sections = [txt_bool, excel_bool, pptx_bool, pptx_bool].count(True)
        self.temp_vert_prgbar.section_report.set("Start export")
        self.temp_vert_prgbar.section_bar["maximum"] = sections+1
        if txt_bool:
            self.temp_vert_prgbar.each_prgbar["value"] = 0
            self.temp_vert_prgbar.each_prgbar["maximum"] = 100
            self.temp_vert_prgbar.section_report.set(
                "Exporting result table of text format")
            time.sleep(2)
            self.temp_vert_prgbar.section_bar.step(1)
            text_sheet_exporter(path=export_path, 
                target_table=self.target_table, cid=self.cid_result_dict, 
                oad=self.oad_result_dict, 
                structure_info=self.lipid_structural_info_dict, 
                normalized=self.normalized_data, stamp=self.user_prefix)
            self.temp_vert_prgbar.each_prgbar.step(99)
            time.sleep(2)
        if excel_bool:
            self.temp_vert_prgbar.each_prgbar["value"] = 0
            self.temp_vert_prgbar.each_prgbar["maximum"] \
                = len(self.target_table)+1
            self.temp_vert_prgbar.section_report.set(
                "Exporting analysis report (.xlsx)")
            excel_sheet_exporter2(path=export_path, 
                target_table=self.target_table,
                msms=self.msms_dict, cid=self.cid_result_dict, 
                oad=self.oad_result_dict, 
                structure_info=self.lipid_structural_info_dict, 
                normalized=self.normalized_data, stamp=self.user_prefix,
                each_bar=self.temp_vert_prgbar.each_prgbar,
                each_rep=self.temp_vert_prgbar.each_report)
            self.temp_vert_prgbar.section_bar.step(1)
            time.sleep(2)
        if pptx_bool:
            self.temp_vert_prgbar.section_bar.step(1)
            self.temp_vert_prgbar.each_prgbar["value"] = 0
            self.temp_vert_prgbar.each_prgbar["maximum"] \
                = 2*len(self.target_table)+1
            self.temp_vert_prgbar.section_report.set(
                "Generating OAD-MS/MS figures")
            generate_msms_figures(path=export_path, msms_dict=self.msms_dict, 
                cid=self.cid_result_dict, oad=self.oad_result_dict,
                structure_info=self.lipid_structural_info_dict, 
                graph=self.graph_dict, stamp=self.user_prefix,
                each_bar=self.temp_vert_prgbar.each_prgbar,
                each_rep=self.temp_vert_prgbar.each_report)
            self.temp_vert_prgbar.section_report.set(
                "Saving OAD-MS/MS figures (.pptx)")
            self.temp_vert_prgbar.section_bar.step(1)
            save_msms_fig_as_pptx(path=export_path, stamp=self.user_prefix, 
                target_table=self.target_table,
                each_bar=self.temp_vert_prgbar.each_prgbar,
                each_rep=self.temp_vert_prgbar.each_report)
            time.sleep(1)
        self.temp_vert_prgbar.section_report.set("Export completed.")
        self.temp_vert_prgbar.section_bar.step(0.99)
        time.sleep(2)
        print(f'Export required time: {time.time() - start:.3f} [sec]')
        self.temp_vert_prgbar.running = False
        self.temp_vert_prgbar.close_window()
        # self.temp_prgbar.close_window()
    
    def start_vertical_prgbar_window(self):
        self.temp_vert_prgbar.process_timer()
        
    def create_export_prgbar_window(self):
        self.temp_prgbar = ProgressBar(self.master,
            title="Result export in progress", 
            detail="--- Result exporting ---")

    #endregion

    #region ProgressBar window
    def create_progress_bar_window(self):
        self.progressbar_window = tk.Toplevel(self.master)
        self.progressbar_window.attributes('-topmost', True)
        self.center_position(self.progressbar_window, width=690, height=155)
        self.progressbar_window.resizable(width=False, height=False)
        self.progressbar_window.title("Data analysis in progress")
        self.progressbar_window.grab_set()
        self.progressbar_window.iconphoto(False, tk.PhotoImage(file=IconPath))
        self.progressbar_window.protocol(
            'WM_DELETE_WINDOW', lambda : self.close_window(self.setting_window))
        self.create_progressbar_widgets()
        self.data_analysis_running = True
        self.analysis_start = True
        self.process_timer()
    
    def create_progressbar_widgets(self):
        # Section Bar
        self.section_bar = ttk.Progressbar(self.progressbar_window, 
            orient='horizontal', length=534, mode='determinate')
        self.section_bar.place(x=30, y=30)
        self.section_report = tk.StringVar()
        self.section_report.set("Start analysis")
        section_label = ttk.Label(self.progressbar_window, 
            textvariable=self.section_report, font=("", 14))
        section_label.place(x=30, y=55)
        # Timer
        self.section_timer = tk.StringVar()
        self.section_timer_label = ttk.Label(
            self.progressbar_window, 
            textvariable=self.section_timer, font=("", 14))
        self.section_timer_label.place(x=580, y=28)
        self.section_timer.set("00:00:00")
        # Each progress
        self.each_prgbar = ttk.Progressbar(self.progressbar_window, 
            orient='horizontal', length=620, mode='determinate')
        self.each_prgbar.place(x=30, y=90)
        self.each_report = tk.StringVar()
        self.each_report.set("")
        each_label = ttk.Label(self.progressbar_window, 
            textvariable=self.each_report, font=("", 14))
        each_label.place(x=30, y=115)
        # Timer
        self.small_timer = tk.StringVar()
        # self.small_timer_label = ttk.Label(self.progressbar_window, textvariable=self.small_timer, font=("", 14))
        # self.small_timer_label.place(x=580, y=88)
        # self.small_timer.set("00:00:00")

    def process_timer(self):
        hour = 0
        minute = 0
        second = 0
        while self.data_analysis_running:
            second += 1
            time.sleep(1)
            if second == 60:
                minute += 1
                second = 0
                if minute == 60:
                    minute = 0
                    hour += 1
            if minute >= 10 and second >= 10:
                self.section_timer.set("0{h}:{m}:{s}".format(h=str(hour),
                m=str(minute), s=str(second)))
            elif minute < 10 and second >= 10:
                self.section_timer.set("0{h}:0{m}:{s}".format(h=str(hour),
                m=str(minute), s=str(second)))
            elif minute >= 10 and second < 10:
                self.section_timer.set("0{h}:{m}:0{s}".format(h=str(hour),
                m=str(minute), s=str(second)))
            else:
                self.section_timer.set("0{h}:0{m}:0{s}".format(h=str(hour),
                m=str(minute), s=str(second)))
        if not self.data_analysis_running:
            self.progressbar_window.destroy()
    #endregion

    #region Updating main window
    def set_analysis_result_into_table(self):
        self.create_result_table_tree(self.result_table_pwin)
        self.create_db_1_table_tree(self.db_1_pwin)
        self.create_db_2_table_tree(self.db_2_pwin)
        self.create_db_3_table_tree(self.db_3_pwin)
        self.create_matched_fragments_table_tree(self.matched_fragments_pwin)
        self.comment_dict = {}
        tag = 'white'
        ion_v_col = 'pmol/mg tissue' if self.normalized_data else 'Height'
        for i, (row, df) in enumerate(self.target_table.iterrows()):
            # tree = ['ID', 'RT', 'm/z', 'm/z type', 'Solved','Result name', 
            #         'Height', 'Ontology', 'Comment']
            idx, rt, mz = df['ID'],math_floor(df['RT(min)'],3),df['Precise m/z']
            mz_type, solved = df['Precise m/z type'], df['Solved level']
            result_name, ion_v = df['OAD result name'], df[ion_v_col]
            data_from = df['Data from']
            ontology = df['Ontology']
            comment = df['User comment']
            tag = 'colored' if tag == 'white' else 'white'
            self.result_table_tree.insert("", "end", tags=tag,
                values=[idx, rt, mz, mz_type, solved, result_name, 
                        data_from, ion_v, ontology, comment])
        new_label = "Result table (molecules={})".format(i+1)
        self.left_wrap_label_change(new_label)
        self.result_table_tree.tag_configure('colored', background="#dcdcdc")
        self.result_table_pwin.add(self.result_table_tree)
        first_iid = self.result_table_tree.get_children()[0]
        self.result_table_tree.focus(first_iid)
        self.result_table_tree.selection_set(first_iid)
        self.result_table_tree.bind("<<TreeviewSelect>>", self.data_setting)

    def left_wrap_label_change(self, txt):
        self.left_wrap_lebel.set(txt)
        self.left_wrap_frame.configure(text=self.left_wrap_lebel.get())

    def matched_ions_label_label_change(self, txt):
        self.matched_ions_label.set(txt)
        self.matched_fragments_frame.configure(text=self.matched_ions_label.get())

    def data_setting(self, event, is_first=False):
        #region CID result dict structure
        # cid_dict = {'Lipid subclass': {'Glycine': [ref_mz, measured_mz, ppm],
        #                                'Presence': 50.00},
        #             'Moiety':         {'acyl-2': [ref_mz, measured_mz, ppm],
        #                                'Presence': 50.00}
        #            }
        #endregion
        #region OAD result dict structure
        # result_dict = {'Resolved level': 'All' or 'Partial' or 'None', 
        #                'Validated num': 0 ~ 3,
        #                'Each bools': [True, False, ...], 
        #                'Moiety-1': {0: {'Positions': '',
        #                                 'N-description': '',
        #                                 'Score': float,
        #                                 'Ratio sum': float,
        #                                 'Presence': float,
        #                                 'Notice': '',
        #                                 'Peaks dict': {'n-9/dis@n-8/+O/': [Ref m/z, Ref delta, Measured m/z, Measured ratio, ppm]}
        #                                },
        #                             1: {....},
        #                             'Determined db: {'Positions': '',
        #                                              'N-description': '',
        #                                              'Score': float,
        #                                              'Ratio sum': float,
        #                                              'Presence': float,
        #                                              'Notice': '',
        #                                              'Measured peaks': [[Measured m/z, Measured ratio, ppm], [...]],
        #                                              'Ref peaks': [[OAD type, Ref m/z, Ref NL, Ref ratio], [...]]
        #                                             }
        #                             }
        #               }
        #endregion
        #region Extract target ID
        # if is_first:
        #     first_iid = self.result_table_tree.get_children()[0]
        #     str_idx = self.result_table_tree.item(first_iid, 'value')[0]
        #     idx = int(str_idx)
        # else:
        #     selected_items = self.result_table_tree.selection()
        #     selected_values = self.result_table_tree.item(
        #         selected_items[0], 'values')
        #     str_idx = selected_values[0]
        #     idx = int(str_idx)
        #endregion
        focused_item = self.get_focused_item(is_first=is_first)
        idx = int(focused_item[0])
        self.selected_idx = idx
        result_dict = self.oad_result_dict[idx]
        result_dict_len = len(result_dict)
        self.temp_msms_df = self.msms_dict[idx]
        self.temp_lipid_info = self.lipid_structural_info_dict[idx]
        self.temp_graph_info = self.graph_dict[idx]
        resolved_level = result_dict['Resolved level']
        each_bools = result_dict['Each bools']
        moiety_1, moiety_2, moiety_3 = {}, {}, {}
        #region Update db trees and set current C=C pos
        if result_dict_len == 4:
            moiety_1 = result_dict['Moiety-1']
        elif result_dict_len == 5:
            moiety_1 = result_dict['Moiety-1']
            moiety_2 = result_dict['Moiety-2']
        elif result_dict_len == 6:
            moiety_1 = result_dict['Moiety-1']
            moiety_2 = result_dict['Moiety-2']
            moiety_3 = result_dict['Moiety-3']
        self.update_db_trees(result_dict=result_dict, is_first=is_first)
        if moiety_1 and moiety_1[0]:
            self.db_1 = moiety_1[0]['N-description']
        else:
            self.db_1 = ''
        if moiety_2 and moiety_2[0]:
            self.db_2 = moiety_2[0]['N-description']
        else:
            self.db_2 = ''
        if moiety_3 and moiety_3[0]:
            self.db_3 = moiety_3[0]['N-description']
        else:
            self.db_3 = ''
        self.focus_and_select_db_trees()
        #endregion
        #region Update total MS/MS fig
        cid_dict = self.cid_result_dict[idx]
        self.update_msms_fig(cid_info=cid_dict)
        #endregion
        #region Update OAD-MS/MS fig with reference
        self.update_oad_msms_fig_with_ref(
            dict_1=moiety_1, dict_2=moiety_2, dict_3=moiety_3)
        if is_first:
            # func(self.partially_update_data) occurs 3 times unnecessarily
            # due to the binding problem
            self.db_candidates_frame.bind_class(
                "Treeview", "<<TreeviewSelect>>", self.partially_update_data)
        #endregion

    def update_db_trees(self, result_dict, is_first=False, masstable=True):
        result_dict_len = len(result_dict)
        moiety_1, moiety_2, moiety_3 = {}, {}, {}
        if result_dict_len == 4:
            moiety_1 = result_dict['Moiety-1']
            if moiety_1:
                self.fill_in_data_on_db_tree(
                    tree_no=1, data=moiety_1, is_first=is_first)
                self.fill_in_nan_on_db_trees(
                    db_2=True, db_3=True, is_first=is_first)
                if masstable:
                    self.fill_in_matched_fragment_ions_tree(
                        moiety_1=moiety_1[0], is_first=is_first)
            else:
                self.fill_in_nan_on_db_trees(
                    db_1=True, db_2=True, db_3=True, is_first=is_first)
                if masstable:
                    self.fill_in_matched_fragment_ions_tree(is_first=is_first)
        elif result_dict_len == 5:
            moiety_1 = result_dict['Moiety-1']
            moiety_2 = result_dict['Moiety-2']
            if moiety_1 and moiety_2:
                self.fill_in_data_on_db_tree(
                    tree_no=1, data=moiety_1, is_first=is_first)
                self.fill_in_data_on_db_tree(
                    tree_no=2, data=moiety_2, is_first=is_first)
                self.fill_in_nan_on_db_trees(
                    db_3=True, is_first=is_first)
                if masstable:
                    self.fill_in_matched_fragment_ions_tree(
                        moiety_1=moiety_1[0], moiety_2=moiety_2[0], 
                        is_first=is_first)
        elif result_dict_len == 6:
            moiety_1 = result_dict['Moiety-1']
            moiety_2 = result_dict['Moiety-2']
            moiety_3 = result_dict['Moiety-3']
            if moiety_1 and moiety_2 and moiety_3:
                self.fill_in_data_on_db_tree(
                    tree_no=1, data=moiety_1, is_first=is_first)
                self.fill_in_data_on_db_tree(
                    tree_no=2, data=moiety_2, is_first=is_first)
                self.fill_in_data_on_db_tree(
                    tree_no=3, data=moiety_3, is_first=is_first)
                if masstable:
                    self.fill_in_matched_fragment_ions_tree(
                        moiety_1=moiety_1[0], moiety_2=moiety_2[0], 
                        moiety_3=moiety_3[0], is_first=is_first)

    def get_focused_item(self, is_first=False):
        if is_first:
            first_iid = self.result_table_tree.get_children()[0]
            selected_values = self.result_table_tree.item(first_iid, 'value')
            # idx = int(selected_values[0])
        else:
            selected_items = self.result_table_tree.selection()
            selected_values = self.result_table_tree.item(
                selected_items[0], 'values')
            # str_idx = selected_values[0]
            # idx = int(str_idx)
        return selected_values

    def fill_in_nan_on_db_trees(
        self, db_1=False, db_2=False, db_3=False, is_first=False):
        # tree.column = ['Rank', 'moiety', 'Score', 'Presence', 'Rel. Int']
        li=['N/A', 'N/A', 'N/A', 'N/A', 'N/A']
        if db_1:
            self.db_1_tree.delete(*self.db_1_tree.get_children())
            # if db_1: li=['###', 'Unresolved', '###', '###', '###']
            self.db_1_tree.insert("", "end", values=li)
            if is_first:
                self.db_1_pwin.add(self.db_1_tree)
        if db_2:
            self.db_2_tree.delete(*self.db_2_tree.get_children())
            # if db_2: li=['###', 'Unresolved', '###', '###', '###']
            self.db_2_tree.insert("", "end", values=li)
            if is_first:
                self.db_2_pwin.add(self.db_2_tree)
        if db_3:
            self.db_3_tree.delete(*self.db_3_tree.get_children())
            # if db_3: li=['###', 'Unresolved', '###', '###', '###']
            self.db_3_tree.insert("", "end", values=li)
            if is_first:
                self.db_3_pwin.add(self.db_3_tree)

    def fill_in_data_on_db_tree(self, tree_no, data, is_first=False):
        # d = {Positions:~, N-description:~, Score:~, Ratio sum:~, 
        #      Presence:~, Notice:~, Peaks dict:{}}
        # tree.column = ['Rank', 'moiety-1', 'Score', 'Presence', 'Rel. Int']
        def check_no_unresolved_in_data(data_d):
            bools = []
            for rank, d in data_d.items():
                if rank == 'Determined db':
                    continue
                if d:
                    if 'Unresolved' not in d['Notice']:
                        bools.append(True)
                    else:
                        bools.append(False)
                else:
                    bools.append(True)
            return all(bools)
        li=['###', 'Unresolved', '###', '###', '###']
        last = len(data)
        if tree_no == 1:
            self.db_1_tree.delete(*self.db_1_tree.get_children())
            bg = "#fdedec" #red
            for key, d in data.items():
                if key == 'Determined db':
                    continue
                tag = 'colored' if key%2 == 1 else 'white'
                # if 'Unresolved' not in d['Notice']:
                self.db_1_tree.insert("", "end", tags=tag, 
                    values=[key+1, d['N-description'], d['Score'], 
                    d['Presence'], d['Ratio sum']])
                # else:
                #     self.db_1_tree.insert("", "end", tags=tag, values=li)
            # if check_no_unresolved_in_data(data):
            #     tag = 'colored' if last%2 == 1 else 'white'
            #     self.db_1_tree.insert("", "end", tags=tag, values=li)
            if is_first:
                self.db_1_pwin.add(self.db_1_tree)
            self.db_1_tree.tag_configure('colored', background=bg)
        if tree_no == 2:
            self.db_2_tree.delete(*self.db_2_tree.get_children())
            bg = "#ebf5fb" #blue
            for key, d in data.items():
                if key == 'Determined db':
                    continue
                tag = 'colored' if key%2 == 1 else 'white'
                # if 'Unresolved' not in d['Notice']:
                self.db_2_tree.insert("", "end", tags=tag, 
                    values=[key+1, d['N-description'], d['Score'], 
                    d['Presence'], d['Ratio sum']])
                # else:
                #     self.db_2_tree.insert("", "end", tags=tag, values=li)
            # if check_no_unresolved_in_data(data):
            #     tag = 'colored' if last%2 == 1 else 'white'
            #     self.db_2_tree.insert("", "end", tags=tag, values=li)
            if is_first:
                self.db_2_pwin.add(self.db_2_tree)
            self.db_2_tree.tag_configure('colored', background=bg)
        if tree_no == 3:
            self.db_3_tree.delete(*self.db_3_tree.get_children())
            bg = "#e9f7ef" #green
            for key, d in data.items():
                if key == 'Determined db':
                    continue
                tag = 'colored' if key%2 == 1 else 'white'
                # if 'Unresolved' not in d['Notice']:
                self.db_3_tree.insert("", "end", tags=tag, 
                    values=[key+1, d['N-description'], d['Score'], 
                    d['Presence'], d['Ratio sum']])
            #     else:
            #         self.db_3_tree.insert("", "end", tags=tag, values=li)
            # if check_no_unresolved_in_data(data):
            #     tag = 'colored' if last%2 == 1 else 'white'
            #     self.db_3_tree.insert("", "end", tags=tag, values=li)
            if is_first:
                self.db_3_pwin.add(self.db_3_tree)
            self.db_3_tree.tag_configure('colored', background=bg)

    def fill_in_matched_fragment_ions_tree(self, moiety_1=False, moiety_2=False, 
        moiety_3=False, is_first=False):
        # tree = ['m/z', 'Intensity', 'Ratio(%)', 'Delta', 
        #         'Ref m/z1', 'ppm1', 'NL type1', 
        #         'Ref m/z2', 'ppm2', 'NL type2', 
        #         'Ref m/z3', 'ppm3', 'NL type3']
        # 'Peaks dict': {'n-9/dis@n-8/+O/': [Ref m/z, Ref delta, Measured m/z, 
        #                                    Measured ratio, ppm]}
        msms = self.temp_msms_df
        self.matched_fragments_tree.delete(
            *self.matched_fragments_tree.get_children())
        def get_peak_dict(d):
            return d['Peaks dict'] if d else {'none': [0, 0, 0, 0, 0]}
        def get_input_values(mz, d):
            ref_mz, ppm, oad_type, count = '', '', '', 0
            for key, li in d.items():
                if mz == li[2]:
                    if count == 0: oad_type = key
                    else: oad_type = f'{oad_type}|{key}'
                    ref_mz, ppm = li[0], li[4]
                    count += 1
            return ref_mz, ppm, oad_type
        peak_dict_1 = get_peak_dict(moiety_1)
        peak_dict_2 = get_peak_dict(moiety_2)
        peak_dict_3 = get_peak_dict(moiety_3)
        mzs = msms['frag m/z'].values.tolist()
        intensities = msms['intensity'].values.tolist()
        ratios = msms['Ratio(%)'].values.tolist()
        deltas = msms['Delta'].values.tolist()
        for i, (mz, intensity, ratio, delta) in enumerate(
            zip(mzs, intensities, ratios, deltas)):
            tag = 'colored' if i%2 == 1 else 'white'
            ref_mz_1, ppm_1, oad_type_1 = get_input_values(mz, peak_dict_1)
            ref_mz_2, ppm_2, oad_type_2 = get_input_values(mz, peak_dict_2)
            ref_mz_3, ppm_3, oad_type_3 = get_input_values(mz, peak_dict_3)
            iid = self.matched_fragments_tree.insert("", "end", tags=tag, 
            values=[mz, intensity, ratio, delta, ref_mz_1, ppm_1, oad_type_1,
                    ref_mz_2, ppm_2, oad_type_2, ref_mz_3, ppm_3, oad_type_3])
        new_label = "Mass table (ions={})".format(i+1)
        self.matched_ions_label_label_change(new_label)
        self.matched_fragments_tree.tag_configure('colored', background="#dcdcdc")
        if is_first:
            self.matched_fragments_pwin.add(self.matched_fragments_tree)

    def switch_matched_fragment_ions_tree(self):
        switch = self.show_only_matched_ions.get()
        db_1_pos, db_2_pos, db_3_pos = self.db_1, self.db_2, self.db_3
        result_dict = self.oad_result_dict[self.selected_idx]
        result_dict_len = len(result_dict)
        msms = self.temp_msms_df
        moiety_1, moiety_2, moiety_3 = {}, {}, {}
        if result_dict_len >= 4:
            for rank, d in result_dict['Moiety-1'].items():
                if db_1_pos == d['N-description']:
                    moiety_1 = d
        if result_dict_len >= 5:
            for rank, d in result_dict['Moiety-2'].items():
                if db_2_pos == d['N-description']:
                    moiety_2 = d
        if result_dict_len >= 6:
            for rank, d in result_dict['Moiety-3'].items():
                if db_3_pos == d['N-description']:
                    moiety_3 = d
        self.matched_fragments_tree.delete(*self.matched_fragments_tree.get_children())
        # 'Peaks dict': {'n-9/dis@n-8/+O/': [Ref m/z, Ref delta, Measured m/z, Measured ratio, ppm]}
        def get_peak_dict(d):
            return d['Peaks dict'] if d else {'none': [0, 0, 0, 0, 0]}
        def get_filling_values(mz, d, num):
            ref_mz, ppm, oad_type, count = '', '', '', 0
            for key, li in d.items():
                if mz == li[2]:
                    if count == 0: oad_type = key
                    else: oad_type = f'{oad_type}|{key}'
                    ref_mz, ppm = li[0], li[4]
                    count += 1
            return ref_mz, ppm, oad_type
        peak_dict_1 = get_peak_dict(moiety_1)
        peak_dict_2 = get_peak_dict(moiety_2)
        peak_dict_3 = get_peak_dict(moiety_3)
        if switch:
            tag = 'white'
            count = 0
            for i, (row, df) in enumerate(msms.iterrows()):
                mz = df['frag m/z']
                intensity = int(df['intensity'])
                ratio = df['Ratio(%)']
                delta = df['Delta']
                ref_mz_1, ppm_1, oad_type_1 = get_filling_values(mz, peak_dict_1, num=1)
                ref_mz_2, ppm_2, oad_type_2 = get_filling_values(mz, peak_dict_2, num=2)
                ref_mz_3, ppm_3, oad_type_3 = get_filling_values(mz, peak_dict_3, num=3)
                if ref_mz_1 != '' or ref_mz_2 != '' or ref_mz_3 != '':
                    tag = 'colored' if tag == 'white' else 'white'
                    count += 1
                    iid = self.matched_fragments_tree.insert("", "end", tags=tag, 
                    values=[mz, intensity, ratio, delta, ref_mz_1, ppm_1, oad_type_1,
                            ref_mz_2, ppm_2, oad_type_2, ref_mz_3, ppm_3, oad_type_3])
            new_label = f'Mass table (ions={count})'
            self.matched_ions_label_label_change(new_label)
            self.matched_fragments_tree.tag_configure('colored', background="#dcdcdc")
        else:
            for i, (row, df) in enumerate(msms.iterrows()):
                mz = df['frag m/z']
                intensity = int(df['intensity'])
                ratio = df['Ratio(%)']
                delta = df['Delta']
                tag = 'colored' if i%2 == 1 else 'white'
                ref_mz_1, ppm_1, oad_type_1 = get_filling_values(mz, peak_dict_1, num=1)
                ref_mz_2, ppm_2, oad_type_2 = get_filling_values(mz, peak_dict_2, num=2)
                ref_mz_3, ppm_3, oad_type_3 = get_filling_values(mz, peak_dict_3, num=3)
                iid = self.matched_fragments_tree.insert("", "end", tags=tag, 
                values=[mz, intensity, ratio, delta, ref_mz_1, ppm_1, oad_type_1,
                        ref_mz_2, ppm_2, oad_type_2, ref_mz_3, ppm_3, oad_type_3])
            new_label = f'Mass table (ions={i+1})'
            self.matched_ions_label_label_change(new_label)
            self.matched_fragments_tree.tag_configure('colored', background="#dcdcdc")
    
    def switch_result_table_tree(self):
        switch = self.show_only_solved_mol.get()
        target_idx = self.selected_idx
        selected_iid = self.result_table_tree.get_children()[0]
        self.result_table_tree.delete(*self.result_table_tree.get_children())
        ion_v_col = 'pmol/mg tissue' if self.normalized_data else 'Height'
        tag = 'white'
        if switch:
            count = 0
            for i, (row, df) in enumerate(self.target_table.iterrows()):
                # tree = ['ID', 'RT', 'm/z', 'm/z type', 'Solved','Result name', 'Height', 'Ontology', 'Comment']
                solved = df['Solved level']
                if solved == 'None': continue
                idx = df['ID']
                rt = math_floor(df['RT(min)'], 3)
                mz = df['Precise m/z']
                mz_type = df['Precise m/z type']
                result_name = df['OAD result name']
                data_from = df['Data from']
                ion_v = df[ion_v_col]
                ontology = df['Ontology']
                comment = df['User comment']
                tag = 'colored' if tag == 'white' else 'white'
                count += 1
                iid = self.result_table_tree.insert("", "end", tags=tag,
                        values=[idx, rt, mz, mz_type, solved, result_name, 
                                data_from, ion_v, ontology, comment])
                if idx == target_idx:
                    selected_iid = iid
            new_label = f'Result table (molecules={count})'
        else:
            for i, (row, df) in enumerate(self.target_table.iterrows()):
                # tree = ['ID', 'RT', 'm/z', 'm/z type', 'Solved','Result name', 'Height', 'Ontology', 'Comment']
                idx = df['ID']
                rt = math_floor(df['RT(min)'], 3)
                mz = df['Precise m/z']
                mz_type = df['Precise m/z type']
                solved = df['Solved level']
                result_name = df['OAD result name']
                data_from = df['Data from']
                ion_v = df[ion_v_col]
                ontology = df['Ontology']
                comment = df['User comment']
                tag = 'colored' if tag == 'white' else 'white'
                iid = self.result_table_tree.insert("", "end", tags=tag,
                        values=[idx, rt, mz, mz_type, solved, result_name, 
                                data_from, ion_v, ontology, comment])
                if idx == target_idx:
                    selected_iid = iid
            new_label = f'Result table (molecules={i+1})'
        self.left_wrap_label_change(new_label)
        self.result_table_tree.tag_configure('colored', background="#d3d3d3")
        try:
            self.result_table_tree.focus(selected_iid)
            self.result_table_tree.selection_set(selected_iid)
        except tk.TclError:
            first_iid = self.result_table_tree.get_children()[0]
            self.result_table_tree.focus(first_iid)
            self.result_table_tree.selection_set(first_iid)
        self.result_table_tree.bind("<<TreeviewSelect>>", self.data_setting)

    def focus_and_select_db_trees(self):
        iid_1 = self.db_1_tree.get_children()[0]
        self.db_1_tree.focus(iid_1)
        self.db_1_tree.selection_set(iid_1)
        iid_2 = self.db_2_tree.get_children()[0]
        self.db_2_tree.focus(iid_2)
        self.db_2_tree.selection_set(iid_2)
        iid_3 = self.db_3_tree.get_children()[0]
        self.db_3_tree.focus(iid_3)
        self.db_3_tree.selection_set(iid_3)

    def partially_update_data(self, event):
        db_list = self.get_current_db_pos_and_rank()
        db_1_pos = db_list[0][1]
        db_2_pos = db_list[1][1]
        db_3_pos = db_list[2][1]
        db_1_changed = (db_1_pos != self.db_1)
        db_2_changed = (db_2_pos != self.db_2)
        db_3_changed = (db_3_pos != self.db_3)
        if any([db_1_changed, db_2_changed, db_3_changed]):
            self.db_1, self.db_2, self.db_3 = db_1_pos, db_2_pos, db_3_pos
            result_dict = self.oad_result_dict[self.selected_idx]
            lipid_info = self.lipid_structural_info_dict[self.selected_idx]
            result_dict_len = len(result_dict)
            new_db_1, new_db_2, new_db_3 = {}, {}, {}
            moiety_1, moiety_2, moiety_3 = {}, {}, {}
            if result_dict_len >= 4:
                for rank, d in result_dict['Moiety-1'].items():
                    if rank == 'Determined db':
                        continue
                    if db_1_pos == d['N-description']:
                        new_db_1 = d
                        moiety_1[0] = d
            if result_dict_len >= 5:
                for rank, d in result_dict['Moiety-2'].items():
                    if rank == 'Determined db':
                        continue
                    if db_2_pos == d['N-description']:
                        new_db_2 = d
                        moiety_2[0] = d
            if result_dict_len >= 6:
                for rank, d in result_dict['Moiety-3'].items():
                    if rank == 'Determined db':
                        continue
                    if db_3_pos == d['N-description']:
                        new_db_3 = d
                        moiety_3[0] = d
            self.fill_in_matched_fragment_ions_tree(
                moiety_1=new_db_1, moiety_2=new_db_2, moiety_3=new_db_3)
            self.update_oad_msms_fig_with_ref(
                dict_1=moiety_1, dict_2=moiety_2, dict_3=moiety_3)

    def get_current_db_pos_and_rank(self):
        # tree.column = ['Rank', 'moiety-1', 'Score', 'Presence', 'Rel. Int']
        def find_rank(pos, moiety):
            for rank, d in moiety.items():
                if rank != 'Determined db' and d['N-description'] == pos:
                    return rank
            return -1
        try:
            tree1_item = self.db_1_tree.item(self.db_1_tree.selection()[0],'values')
            db_1_pos = tree1_item[1]
        except:
            db_1_pos = ''
        try:
            tree2_item = self.db_2_tree.item(self.db_2_tree.selection()[0],'values')
            db_2_pos = tree2_item[1]
        except:
            db_2_pos = ''
        try:
            tree3_item = self.db_3_tree.item(self.db_3_tree.selection()[0],'values')
            db_3_pos = tree3_item[1]
        except:
            db_3_pos = ''
        res = self.oad_result_dict[self.selected_idx]
        rank_1, rank_2, rank_3 = -1, -1, -1
        if len(res) == 4:
            rank_1 = find_rank(db_1_pos, res['Moiety-1'])
        elif len(res) == 5:
            rank_1 = find_rank(db_1_pos, res['Moiety-1'])
            rank_2 = find_rank(db_2_pos, res['Moiety-2'])
        elif len(res) == 6:
            rank_1 = find_rank(db_1_pos, res['Moiety-1'])
            rank_2 = find_rank(db_2_pos, res['Moiety-2'])
            rank_3 = find_rank(db_3_pos, res['Moiety-3'])
        return [[rank_1, db_1_pos], [rank_2, db_2_pos], [rank_3, db_3_pos]]

    def judge_db_change(self, new_db_1, new_db_2, new_db_3):
        old_db_1, old_db_2, old_db_3 = '', '', ''
        result_dict = self.oad_result_dict[self.selected_idx]
        result_dict_len = len(result_dict)
        if result_dict_len >= 4 and result_dict['Moiety-1']:
            try:
                old_db_1 \
                = result_dict['Moiety-1']['Determined db']['N-description']
            except:
                pass
        if result_dict_len >= 5 and result_dict['Moiety-2']:
            try:
                old_db_2 \
                = result_dict['Moiety-2']['Determined db']['N-description']
            except:
                pass
        if result_dict_len >= 6 and result_dict['Moiety-3']:
            try:
                old_db_3 \
                = result_dict['Moiety-3']['Determined db']['N-description']
            except:
                pass
        if new_db_1 == 'N/A':
            db_1_changed = [False, old_db_1]
        elif new_db_1 == old_db_1:
            db_1_changed = [False, old_db_1]
        else:
            db_1_changed = [True, old_db_1]
        if new_db_2 == 'N/A':
            db_2_changed = [False, old_db_2]
        elif new_db_2 == old_db_2:
            db_2_changed = [False, old_db_2]
        else:
            db_2_changed = [True, old_db_2]
        if new_db_3 == 'N/A':
            db_3_changed = [False, old_db_3]
        elif new_db_3 == old_db_3:
            db_3_changed = [False, old_db_3]
        else:
            db_3_changed = [True, old_db_3]
        return [db_1_changed, db_2_changed, db_3_changed]

    def modify_db_pos(self):
        db_list = self.get_current_db_pos_and_rank()
        rank_1, rank_2, rank_3 = db_list[0][0], db_list[1][0], db_list[2][0]
        new_db_1, new_db_2, new_db_3 \
            = db_list[0][1], db_list[1][1], db_list[2][1]
        bool_list = self.judge_db_change(new_db_1, new_db_2, new_db_3)
        old_db_1, old_db_2, old_db_3 \
            = bool_list[0][1], bool_list[1][1], bool_list[2][1]
        db_1_changed, db_2_changed, db_3_changed \
            = bool_list[0][0], bool_list[1][0], bool_list[2][0]
        idx = self.selected_idx
        if any([db_1_changed, db_2_changed, db_3_changed]):
            if db_1_changed:
                self.re_const_oad_res_dict(idx=idx, moiety_num=1, rank=rank_1,
                    old_db=old_db_1, new_db=new_db_1)
            if db_2_changed:
                self.re_const_oad_res_dict(idx=idx, moiety_num=2, rank=rank_2,
                    old_db=old_db_2, new_db=new_db_2)
            if db_3_changed:
                self.re_const_oad_res_dict(idx=idx, moiety_num=3, rank=rank_3,
                    old_db=old_db_3, new_db=new_db_3)
            df = self.target_table[self.target_table['ID'] == idx]
            row = list(df.index)[0]
            old_name = df['OAD result name'].values[0]
            new_name = self.update_oad_result_name(idx=idx, 
                old_name=old_name, db_list=db_list, bool_list=bool_list)
            self.target_table.loc[row:row, ['OAD result name']] = new_name
            self.update_result_table(idx)

    def update_oad_result_name(self, idx, old_name, db_list, bool_list):
        new_db_1, new_db_2, new_db_3 \
            = db_list[0][1], db_list[1][1], db_list[2][1]
        old_db_1, old_db_2, old_db_3 \
            = bool_list[0][1], bool_list[1][1], bool_list[2][1]
        db_1_changed, db_2_changed, db_3_changed \
            = bool_list[0][0], bool_list[1][0], bool_list[2][0]
        moiety_info = self.temp_lipid_info['Each moiety info']
        unsaturated_moiety_num = self.temp_lipid_info['Unsaturated moiety']
        def set_resolved_bool(idx, new_db, acyl_num):
            if new_db == 'Unresolved':
                self.oad_result_dict[idx]['Each bools'][acyl_num] = False
            else:
                self.oad_result_dict[idx]['Each bools'][acyl_num] = True
        if len(moiety_info) == 2:
            if db_1_changed:
                new_name = self.replace_one_pos_in_moiety(info=moiety_info, 
                    old_name=old_name, old_db=old_db_1, new_db=new_db_1)
                set_resolved_bool(idx=idx, new_db=new_db_1, acyl_num=0)
        elif len(moiety_info) == 4:
            if unsaturated_moiety_num == 1:
                if db_1_changed:
                    if moiety_info['db-1'] > 0:
                        new_name = self.replace_one_pos_in_dimoieties(
                                num=1, old_name=old_name, info=moiety_info, 
                                old_db=old_db_1, new_db=new_db_1)
                    else:
                        new_name = self.replace_one_pos_in_dimoieties(
                                num=2, old_name=old_name, info=moiety_info, 
                                old_db=old_db_1, new_db=new_db_1)
                    set_resolved_bool(idx=idx, new_db=new_db_1, acyl_num=0)
            elif unsaturated_moiety_num == 2:
                if db_1_changed and db_2_changed:
                    new_name = self.replace_two_pos_in_dimoieties(
                        info=moiety_info, old_name=old_name, 
                        old_db_1=old_db_1, new_db_1=new_db_1, 
                        old_db_2=old_db_2, new_db_2=new_db_2)
                    set_resolved_bool(idx=idx, new_db=new_db_1, acyl_num=0)
                    set_resolved_bool(idx=idx, new_db=new_db_2, acyl_num=1)
                elif db_1_changed:
                    new_name = self.replace_one_pos_in_dimoieties(
                                num=1, old_name=old_name, info=moiety_info, 
                                old_db=old_db_1, new_db=new_db_1)
                    set_resolved_bool(idx=idx, new_db=new_db_1, acyl_num=0)
                elif db_2_changed:
                    new_name = self.replace_one_pos_in_dimoieties(
                                num=2, old_name=old_name, info=moiety_info, 
                                old_db=old_db_2, new_db=new_db_2)
                    set_resolved_bool(idx=idx, new_db=new_db_2, acyl_num=1)
        elif len(moiety_info) == 6:
            if unsaturated_moiety_num == 1:
                if db_1_changed:
                    if moiety_info['db-1'] > 0:
                        new_name = self.replace_one_pos_in_trimoieties(
                                    num=1, old_name=old_name, info=moiety_info, 
                                    old_db=old_db_1, new_db=new_db_1)
                    elif moiety_info['db-2'] > 0:
                        new_name = self.replace_one_pos_in_trimoieties(
                                    num=2, old_name=old_name, info=moiety_info, 
                                    old_db=old_db_1, new_db=new_db_1)
                    else:
                        new_name = self.replace_one_pos_in_trimoieties(
                                    num=3, old_name=old_name, info=moiety_info, 
                                    old_db=old_db_1, new_db=new_db_1)
                    set_resolved_bool(idx=idx, new_db=new_db_1, acyl_num=0)
            if unsaturated_moiety_num == 2:
                if db_1_changed and db_2_changed:
                    if moiety_info['db-1'] > 0 and moiety_info['db-2'] > 0:
                        new_name = self.replace_two_pos_in_trimoieties(
                                    num1=1, num2=2, old_name=old_name, 
                                    info=moiety_info, 
                                    old_db_1=old_db_1, new_db_1=new_db_1,
                                    old_db_2=old_db_2, new_db_2=new_db_2)
                    elif moiety_info['db-1'] > 0 and moiety_info['db-3'] > 0:
                        new_name = self.replace_two_pos_in_trimoieties(
                                    num1=1, num2=3, old_name=old_name, 
                                    info=moiety_info, 
                                    old_db_1=old_db_1, new_db_1=new_db_1,
                                    old_db_2=old_db_2, new_db_2=new_db_2)
                    else: #moiety_info['db-2'] > 0 and moiety_info['db-3'] > 0:
                        new_name = self.replace_two_pos_in_trimoieties(
                                    num1=2, num2=3, old_name=old_name, 
                                    info=moiety_info, 
                                    old_db_1=old_db_1, new_db_1=new_db_1,
                                    old_db_2=old_db_2, new_db_2=new_db_2)
                    set_resolved_bool(idx=idx, new_db=new_db_1, acyl_num=0)
                    set_resolved_bool(idx=idx, new_db=new_db_2, acyl_num=1)
                elif db_1_changed:
                    if moiety_info['db-1'] > 0 and moiety_info['db-2'] > 0:
                        new_name = self.replace_one_pos_in_trimoieties(
                                    num=1, old_name=old_name, info=moiety_info, 
                                    old_db=old_db_1, new_db=new_db_1)
                    elif moiety_info['db-1'] > 0 and moiety_info['db-3'] > 0:
                        new_name = self.replace_one_pos_in_trimoieties(
                                    num=1, old_name=old_name, info=moiety_info, 
                                    old_db=old_db_1, new_db=new_db_1)
                    else: #moiety_info['db-2'] > 0 and moiety_info['db-3'] > 0:
                        new_name = self.replace_one_pos_in_trimoieties(
                                    num=2, old_name=old_name, info=moiety_info, 
                                    old_db=old_db_1, new_db=new_db_1)
                    set_resolved_bool(idx=idx, new_db=new_db_1, acyl_num=0)
                else: #db_2_changed
                    if moiety_info['db-1'] > 0 and moiety_info['db-2'] > 0:
                        new_name = self.replace_one_pos_in_trimoieties(
                                    num=2, old_name=old_name, info=moiety_info, 
                                    old_db=old_db_2, new_db=new_db_2)
                    elif moiety_info['db-1'] > 0 and moiety_info['db-3'] > 0:
                        new_name = self.replace_one_pos_in_trimoieties(
                                    num=3, old_name=old_name, info=moiety_info, 
                                    old_db=old_db_2, new_db=new_db_2)
                    else: #moiety_info['db-2'] > 0 and moiety_info['db-3'] > 0:
                        new_name = self.replace_one_pos_in_trimoieties(
                                    num=3, old_name=old_name, info=moiety_info, 
                                    old_db=old_db_2, new_db=new_db_2)
                    set_resolved_bool(idx=idx, new_db=new_db_2, acyl_num=1)
            if unsaturated_moiety_num == 3:
                if db_1_changed and db_2_changed and db_3_changed:
                    new_name = self.replace_three_pos_in_trimoieties(
                        info=moiety_info, old_name=old_name, 
                        old_db_1=old_db_1, new_db_1=new_db_1, 
                        old_db_2=old_db_2, new_db_2=new_db_2,
                        old_db_3=old_db_3, new_db_3=new_db_3)
                    set_resolved_bool(idx=idx, new_db=new_db_1, acyl_num=0)
                    set_resolved_bool(idx=idx, new_db=new_db_2, acyl_num=1)
                    set_resolved_bool(idx=idx, new_db=new_db_3, acyl_num=2)
                elif db_1_changed and db_2_changed:
                    new_name = self.replace_two_pos_in_trimoieties(
                                    num1=1, num2=2, old_name=old_name, 
                                    info=moiety_info, 
                                    old_db_1=old_db_1, new_db_1=new_db_1,
                                    old_db_2=old_db_2, new_db_2=new_db_2)
                    set_resolved_bool(idx=idx, new_db=new_db_1, acyl_num=0)
                    set_resolved_bool(idx=idx, new_db=new_db_2, acyl_num=1)
                elif db_1_changed and db_3_changed:
                    new_name = self.replace_two_pos_in_trimoieties(
                                    num1=1, num2=3, old_name=old_name, 
                                    info=moiety_info, 
                                    old_db_1=old_db_1, new_db_1=new_db_1,
                                    old_db_2=old_db_3, new_db_2=new_db_3)
                    set_resolved_bool(idx=idx, new_db=new_db_1, acyl_num=0)
                    set_resolved_bool(idx=idx, new_db=new_db_3, acyl_num=2)
                elif db_2_changed and db_3_changed:
                    new_name = self.replace_two_pos_in_trimoieties(
                                    num1=2, num2=3, old_name=old_name, 
                                    info=moiety_info, 
                                    old_db_1=old_db_2, new_db_1=new_db_2,
                                    old_db_2=old_db_3, new_db_2=new_db_3)
                    set_resolved_bool(idx=idx, new_db=new_db_2, acyl_num=1)
                    set_resolved_bool(idx=idx, new_db=new_db_3, acyl_num=2)
                elif db_1_changed:
                    new_name = self.replace_one_pos_in_trimoieties(
                                    num=1, old_name=old_name, info=moiety_info, 
                                    old_db=old_db_1, new_db=new_db_1)
                    set_resolved_bool(idx=idx, new_db=new_db_1, acyl_num=0)
                elif db_2_changed:
                    new_name = self.replace_one_pos_in_trimoieties(
                                    num=2, old_name=old_name, info=moiety_info, 
                                    old_db=old_db_2, new_db=new_db_2)
                    set_resolved_bool(idx=idx, new_db=new_db_2, acyl_num=1)
                else: #db_3_changed
                    new_name = self.replace_one_pos_in_trimoieties(
                                    num=3, old_name=old_name, info=moiety_info, 
                                    old_db=old_db_3, new_db=new_db_3)
                    set_resolved_bool(idx=idx, new_db=new_db_3, acyl_num=2)
        else:
            new_name = old_name
        return new_name

    def replace_one_pos_in_moiety(self, old_name, info, old_db, new_db):
        chain, db = 'chain-1', 'db-1'
        moiety = f'{info[chain]}:{info[db]}'
        sign = 'one'
        signed_moiety = f'{info[chain]}:{sign}{info[db]}'
        if old_db == 'Unresolved':
            old = f'{signed_moiety}'
        else:
            old = f'{signed_moiety}({old_db})'
        if new_db == 'Unresolved':
            new = f'{signed_moiety}'
        else:
            new = f'{signed_moiety}({new_db})'
        old_name = old_name.replace(moiety, signed_moiety, 1)
        new_name = old_name.replace(old, new, 1)
        new_name = new_name.replace('one', '', 1)
        return new_name

    def replace_one_pos_in_dimoieties(self, num, old_name, info, old_db, new_db):
        chain_1, db_1 = 'chain-1', 'db-1'
        chain_2, db_2 = 'chain-2', 'db-2'
        moiety_1 = f'{info[chain_1]}:{info[db_1]}'
        moiety_2 = f'{info[chain_2]}:{info[db_2]}'
        signed_moiety_1 = f'{info[chain_1]}:one{info[db_1]}'
        signed_moiety_2 = f'{info[chain_2]}:two{info[db_2]}'
        if num == 1:
            if old_db == 'Unresolved':
                old = f'{signed_moiety_1}'
            else:
                old = f'{signed_moiety_1}({old_db})'
            if new_db == 'Unresolved':
                new = f'{signed_moiety_1}'
            else:
                new = f'{signed_moiety_1}({new_db})'
        else:
            if old_db == 'Unresolved':
                old = f'{signed_moiety_2}'
            else:
                old = f'{signed_moiety_2}({old_db})'
            if new_db == 'Unresolved':
                new = f'{signed_moiety_2}'
            else:
                new = f'{signed_moiety_2}({new_db})'
        old_name = old_name.replace(moiety_1, signed_moiety_1, 1)
        old_name = old_name.replace(moiety_2, signed_moiety_2, 1)
        new_name = old_name.replace(old, new, 1)
        new_name = new_name.replace('one', '', 1).replace('two', '', 1)
        return new_name

    def replace_one_pos_in_trimoieties(
        self, num, old_name, info, old_db, new_db):
        chain_1, db_1 = 'chain-1', 'db-1'
        chain_2, db_2 = 'chain-2', 'db-2'
        chain_3, db_3 = 'chain-3', 'db-3'
        moiety_1 = f'{info[chain_1]}:{info[db_1]}'
        moiety_2 = f'{info[chain_2]}:{info[db_2]}'
        moiety_3 = f'{info[chain_3]}:{info[db_3]}'
        signed_moiety_1 = f'{info[chain_1]}:one{info[db_1]}'
        signed_moiety_2 = f'{info[chain_2]}:two{info[db_2]}'
        signed_moiety_3 = f'{info[chain_3]}:three{info[db_3]}'
        if num == 1:
            if old_db == 'Unresolved':
                old = f'{signed_moiety_1}'
            else:
                old = f'{signed_moiety_1}({old_db})'
            if new_db == 'Unresolved':
                new = f'{signed_moiety_1}'
            else:
                new = f'{signed_moiety_1}({new_db})'
        elif num == 2:
            if old_db == 'Unresolved':
                old = f'{signed_moiety_2}'
            else:
                old = f'{signed_moiety_2}({old_db})'
            if new_db == 'Unresolved':
                new = f'{signed_moiety_2}'
            else:
                new = f'{signed_moiety_2}({new_db})'
        else:
            if old_db == 'Unresolved':
                old = f'{signed_moiety_3}'
            else:
                old = f'{signed_moiety_3}({old_db})'
            if new_db == 'Unresolved':
                new = f'{signed_moiety_3}'
            else:
                new = f'{signed_moiety_3}({new_db})'
        old_name = old_name.replace(moiety_1, signed_moiety_1, 1)
        old_name = old_name.replace(moiety_2, signed_moiety_2, 1)
        old_name = old_name.replace(moiety_3, signed_moiety_3, 1)
        new_name = old_name.replace(old, new, 1)
        new_name = new_name.replace('one','',1).replace('two','',1).replace('three','',1)
        return new_name

    def replace_two_pos_in_dimoieties(
        self, old_name, info, old_db_1, new_db_1, old_db_2, new_db_2):
        chain_1, db_1 = 'chain-1', 'db-1'
        chain_2, db_2 = 'chain-2', 'db-2'
        moiety_1 = f'{info[chain_1]}:{info[db_1]}'
        moiety_2 = f'{info[chain_2]}:{info[db_2]}'
        signed_moiety_1 = f'{info[chain_1]}:one{info[db_1]}'
        signed_moiety_2 = f'{info[chain_2]}:two{info[db_2]}'
        if old_db_1 == 'Unresolved':
            old_1 = f'{signed_moiety_1}'
        else:
            old_1 = f'{signed_moiety_1}({old_db_1})'
        if old_db_2 == 'Unresolved':
            old_2 = f'{signed_moiety_2}'
        else:
            old_2 = f'{signed_moiety_2}({old_db_2})'
        if new_db_1 == 'Unresolved':
            new_1 = f'{signed_moiety_1}'
        else:
            new_1 = f'{signed_moiety_1}({new_db_1})'
        if new_db_2 == 'Unresolved':
            new_2 = f'{signed_moiety_2}'
        else:
            new_2 = f'{signed_moiety_2}({new_db_2})'
        old_name = old_name.replace(moiety_1, signed_moiety_1, 1)
        old_name = old_name.replace(moiety_2, signed_moiety_2, 1)
        new_name = old_name.replace(old_1, new_1, 1)
        new_name = new_name.replace(old_2, new_2, 1)
        new_name = new_name.replace('one', '', 1).replace('two', '', 1)
        return new_name

    def replace_two_pos_in_trimoieties(self, num1, num2, old_name, 
        info, old_db_1, new_db_1, old_db_2, new_db_2):
        chain_1, db_1 = 'chain-1', 'db-1'
        chain_2, db_2 = 'chain-2', 'db-2'
        chain_3, db_3 = 'chain-3', 'db-3'
        moiety_1 = f'{info[chain_1]}:{info[db_1]}'
        moiety_2 = f'{info[chain_2]}:{info[db_2]}'
        moiety_3 = f'{info[chain_3]}:{info[db_3]}'
        signed_moiety_1 = f'{info[chain_1]}:one{info[db_1]}'
        signed_moiety_2 = f'{info[chain_2]}:two{info[db_2]}'
        signed_moiety_3 = f'{info[chain_3]}:three{info[db_3]}'
        if num1 == 1 and num2 == 2:
            if old_db_1 == 'Unresolved':
                old_1 = f'{signed_moiety_1}'
            else:
                old_1 = f'{signed_moiety_1}({old_db_1})'
            if new_db_1 == 'Unresolved':
                new_1 = f'{signed_moiety_1}'
            else:
                new_1 = f'{signed_moiety_1}({new_db_1})'
            if old_db_2 == 'Unresolved':
                old_2 = f'{signed_moiety_2}'
            else:
                old_2 = f'{signed_moiety_2}({old_db_2})'
            if new_db_2 == 'Unresolved':
                new_2 = f'{signed_moiety_2}'
            else:
                new_2 = f'{signed_moiety_2}({new_db_2})'
        elif num1 == 1 and num2 == 3:
            if old_db_1 == 'Unresolved':
                old_1 = f'{signed_moiety_1}'
            else:
                old_1 = f'{signed_moiety_1}({old_db_1})'
            if new_db_1 == 'Unresolved':
                new_1 = f'{signed_moiety_1}'
            else:
                new_1 = f'{signed_moiety_1}({new_db_1})'
            if old_db_2 == 'Unresolved':
                old_2 = f'{signed_moiety_3}'
            else:
                old_2 = f'{signed_moiety_3}({old_db_2})'
            if new_db_2 == 'Unresolved':
                new_2 = f'{signed_moiety_3}'
            else:
                new_2 = f'{signed_moiety_3}({new_db_2})'
        else: #num1 == 2 and num2 == 3:
            if old_db_1 == 'Unresolved':
                old_1 = f'{signed_moiety_2}'
            else:
                old_1 = f'{signed_moiety_2}({old_db_1})'
            if new_db_1 == 'Unresolved':
                new_1 = f'{signed_moiety_2}'
            else:
                new_1 = f'{signed_moiety_2}({new_db_1})'
            if old_db_2 == 'Unresolved':
                old_2 = f'{signed_moiety_3}'
            else:
                old_2 = f'{signed_moiety_3}({old_db_2})'
            if new_db_2 == 'Unresolved':
                new_2 = f'{signed_moiety_3}'
            else:
                new_2 = f'{signed_moiety_3}({new_db_2})'
        old_name = old_name.replace(moiety_1, signed_moiety_1, 1)
        old_name = old_name.replace(moiety_2, signed_moiety_2, 1)
        old_name = old_name.replace(moiety_3, signed_moiety_3, 1)
        new_name = old_name.replace(old_1, new_1, 1)
        new_name = new_name.replace(old_2, new_2, 1)
        new_name = new_name.replace('one','',1).replace('two','',1).replace('three','',1)
        return new_name

    def replace_three_pos_in_trimoieties(self, old_name, info, 
        old_db_1, new_db_1, old_db_2, new_db_2, old_db_3, new_db_3):
        chain_1, db_1 = 'chain-1', 'db-1'
        chain_2, db_2 = 'chain-2', 'db-2'
        chain_3, db_3 = 'chain-3', 'db-3'
        moiety_1 = f'{info[chain_1]}:{info[db_1]}'
        moiety_2 = f'{info[chain_2]}:{info[db_2]}'
        moiety_3 = f'{info[chain_3]}:{info[db_3]}'
        signed_moiety_1 = f'{info[chain_1]}:one{info[db_1]}'
        signed_moiety_2 = f'{info[chain_2]}:two{info[db_2]}'
        signed_moiety_3 = f'{info[chain_3]}:three{info[db_3]}'
        if old_db_1 == 'Unresolved':
            old_1 = f'{signed_moiety_1}'
        else:
            old_1 = f'{signed_moiety_1}({old_db_1})'
        if old_db_2 == 'Unresolved':
            old_2 = f'{signed_moiety_2}'
        else:
            old_2 = f'{signed_moiety_2}({old_db_2})'
        if old_db_3 == 'Unresolved':
            old_3 = f'{signed_moiety_3}'
        else:
            old_3 = f'{signed_moiety_3}({old_db_3})'
        if new_db_1 == 'Unresolved':
            new_1 = f'{signed_moiety_1}'
        else:
            new_1 = f'{signed_moiety_1}({new_db_1})'
        if new_db_2 == 'Unresolved':
            new_2 = f'{signed_moiety_2}'
        else:
            new_2 = f'{signed_moiety_2}({new_db_2})'
        if new_db_3 == 'Unresolved':
            new_3 = f'{signed_moiety_3}'
        else:
            new_3 = f'{signed_moiety_3}({new_db_3})'
        old_name = old_name.replace(moiety_1, signed_moiety_1, 1)
        old_name = old_name.replace(moiety_2, signed_moiety_2, 1)
        old_name = old_name.replace(moiety_3, signed_moiety_3, 1)
        new_name = old_name.replace(old_1, new_1, 1)
        new_name = new_name.replace(old_2, new_2, 1)
        new_name = new_name.replace(old_3, new_3, 1)
        new_name \
        = new_name.replace('one','',1).replace('two','',1).replace('three','',1)
        return new_name

    def re_const_oad_res_dict(self, idx, moiety_num, rank, old_db, new_db):
        unresolved_d = {'Positions': '', 'N-description': 'Unresolved',
                        'Score': '###', 'Ratio sum': '###', 'Presence': '###',
                        'Notice': 'Unresolved', 
                        'Measured peaks': [[0, 0, 0]], 
                        'Ref peaks': [['', 0, 0, 0]], 
                        'Peaks dict': {'none': [0, 0, 0, 0, 0]}
                        }
        moiety = f'Moiety-{moiety_num}'
        old_no1 = self.oad_result_dict[idx][moiety][0].copy()
        new_no1 = self.oad_result_dict[idx][moiety][rank].copy()
        moiety_d = self.oad_result_dict[idx][moiety].copy()
        def get_unresolved_rank(d):
            for key, v in d.items():
                if v['N-description'] == 'Unresolved':
                    return key
        if 'Modified' not in new_no1['Notice']:
            new_no1['Notice'] += ', Modified'
        new_d = {0: new_no1}
        if new_db == 'Unresolved':
            for del_key in [0, rank, 'Determined db']:
                    removed = moiety_d.pop(del_key)
            new_d[1] = old_no1
            for i, (rank, v) in enumerate(moiety_d.items(), start=2):
                new_d[i] = v
        else:
            if old_db == 'Unresolved':
                for del_key in [0, rank, 'Determined db']:
                    removed = moiety_d.pop(del_key)
                for i, (rank, v) in enumerate(moiety_d.items(), start=1):
                    new_d[i] = v
            else:
                unresolved_rank = get_unresolved_rank(moiety_d)
                for del_key in [0, rank, unresolved_rank, 'Determined db']:
                    removed = moiety_d.pop(del_key)
                new_d[1] = old_no1
                for i, (rank, v) in enumerate(moiety_d.items(), start=2):
                    new_d[i] = v
            unresolved_d['Notice'] += ', Modified'
            new_d[len(new_d)] = unresolved_d
        new_d['Determined db'] = new_no1
        self.oad_result_dict[idx][moiety] = new_d

    def update_result_table(self, target_idx):
        self.result_table_tree.delete(*self.result_table_tree.get_children())
        switch = self.show_only_solved_mol.get()
        ion_v_col = 'pmol/mg tissue' if self.normalized_data else 'Height'
        tag = 'white'
        if switch:
            count = 0
            for i, (row, df) in enumerate(self.target_table.iterrows()):
                # tree = ['ID', 'RT', 'm/z', 'm/z type', 'Solved','Result name', 'Height', 'Ontology', 'Comment']
                solved = df['Solved level']
                if solved == 'None': continue
                idx = df['ID']
                rt = math_floor(df['RT(min)'], 3)
                mz = df['Precise m/z']
                mz_type = df['Precise m/z type']
                result_name = df['OAD result name']
                data_from = df['Data from']
                ion_v = df[ion_v_col]
                ontology = df['Ontology']
                comment = df['User comment']
                tag = 'colored' if tag == 'white' else 'white'
                count += 1
                iid = self.result_table_tree.insert("", "end", tags=tag,
                        values=[idx, rt, mz, mz_type, solved, result_name, 
                                data_from, ion_v, ontology, comment])
                if idx == target_idx:
                    selected_iid = iid
            new_label = f'Result table (molecules={count})'
        else:
            for i, (row, df) in enumerate(self.target_table.iterrows()):
                # tree = ['ID', 'RT', 'm/z', 'm/z type', 'Solved','Result name', 'Height', 'Ontology', 'Comment']
                idx = df['ID']
                rt = math_floor(df['RT(min)'], 3)
                mz = df['Precise m/z']
                mz_type = df['Precise m/z type']
                solved = df['Solved level']
                result_name = df['OAD result name']
                data_from = df['Data from']
                ion_v = df[ion_v_col]
                ontology = df['Ontology']
                comment = df['User comment']
                tag = 'colored' if tag == 'white' else 'white'
                iid = self.result_table_tree.insert("", "end", tags=tag,
                        values=[idx, rt, mz, mz_type, solved, result_name, 
                                data_from, ion_v, ontology, comment])
                if idx == target_idx:
                    selected_iid = iid
            new_label = f'Result table (molecules={i+1})'
        self.left_wrap_label_change(new_label)
        self.result_table_tree.tag_configure('colored', background="#d3d3d3")
        try:
            self.result_table_tree.focus(selected_iid)
            self.result_table_tree.selection_set(selected_iid)
        except tk.TclError:
            first_iid = self.result_table_tree.get_children()[0]
            self.result_table_tree.focus(first_iid)
            self.result_table_tree.selection_set(first_iid)
        self.result_table_tree.bind("<<TreeviewSelect>>", self.data_setting)

    def enter_comment_on_result_table(self):
        entry_comment = self.comment_entry.get()
        if entry_comment != "":
            idx = self.selected_idx
            df = self.target_table[self.target_table['ID'] == idx]
            target_row = list(df.index)[0]
            previous_comment = df['User comment'].values[0]
            new_comment = previous_comment + '| ' + entry_comment
            self.target_table.loc[target_row:target_row, ['User comment']] = new_comment
            self.comment_str.set("")
            self.update_result_table(idx)

    def update_msms_fig(self, cid_info):
        #region CID result dict structure
        # cid_dict = {'Lipid subclass': {'Glycine': [ref_mz, measured_mz, intensity, ratio, ppm],
        #                                'Presence': 50.00},
        #             'Moiety':         {'acyl-2': [ref_mz, measured_mz, intensity, ratio, ppm],
        #                                'Presence': 50.00}
        #            }
        #endregion
        #region base settings
        msms_df = self.temp_msms_df
        lipid_info = self.temp_lipid_info
        dpi_setting = 150
        fig_x, fig_y = 3.6, 2.4
        y_max = 100
        y_min = 0
        precursor_mz_font_size = 10
        frag_mz_font_size = 5
        x_label_font_size, y_label_font_size = 10, 10
        small_font_size = 8
        middle_font_size = 10
        line_color = 'black'
        cid_color = 'red'
        xtick_num = 10
        axes = ['top', 'bottom', 'left', 'right']
        axes_width = [0, 1, 1, 0]
        tick_width = 1
        tick_pad = 1
        total_ions = msms_df['frag m/z'].values.tolist()
        min_mz = min(total_ions)
        ex_msms = msms_df[msms_df['Ratio(%)'] >= 1]
        total_fig_x = ex_msms['frag m/z'].values.tolist()
        total_fig_y = ex_msms['Ratio(%)'].values.tolist()
        # total_fig_x = msms_df['frag m/z'].values.tolist()
        # total_fig_y = msms_df['Ratio(%)'].values.tolist()
        if lipid_info['MS2 Mz'] > 0:
            precursor_mz = lipid_info['MS2 Mz']
        else:
            precursor_mz = lipid_info['Precursor Mz']
        x_max = math.ceil(precursor_mz/xtick_num)*xtick_num+5
        x_min = math.floor(min_mz/xtick_num)*xtick_num
        # x ragne m/z 100-500 -> bar_width=2
        bar_width = math.floor(10*(x_max-x_min)/200)/10
        #endregion
        #region matplotlib base setting
        fig = Figure(figsize=(fig_x, fig_y), dpi=dpi_setting)
        fig.patch.set_facecolor('white')
        fig.patch.set_alpha(0)
        ax = fig.add_subplot(1,1,1)
        ax.patch.set_facecolor('white')
        ax.patch.set_alpha(0)
        ax.set_xlim([x_min, x_max])
        ax.set_ylim([y_min, y_max])
        for axis, width in zip(axes, axes_width):
            ax.spines[axis].set_linewidth(width)
        ax.tick_params(labelsize=small_font_size, width=tick_width, pad=tick_pad)
        ax.set_yticks([0, 50, 100])
        ax.set_yticklabels(['0', '50', '100'])
        #endregion
        #region Plotting measured MS/MS
        ax.bar(total_fig_x, total_fig_y, color=line_color, 
            width=bar_width, linewidth=0, zorder=1)
        ax.text(precursor_mz, 100, str(precursor_mz), 
            fontsize=precursor_mz_font_size, ha='center', va='bottom')
        def check_near_txt(others, mz, ratio):
            space = bar_width if ratio > 10 else 0
            for other in others:
                if other <= mz+25:
                    space += bar_width*4
            if mz-20 <= min_mz:
                space = -10
            return space
        if len(cid_info['Lipid subclass']) > 1:
            x_list, y_list, info_list, others = [], [], [], []
            for key, v in cid_info['Lipid subclass'].items():
                if key != 'Presence' and v[1] > 0:
                    x_list.append(v[1]), y_list.append(v[3])
                    info_list.append([key, v[1], v[2], v[3], v[4]])
                    #['Glycine', measured_mz, intensity, ratio, ppm]
            ax.bar(x_list, y_list, color=cid_color, width=bar_width, 
                linewidth=0, zorder=2)
            for li in info_list:
                annotation = 'm/z {}, ppm={}, {}'.format(li[1], li[4], li[0])
                space = check_near_txt(others, li[1], li[3])
                ax.text(li[1]-space, 10, annotation, fontsize=frag_mz_font_size, 
                color=cid_color, va='baseline', rotation=90, rotation_mode='anchor')
                others.append(li[1])
        if len(cid_info['Moiety']) > 1:
            x_list, y_list, info_list, others = [], [], [], []
            for key, v in cid_info['Moiety'].items():
                if key != 'Presence' and v[1] > 0:
                    x_list.append(v[1]), y_list.append(v[3])
                    info_list.append([key, v[1], v[2], v[3], v[4]])
                    #['Glycine', measured_mz, intensity, ratio, ppm]
            ax.bar(x_list, y_list, color=cid_color, width=bar_width, 
                linewidth=0, zorder=2)
            for li in info_list:
                annotation = 'm/z {}, ppm={}, {}'.format(li[1], li[4], li[0])
                space = check_near_txt(others, li[1], li[3])
                ax.text(li[1]-space, 10, annotation, fontsize=frag_mz_font_size, 
                color=cid_color, va='baseline', rotation=90, rotation_mode='anchor')
                others.append(li[1])
        #endregion
        #region matplotlib regend setting
        msms_label_measured = '$\it{Measurement}$'
        ax.text(x_min+5, 100, msms_label_measured, fontsize=small_font_size)
        y_label = '$\it{Relative}$' + ' ' + '$\it{abundance}$'
        ax.set_xlabel('$\it{m/z}$', fontsize=x_label_font_size, labelpad=tick_pad)
        ax.set_ylabel(y_label, fontsize=y_label_font_size, labelpad=tick_pad)
        # fig.tight_layout()
        self.centoird_msms_canvas = FigureCanvasTkAgg(fig, master=self.centroid_msms_frame)
        self.centoird_msms_canvas.get_tk_widget().place(relx=0, rely=0, relwidth=1, relheight=1)
        #endregion

    def update_oad_msms_fig_with_ref(self, dict_1, dict_2, dict_3):
        msms_df = self.temp_msms_df
        lipid_info = self.temp_lipid_info
        graph_info = self.temp_graph_info['OAD']
        #region base settings
        dpi_setting = 150
        fig_x, fig_y = 4.8, 2.4
        precursor_mz_font_size = 10
        frag_mz_font_size = 6
        x_label_font_size, y_label_font_size = 10, 10
        small_font_size = 8
        middle_font_size = 10
        line_color = 'black'
        moiety_1_color = 'red'
        moiety_2_color = 'blue'
        moiety_3_color = 'green'
        xtick_num = 10
        axes = ['top', 'bottom', 'left', 'right']
        axes_width = [0, 1, 1, 0]
        tick_width = 1
        tick_pad = 1
        #endregion
        #region Data structure
        # dict = {0:              {'Positions': '', 'N-description': '', 'Score': float, 
        #                          'Ratio sum': float, 'Presence': float, 'Notice': '',
        #                          'Peaks dict': {'n-9/dis@n-8/+O/': [Ref m/z, Ref delta, Measured m/z, Measured ratio, ppm]}
        #                         },
        #         1:              {....},
        #         Determined db:  {'Positions': '', 'N-description': '', 'Score': float,
        #                          'Ratio sum': float, 'Presence': float, 'Notice': '',
        #                          'Measured peaks': [[Measured m/z, Measured ratio, ppm], [...]],
        #                          'Ref peaks': [[OAD type, Ref m/z, Ref NL, Ref ratio], [...]]
        #                         }
        #         }
        # 'Peaks dict': {'n-9/dis@n-8/+O/': [Ref m/z, Ref delta, Measured m/z, Measured ratio, ppm]}
        #endregion
        #region Detail setting
        positions_list = []
        # peaks_dict_1, peaks_dict_2, peaks_dict_3 = {}, {}, {}
        ref_mz_1, ref_ratio_1, ref_mz_2, ref_ratio_2, ref_mz_3, ref_ratio_3 \
        = [], [], [], [], [], []
        # measured_oad_ions_1, measured_oad_ions_2, measured_oad_ions_3 = [], [], []
        # oad04_ions_1, oad04_ions_2, oad04_ions_3 = [], [], []
        # oad10_ions_1, oad10_ions_2, oad10_ions_3 = [], [], []
        is_this_plasmalogen = 'Plasm' in lipid_info['Ontology']
        # ref_precursor_mz = lipid_info['Ref precursor Mz']
        if dict_3:
            positions_list.append(dict_3[0]['Positions'])
            # peaks_dict_3 = dict_3[0]['Peaks dict']
        elif dict_2:
            positions_list.append(dict_2[0]['Positions'])
            # peaks_dict_2 = dict_2[0]['Peaks dict']
        elif dict_1:
            positions_list.append(dict_1[0]['Positions'])
            # peaks_dict_1 = dict_1[0]['Peaks dict']
        if dict_3:
            ref_mz_3 = [li[1] for li in dict_3[0]['Ref peaks']]
            ref_ratio_3 = [li[3]*(-1) for li in dict_3[0]['Ref peaks']]
            # measured_oad_ions_3 = [li[1] for li in dict_3[0]['Measured peaks']]
            # oad04_ions_3 = [[li[1], li[3]] for li in dict_3[0]['Ref peaks'] if 'OAD04' in li[0]]
            # oad10_ions_3 = [[li[1], li[3]] for li in dict_3[0]['Ref peaks'] if 'OAD10' in li[0]]
        if dict_2:
            ref_mz_2 = [li[1] for li in dict_2[0]['Ref peaks']]
            ref_ratio_2 = [li[3]*(-1) for li in dict_2[0]['Ref peaks']]
            # measured_oad_ions_2 = [li[1] for li in dict_2[0]['Measured peaks']]
            # oad04_ions_2 = [[li[1], li[3]] for li in dict_2[0]['Ref peaks'] if 'OAD04' in li[0]]
            # oad10_ions_2 = [[li[1], li[3]] for li in dict_2[0]['Ref peaks'] if 'OAD10' in li[0]]
        if dict_1:
            ref_mz_1 = [li[1] for li in dict_1[0]['Ref peaks']]
            ref_ratio_1 = [li[3]*(-1) for li in dict_1[0]['Ref peaks']]
            # measured_oad_ions_1 = [li[1] for li in dict_1[0]['Measured peaks']]
            # oad04_ions_1 = [[li[1], li[3]] for li in dict_1[0]['Ref peaks'] if 'OAD04' in li[0]]
            # oad10_ions_1 = [[li[1], li[3]] for li in dict_1[0]['Ref peaks'] if 'OAD10' in li[0]]
        precursor_mz, ref_precursor_mz = graph_info['MS2 Mz'], graph_info['Ref precursor Mz']
        x_min, x_max = graph_info['x-range'][0], graph_info['x-range'][1]
        bar_width = graph_info['Bar_width']
        y_max = 100
        y_min = -100
        #endregion
        #region magnification setting
        final_magnification = graph_info['Magnification']
        #endregion
        #region matplotlib base setting
        fig = Figure(figsize=(fig_x, fig_y), dpi=dpi_setting)
        fig.patch.set_facecolor('white')
        fig.patch.set_alpha(0)
        ax = fig.add_subplot(1,1,1)
        ax.patch.set_facecolor('white')
        ax.patch.set_alpha(0)
        ax.set_xlim([x_min, x_max])
        ax.set_ylim([y_min, y_max])
        for axis, width in zip(axes, axes_width):
            ax.spines[axis].set_linewidth(width)
        ax.tick_params(labelsize=small_font_size, width=tick_width, pad=tick_pad)
        ax.set_yticks([-100, -50, 0, 50, 100])
        ax.set_yticklabels(['100', '50', '0', '50', '100'])
        # ax.spines['right'].set_visible(False)
        # ax.spines['top'].set_visible(False)
        #endregion
        #region Plotting OAD-MS/MS
        ex_msms = msms_df[msms_df['frag m/z'] >= x_min]
        total_fig_x = ex_msms['frag m/z'].values.tolist()
        total_fig_y = ex_msms['Ratio(%)'].values.tolist()
        # total_fig_x = msms_df['frag m/z'].values.tolist()
        # total_fig_y = msms_df['Ratio(%)'].values.tolist()
        total_fig_y = [v*final_magnification for v in total_fig_y]
        ref_ratio_1 = [v*final_magnification for v in ref_ratio_1]
        ref_ratio_2 = [v*final_magnification for v in ref_ratio_2]
        ref_ratio_3 = [v*final_magnification for v in ref_ratio_3]
        ax.bar(total_fig_x, total_fig_y, color=line_color, width=bar_width, 
               linewidth=0)
        ax.text(precursor_mz, 101, str(precursor_mz), 
            fontsize=precursor_mz_font_size, ha='center', va='bottom')
        ax.bar(ref_precursor_mz, -99, color=moiety_1_color, width=bar_width, 
               linewidth=0)
        ax.text(ref_precursor_mz, -125, str(ref_precursor_mz), 
            fontsize=precursor_mz_font_size, ha='center', va='bottom', 
            color=moiety_1_color)
        if ref_mz_1:
            ax.bar(ref_mz_1, ref_ratio_1, color=moiety_1_color, 
                width=bar_width, linewidth=0)
        if ref_mz_2:
            ax.bar(ref_mz_2, ref_ratio_2, color=moiety_2_color, 
                width=bar_width, linewidth=0)
        if ref_mz_3:
            ax.bar(ref_mz_3, ref_ratio_3, color=moiety_3_color, 
                width=bar_width, linewidth=0)
        #endregion
        #region matplotlib regend setting
        magnified_mz = round(precursor_mz-10)
        ax.hlines(y=100, xmin=(x_min+2), xmax=magnified_mz, colors='gray', linestyle='dashed')
        ax.hlines(y=0, xmin=x_min, xmax=x_max, linewidths=0.5)
        hlines_label = '×{}'.format(final_magnification)
        label_position_x = statistics.mean([x_min+5, magnified_mz])
        ax.text(label_position_x, 102, hlines_label, fontsize=small_font_size)
        msms_label_measured = '$\it{Measurement}$'
        msms_label_ref = '$\it{Reference}$'
        ax.text(x_min+1, 102, msms_label_measured, fontsize=small_font_size)
        ax.text(x_min+1, -98, msms_label_ref, fontsize=small_font_size)
        y_label = '$\it{Relative}$' + ' ' + '$\it{abundance}$'
        ax.set_xlabel('$\it{m/z}$', fontsize=x_label_font_size, labelpad=tick_pad)
        ax.set_ylabel(y_label, fontsize=y_label_font_size, labelpad=tick_pad)
        # fig.tight_layout()
        self.measured_vs_ref_msms_canvas = FigureCanvasTkAgg(fig, master=self.measured_vs_ref_msms_frame)
        self.measured_vs_ref_msms_canvas.get_tk_widget().place(relx=0, rely=0, relwidth=1, relheight=1)
        #endregion

    def matplotlib_test1(self, parent):
        values = [20, 34, 30, 35, 27]
        x_range = range(0, len(values))
        fig = Figure(figsize=(3.6, 2.4), dpi=150)
        ax = fig.add_subplot(111)
        ax.bar(x_range, values, 0.5)
        fig.tight_layout()
        self.centoird_msms_canvas = FigureCanvasTkAgg(fig, master=parent)
        self.centoird_msms_canvas.get_tk_widget().place(relx=0, rely=0, relwidth=1, relheight=1)

    def matplotlib_test2(self, parent):
        values = [20, 34, 30, 35, 27]
        x_range = range(0, len(values))
        fig = Figure(figsize=(3.6, 2.4), dpi=150)
        ax = fig.add_subplot(111)
        ax.bar(x_range, values, 0.5)
        fig.tight_layout()
        self.measured_vs_ref_msms_canvas = FigureCanvasTkAgg(fig, master=parent)
        self.measured_vs_ref_msms_canvas.get_tk_widget().place(relx=0, rely=0, relwidth=1, relheight=1)

    #endregion

    #region Re-analysis
    def set_re_analysis_condition(self):
        focused_item = self.get_focused_item()
        metabolite_name = focused_item[5]
        self.re_analysis_win = ReAnalysisWindow(
            master=self.master, name=metabolite_name, info=self.temp_lipid_info)
        self.re_analysis_win.start_btn["command"] \
            = self.check_re_analysis_condition
        self.re_analysis_win.close_btn["command"] \
            = self.re_analysis_win.close_window

    def check_re_analysis_condition(self):
        content1 = self.re_analysis_win.text1.get('1.0', 'end -1c')
        content2 = self.re_analysis_win.text2.get('1.0', 'end -1c')
        content3 = self.re_analysis_win.text3.get('1.0', 'end -1c')
        dbs_1 = self.extract_dbs(content1)
        dbs_2 = self.extract_dbs(content2)
        dbs_3 = self.extract_dbs(content3)
        moiety_info = self.temp_lipid_info['Each moiety info']
        unsaturated_num = self.temp_lipid_info['Unsaturated moiety']
        def check_duplicates(new_dbs, oad_result_d, moiety_num):
            moiety = f'Moiety-{moiety_num}'
            old_dbs = [
                v['Positions'] for rank, v in oad_result_d[moiety].items()]
            dup_bools = []
            for db in new_dbs:
                if db in old_dbs:
                    dup_bools.append(True)
                else:
                    dup_bools.append(False)
            return any(dup_bools)
        def check_db_num_matched(new_dbs, db_num):
            unmatched_bools = []
            for db in new_dbs:
                unmatched_bools.append(len(db) != db_num)
            return any(unmatched_bools)
        def check_positions(new_dbs, c_num):
            positions_bools = []
            for db in new_dbs:
                for pos in db:
                    sub = c_num - pos
                    b = (sub < 2) or (pos < 3)
                    positions_bools.append(b)
            if len(new_dbs[0]) > 1:
                for db in new_dbs:
                    for i in range(len(db)-1):
                        sub = db[i+1] - db[i]
                        positions_bools.append(sub < 2)
            return any(positions_bools)
        def show_dup_error_message():
            messagebox.showerror(
                "Duplicate error",
                "Re-analysis candidate already exist.")
        def show_unmatched_error_message():
            messagebox.showerror(
                "Unmatched error", 
                "Number of double bonds is not matched.")
        def show_position_error_message():
            messagebox.showerror(
                "C=C position error", 
                "C=C position is not appropriate.")
        if any([dbs_1, dbs_2, dbs_3]):
            idx = int(self.get_focused_item()[0])
            temp_oad_result_d = self.oad_result_dict[idx]
            dup1, dup2, dup3 = False, False, False
            unmatch1, unmatch2, unmatch3 = False, False, False
            pos1, pos2, pos3 = False, False, False
            if len(moiety_info) == 2:
                if dbs_1:
                    c_num = moiety_info['chain-1']
                    db_num = moiety_info['db-1']
                    dup1 = check_duplicates(new_dbs=dbs_1, 
                        oad_result_d=temp_oad_result_d, moiety_num=1)
                    unmatch1 = check_db_num_matched(
                        new_dbs=dbs_1, db_num=db_num)
                    pos1 = check_positions(new_dbs=dbs_1, c_num=c_num)
            elif len(moiety_info) == 4:
                if unsaturated_num == 1:
                    if dbs_1:
                        c_num = moiety_info['chain-1']
                        db_num = moiety_info['db-1']
                        dup1 = check_duplicates(new_dbs=dbs_1, 
                            oad_result_d=temp_oad_result_d, moiety_num=1)
                        unmatch1 = check_db_num_matched(
                            new_dbs=dbs_1, db_num=db_num)
                        pos1 = check_positions(new_dbs=dbs_1, c_num=c_num)
                    if dbs_2:
                        c_num = moiety_info['chain-2']
                        db_num = moiety_info['db-2']
                        dup2 = check_duplicates(new_dbs=dbs_2, 
                            oad_result_d=temp_oad_result_d, moiety_num=1)
                        unmatch2 = check_db_num_matched(
                            new_dbs=dbs_2, db_num=db_num)
                        pos2 = check_positions(new_dbs=dbs_2, c_num=c_num)
                elif unsaturated_num == 2:
                    if dbs_1:
                        c_num = moiety_info['chain-1']
                        db_num = moiety_info['db-1']
                        dup1 = check_duplicates(new_dbs=dbs_1, 
                            oad_result_d=temp_oad_result_d, moiety_num=1)
                        unmatch1 = check_db_num_matched(
                            new_dbs=dbs_1, db_num=db_num)
                        pos1 = check_positions(new_dbs=dbs_1, c_num=c_num)
                    if dbs_2:
                        c_num = moiety_info['chain-2']
                        db_num = moiety_info['db-2']
                        dup2 = check_duplicates(new_dbs=dbs_2, 
                            oad_result_d=temp_oad_result_d, moiety_num=2) 
                        unmatch2 = check_db_num_matched(
                            new_dbs=dbs_2, db_num=db_num)
                        pos2 = check_positions(new_dbs=dbs_2, c_num=c_num)
            elif len(moiety_info) == 6:
                if unsaturated_num == 1:
                    if dbs_1:
                        # if moiety_info['db-1'] > 0:
                        #     c_num = moiety_info['chain-1']
                        #     num, db_num = 1, moiety_info['db-1']
                        # elif moiety_info['db-2'] > 0:
                        #     c_num = moiety_info['chain-2']
                        #     num, db_num = 2, moiety_info['db-2']
                        # elif moiety_info['db-3'] > 0:
                        #     c_num = moiety_info['chain-3']
                        #     num, db_num = 3, moiety_info['db-3']
                        c_num = moiety_info['chain-1']
                        db_num = moiety_info['db-1']
                        dup1 = check_duplicates(new_dbs=dbs_1, 
                            oad_result_d=temp_oad_result_d, moiety_num=1)
                        unmatch1 = check_db_num_matched(
                            new_dbs=dbs_1, db_num=db_num)
                        pos1 = check_positions(new_dbs=dbs_1, c_num=c_num)
                    if dbs_2:
                        c_num = moiety_info['chain-2']
                        db_num = moiety_info['db-2']
                        dup2 = check_duplicates(new_dbs=dbs_2, 
                            oad_result_d=temp_oad_result_d, moiety_num=1)
                        unmatch2 = check_db_num_matched(
                            new_dbs=dbs_2, db_num=db_num)
                        pos2 = check_positions(new_dbs=dbs_2, c_num=c_num)
                    if dbs_3:
                        c_num = moiety_info['chain-3']
                        db_num = moiety_info['db-3']
                        dup3 = check_duplicates(new_dbs=dbs_3, 
                            oad_result_d=temp_oad_result_d, moiety_num=1)
                        unmatch3 = check_db_num_matched(
                            new_dbs=dbs_3, db_num=db_num)
                        pos3 = check_positions(new_dbs=dbs_3, c_num=c_num)
                elif unsaturated_num == 2:
                    if dbs_1:
                        c_num = moiety_info['chain-1']
                        db_num = moiety_info['db-1']
                        # if moiety_info['db-1'] > 0 and moiety_info['db-2'] > 0:
                        #     c_num = moiety_info['chain-1']
                        #     num, db_num = 1, moiety_info['db-1']
                        # elif moiety_info['db-1'] > 0 and moiety_info['db-3'] > 0:
                        #     c_num = moiety_info['chain-1']
                        #     num, db_num = 1, moiety_info['db-1']
                        # elif moiety_info['db-2'] > 0 and moiety_info['db-3'] > 0:
                        #     c_num = moiety_info['chain-2']
                        #     num, db_num = 2, moiety_info['db-2']
                        dup1 = check_duplicates(new_dbs=dbs_1, 
                            oad_result_d=temp_oad_result_d, moiety_num=1)
                        unmatch1 = check_db_num_matched(
                            new_dbs=dbs_1, db_num=db_num)
                        pos1 = check_positions(new_dbs=dbs_1, c_num=c_num)
                    if dbs_2:
                        # c_num = moiety_info['chain-2']
                        # num, db_num = 2, moiety_info['db-2']
                        if moiety_info['db-1'] > 0 and moiety_info['db-2'] > 0:
                            c_num = moiety_info['chain-2']
                            num, db_num = 2, moiety_info['db-2']
                        # elif moiety_info['db-1'] > 0 and moiety_info['db-3'] > 0:
                        #     c_num = moiety_info['chain-3']
                        #     num, db_num = 3, moiety_info['db-3']
                        elif moiety_info['db-2'] > 0 and moiety_info['db-3'] > 0:
                            c_num = moiety_info['chain-2']
                            num, db_num = 1, moiety_info['db-2']
                        dup2 = check_duplicates(new_dbs=dbs_2, 
                            oad_result_d=temp_oad_result_d, moiety_num=num)
                        unmatch2 = check_db_num_matched(
                            new_dbs=dbs_2, db_num=db_num)
                        pos2 = check_positions(new_dbs=dbs_2, c_num=c_num)
                    if dbs_3:
                        c_num = moiety_info['chain-3']
                        db_num = moiety_info['db-3']
                        dup3 = check_duplicates(new_dbs=dbs_3, 
                            oad_result_d=temp_oad_result_d, moiety_num=2)
                        unmatch3 = check_db_num_matched(
                            new_dbs=dbs_3, db_num=db_num)
                        pos3 = check_positions(new_dbs=dbs_3, c_num=c_num)
                elif unsaturated_num == 3:
                    if dbs_1:
                        c_num = moiety_info['chain-1']
                        db_num = moiety_info['db-1']
                        dup1 = check_duplicates(new_dbs=dbs_1, 
                            oad_result_d=temp_oad_result_d, moiety_num=1)
                        unmatch1 = check_db_num_matched(
                            new_dbs=dbs_1, db_num=db_num)
                        pos1 = check_positions(new_dbs=dbs_1, c_num=c_num)
                    if dbs_2:
                        c_num = moiety_info['chain-2']
                        db_num = moiety_info['db-2']
                        dup2 = check_duplicates(new_dbs=dbs_2, 
                            oad_result_d=temp_oad_result_d, moiety_num=2)
                        unmatch2 = check_db_num_matched(
                            new_dbs=dbs_2, db_num=db_num)
                        pos2 = check_positions(new_dbs=dbs_2, c_num=c_num)
                    if dbs_3:
                        c_num = moiety_info['chain-3']
                        db_num = moiety_info['db-3']
                        dup3 = check_duplicates(new_dbs=dbs_3, 
                            oad_result_d=temp_oad_result_d, moiety_num=3)
                        unmatch3 = check_db_num_matched(
                            new_dbs=dbs_3, db_num=db_num)
                        pos3 = check_positions(new_dbs=dbs_3, c_num=c_num)
            if any([dup1, dup2, dup3]):
                show_dup_error_message()
            elif any([unmatch1, unmatch2, unmatch3]):
                show_unmatched_error_message()
            elif any([pos1, pos2, pos3]):
                show_position_error_message()
            else:
                self.re_analysis_win.close_window()
                self.start_re_analysis(dbs_1, dbs_2, dbs_3)
        else:
            messagebox.showerror(
                "Empty error", """Please set C=C positional candidates.""")

    def extract_dbs(self, text):
        if text:
            comb_list = text.split('\n')
            dbs_list = [re.findall(r'\d+', comb) for comb in comb_list]
            dbs_list = [tuple(map(int, comb)) for comb in dbs_list]
            dbs_list = list(dict.fromkeys(dbs_list))
            return dbs_list
        else:
            return []

    def start_re_analysis(self, dbs_1, dbs_2, dbs_3):
        thread_1 = threading.Thread(
            target=self.re_analysis, args=(dbs_1, dbs_2, dbs_3))
        thread_2 = threading.Thread(
            target=self.create_re_analysis_popup_win)
        thread_1.start()
        thread_2.start()

    def create_re_analysis_popup_win(self):
        self.temp_popup = PopUpWindow(master=self.master,
            title='Re-analysis in progress', message='Re-analizing ...')

    def create_re_analysis_prgbar_window(self):
        self.temp_prgbar = ProgressBar(self.master,
            title="Re-analysis in progress", detail="--- Re analizing ---")

    def re_analysis(self, dbs_1, dbs_2, dbs_3):
        ontology = self.temp_lipid_info['Ontology']
        ref_mz = self.temp_lipid_info['Ref precursor Mz']
        ms_tol = math_floor(
            self.ms_tolerance_ppm*ref_mz/(1000*1000), 6)
        db_in_SPB = self.temp_lipid_info['Unsaturated sphingobase']
        deuterium = 0
        idx = int(self.get_focused_item()[0])
        if dbs_1:
            c_num = self.temp_lipid_info['Each moiety info']['chain-1']
            db_num = self.temp_lipid_info['Each moiety info']['db-1']
            sph_set = [db_in_SPB, c_num]
            for comb in dbs_1:
                ref_oad_d_1 = ReAnalyzer.generate_ref_oad_nl_and_type(
                    each_comb=comb, ontology=ontology, deuterium=deuterium)
                re_anal_bool_1 = ReAnalyzer.query_essential_diagnostic_ions(
                    df=self.temp_msms_df, ref_oad_dict=ref_oad_d_1, 
                    db_in_SPB=db_in_SPB, c_num=c_num, db_num=db_num,
                    tolerance=ms_tol,
                    must_nl_cut_off_dict=self.must_nl_cut_off_dict,
                    structure_dict=self.temp_lipid_info)
                n_description_1 = self.check_ion_lacking_pos(re_anal_bool_1)
                score_dict_1 = ReAnalyzer.calc_presence_ratios_and_score(
                    ref_oad_dict=ref_oad_d_1, cut_df=self.temp_msms_df,
                    ref_precursor_mz=ref_mz, ms_tol_ppm=self.ms_tolerance_ppm,
                    sph_set=sph_set)
                score_dict_1[0]['N-description'] = n_description_1
                key1 = self.set_re_analysis_result(
                    re_analysis_moiety=1, idx=idx, score_d=score_dict_1)
            # self.re_arrange_oad_dict(idx=idx, key=key1)
        if dbs_2:
            c_num = self.temp_lipid_info['Each moiety info']['chain-2']
            db_num = self.temp_lipid_info['Each moiety info']['db-2']
            sph_set = [False, c_num]
            for comb in dbs_2:
                ref_oad_d_2 = ReAnalyzer.generate_ref_oad_nl_and_type(
                    each_comb=comb, ontology=ontology, deuterium=deuterium)
                re_anal_bool_2 = ReAnalyzer.query_essential_diagnostic_ions(
                    df=self.temp_msms_df, ref_oad_dict=ref_oad_d_2, 
                    db_in_SPB=db_in_SPB, c_num=c_num, db_num=db_num,
                    tolerance=ms_tol,
                    must_nl_cut_off_dict=self.must_nl_cut_off_dict,
                    structure_dict=self.temp_lipid_info)
                n_description_2 = self.check_ion_lacking_pos(re_anal_bool_2)
                score_dict_2 = ReAnalyzer.calc_presence_ratios_and_score(
                    ref_oad_dict=ref_oad_d_2, cut_df=self.temp_msms_df,
                    ref_precursor_mz=ref_mz, ms_tol_ppm=self.ms_tolerance_ppm,
                    sph_set=sph_set)
                score_dict_2[0]['N-description'] = n_description_2
                key2= self.set_re_analysis_result(
                    re_analysis_moiety=2, idx=idx, score_d=score_dict_2)
            # self.re_arrange_oad_dict(idx=idx, key=key2)
        if dbs_3:
            c_num = self.temp_lipid_info['Each moiety info']['chain-3']
            db_num = self.temp_lipid_info['Each moiety info']['db-3']
            sph_set = [False, c_num]
            for comb in dbs_3:
                ref_oad_d_3 = ReAnalyzer.generate_ref_oad_nl_and_type(
                    each_comb=comb, ontology=ontology, deuterium=deuterium)
                re_anal_bool_3 = ReAnalyzer.query_essential_diagnostic_ions(
                    df=self.temp_msms_df, ref_oad_dict=ref_oad_d_3, 
                    db_in_SPB=db_in_SPB, c_num=c_num, db_num=db_num,
                    tolerance=ms_tol,
                    must_nl_cut_off_dict=self.must_nl_cut_off_dict,
                    structure_dict=self.temp_lipid_info)
                n_description_3 = self.check_ion_lacking_pos(re_anal_bool_3)
                score_dict_3 = ReAnalyzer.calc_presence_ratios_and_score(
                    ref_oad_dict=ref_oad_d_3, cut_df=self.temp_msms_df,
                    ref_precursor_mz=ref_mz, ms_tol_ppm=self.ms_tolerance_ppm,
                    sph_set=sph_set)
                score_dict_3[0]['N-description'] = n_description_3
                key3 = self.set_re_analysis_result(
                    re_analysis_moiety=3, idx=idx, score_d=score_dict_3)
            # self.re_arrange_oad_dict(idx=idx, key=key3)
        auto_magnification = self.graph_dict[idx]['OAD']['Magnification']
        self.graph_dict[idx]['OAD'] = ReAnalyzer.set_oad_graph_dict_value(
            oad_dict=self.oad_result_dict[idx], lipid_info=self.temp_lipid_info,
            auto_magnification=auto_magnification)
        self.temp_graph_info = self.graph_dict[idx]
        self.update_db_trees(
            result_dict=self.oad_result_dict[idx], masstable=False)
        self.focus_and_select_db_trees()
        self.temp_popup.close_window()

    def re_arrange_oad_dict(self, idx, key):
        unresolved_d = {'Positions': '', 'N-description': 'Unresolved',
                        'Score': '###', 'Ratio sum': '###', 'Presence': '###',
                        'Notice': 'Unresolved', 
                        'Measured peaks': [[0, 0, 0]], 
                        'Ref peaks': [['', 0, 0, 0]], 
                        'Peaks dict': {'none': [0, 0, 0, 0, 0]}
                        }
        moiety_d = self.oad_result_dict[idx][key].copy()
        removed = moiety_d.pop('Determined db')
        temp_d = {i: v for i, (rank, v) in enumerate(moiety_d.items())
                  if v['N-description'] != 'Unresolved'}
        new_d = {j: v for j, (rank, v) in enumerate(temp_d.items())}
        new_d[len(new_d)] = unresolved_d
        new_d['Determined db'] = new_d[0]
        self.oad_result_dict[idx][key] = new_d

    def set_re_analysis_result(self, re_analysis_moiety, idx, score_d):
        unsaturated_moiety_num = self.temp_lipid_info['Unsaturated moiety']
        moiety_info = self.temp_lipid_info['Each moiety info']
        if unsaturated_moiety_num == 1:
            #ex) LPC 18:1, PI 18:0_20:4, TG 14:0_16:0_18:1
            key = 'Moiety-1'
        if unsaturated_moiety_num == 2:
            if re_analysis_moiety == 1:
                #ex) PC 18:1_18:2, TG 14:0_16:1_18:1
                key = 'Moiety-1'
            if re_analysis_moiety == 2:
                if len(moiety_info) == 4:
                    key = 'Moiety-2'
                if len(moiety_info) == 6:
                    if moiety_info['db-1'] > 0:
                        key = 'Moiety-2'
                    else:
                        key = 'Moiety-1'
            if re_analysis_moiety == 3:
                key = 'Moiety-2'
        if unsaturated_moiety_num == 3:
            #ex) TG 14:1_16:1_18:1, Cer 18:1;2O/32:1(O-18:2)
            key = f'Moiety-{re_analysis_moiety}'
        moiety = self.oad_result_dict[idx][key]
        last = len(moiety) -1
        self.oad_result_dict[idx][key][last] = score_d[0]
        return key

    def check_ion_lacking_pos(self, re_anal_bool):
        for key, bools in re_anal_bool.items():
            if len(key) == 1:
                n_description = f'n-{key[0]}*'
            for i, (pos, b) in enumerate(zip(key, bools)):
                if i == 0:
                    if b == True:
                        n_description = f'n-{pos}'
                    else:
                        n_description = f'n-{pos}*'
                else:
                    if b == True:
                        n_description += f',{pos}'
                    else:
                        n_description += f',{pos}*'
        return n_description

    #endregion

class BatchSettingWindow(object):
    def __init__(self, master):
        self.master = master
        self.batch_setwin = tk.Toplevel(master)
        self.batch_setwin.attributes('-topmost', False)
        self.center_position(self.batch_setwin, width=690, height=360)
        self.batch_setwin.resizable(width=False, height=False)
        self.batch_setwin.title("Start up a project")
        self.batch_setwin.grab_set()
        self.batch_setwin.iconphoto(False, tk.PhotoImage(file=IconPath))
        self.create_widgets()
        self.batch_setwin.protocol(
            'WM_DELETE_WINDOW', lambda : self.close_window())
        self.files_info = {}

    def create_widgets(self):
        # Project path
        prj_entry_label = ttk.Label(
            self.batch_setwin, text="Project file path :")
        prj_entry_label.place(x=30, y=30)
        self.prj_path = tk.StringVar()
        self.prj_path_entry = ttk.Entry(
            self.batch_setwin, textvariable=self.prj_path, width=68)
        self.prj_path_entry.place(x=150, y=30)
        self.prj_path_btn = ttk.Button(
            self.batch_setwin, text="Browse", command=self.get_filepath)
        self.prj_path_btn.place(x=580, y=28)
        # Alignment path
        alignment_label = ttk.Label(
            self.batch_setwin, text="Alignment data :")
        alignment_label.place(x=30, y=60)
        self.alignment_path = tk.StringVar()
        self.alignment_entry = ttk.Entry(
            self.batch_setwin, textvariable=self.alignment_path, width=68)
        self.alignment_entry.place(x=150, y=60)
        self.alignment_btn = ttk.Button(
            self.batch_setwin, text="Browse", command=self.get_alignment_path)
        self.alignment_btn.place(x=580, y=58)
        # PeakID path
        peakid_label = ttk.Label(
            self.batch_setwin, text="PeakID data :")
        peakid_label.place(x=30, y=90)
        self.peakid_path = tk.StringVar()
        self.peakid_entry = ttk.Entry(
            self.batch_setwin, textvariable=self.peakid_path, width=68)
        self.peakid_entry.place(x=150, y=90)
        self.peakid_btn = ttk.Button(
            self.batch_setwin, text="Browse", command=self.get_peakid_path)
        self.peakid_btn.place(x=580, y=88)
        # PeakList path
        self.wrapper_label = tk.StringVar()
        self.wrapper_label.set("PeakList files")
        self.wrapper_frame = ttk.LabelFrame(self.batch_setwin, 
            text=self.wrapper_label.get())
        self.wrapper_frame.place(x=26, y=120, width=630, height=140)
        self.peaklist_btn = ttk.Button(
            self.batch_setwin, text="Browse", command=self.get_peaklists_path)
        self.peaklist_btn.place(x=580, y=268)
        # Normalized Check
        self.normalized_data = tk.BooleanVar()
        self.normalized_data.set(True)
        self.normalized_check = ttk.Checkbutton(
            self.batch_setwin, text="Check if the alignment data are normalized", 
            variable=self.normalized_data)
        self.normalized_check.place(x=30, y=270)
        # Clear button
        self.clear_btn = ttk.Button(self.batch_setwin, text="Clear", 
            command=self.clear_button)
        self.clear_btn.place(x=489, y=308)
        # Finish button
        self.finish_btn = ttk.Button(self.batch_setwin, text="Start")
        self.finish_btn.place(x=580, y=308) 

    def create_entry_table(self, paths):
        self.inner_canvas = tk.Canvas(self.wrapper_frame, highlightthickness=0)
        label1 = ttk.Label(self.inner_canvas, text='File name')
        label1.grid(row=1, column=1)
        label2 = ttk.Label(self.inner_canvas, text='Path')
        label2.grid(row=1, column=2)
        file_names = [path.rsplit('/', 1)[-1].replace('.txt', '') 
                      for path in paths]
        self.entries = {col: [] for col in ['Files', 'Paths']}
        y_range = len(file_names)*40
        self.inner_canvas.place(x=0, y=0, width=630, height=y_range)
        for i, (name, path) in enumerate(zip(file_names, paths)):
            f_w, p_w = 20, 80
            self.entries['Files'].append(tk.Entry(self.inner_canvas, width=f_w))
            self.entries['Paths'].append(tk.Entry(self.inner_canvas, width=p_w))
            self.entries['Files'][i].grid(row=i+2, column=1)
            self.entries['Paths'][i].grid(row=i+2, column=2)
            # self.inner_canvas.create_window(
            #     x=0, y=i*40, width=f_w, window=self.entries['Files'][i])
            # self.inner_canvas.create_window(
            #     x=20, y=i*40, width=f_w, window=self.entries['Paths'][i])
            self.entries['Files'][i].insert(0, name)
            self.entries['Paths'][i].insert(0, path)
        self.v_scrollbar = ttk.Scrollbar(self.wrapper_frame, orient=tk.VERTICAL,
            command=self.inner_canvas.yview)
        self.v_scrollbar.pack(side=tk.RIGHT, fill='y')
        self.inner_canvas.configure(yscrollcommand=self.v_scrollbar.set)
        # self.inner_canvas.configure(scrollregion=self.inner_canvas.bbox('all'),
        #     yscrollcommand=self.v_scrollbar.set)
        # self.inner_canvas.config(scrollregion=(0, 0, 100, y_range))

    def get_filepath(self):
        old_path = self.prj_path_entry.get()
        new_path = filedialog.askdirectory(title="Select project folder")
        if new_path == "":
            self.prj_path.set(old_path)
        else:
            timestamp = get_timestamp()
            default_path = f'{new_path}/{timestamp}_analysis_table.pkl'
            self.prj_path.set(default_path)

    def get_alignment_path(self):
        old_path = self.alignment_entry.get()
        filetype_list = [("text file", "*.txt")]
        new_path = filedialog.askopenfilename(
            filetypes=filetype_list, title="Select file")
        if new_path == "": self.alignment_path.set(old_path)
        else: self.alignment_path.set(new_path)

    def get_peakid_path(self):
        old_path = self.peakid_entry.get()
        filetype_list = [("text file", "*.txt")]
        new_path = filedialog.askopenfilename(
            filetypes=filetype_list, title="Select file")
        if new_path == "": self.peakid_path.set(old_path)
        else: self.peakid_path.set(new_path)

    def get_peaklists_path(self):
        # old_path = self.peakid_entry.get()
        filetype_list = [("text file", "*.txt")]
        peaklists = filedialog.askopenfilenames(
            filetypes=filetype_list, title="Select files")
        if peaklists:
            self.create_entry_table(peaklists)

    def clear_button(self):
        self.prj_path.set("")
        self.alignment_path.set("")
        self.peakid_path.set("")
        self.inner_canvas.destroy()
        self.v_scrollbar.destroy()
        self.inner_canvas = tk.Canvas(self.wrapper_frame, highlightthickness=0)
        self.inner_canvas.place(relx=0, rely=0, relwidth=1, relheight=1)
        # self.sample_name.set("")

    def center_position(self, win, width, height):
        x = (win.winfo_screenwidth() // 2) - (width // 2)
        y = (win.winfo_screenheight() // 2) - (height // 2)
        win.geometry('{}x{}+{}+{}'.format(width, height, x, y))

    def close_window(self):
        self.batch_setwin.destroy()

class ProgressBar(object):
    def __init__(self, master, title, detail):
        self.prgbar_win = tk.Toplevel(master)
        self.prgbar_win.attributes('-topmost', True)
        self.center_position(self.prgbar_win, width=690, height=110)
        self.prgbar_win.resizable(width=False, height=False)
        self.prgbar_win.title(title)
        self.prgbar_win.grab_set()
        self.prgbar_win.iconphoto(False, tk.PhotoImage(file=IconPath))
        self.prgbar_win.protocol('WM_DELETE_WINDOW', lambda : self.close_window())
        self.create_prgbar_widgets(detail)

    def center_position(self, win, width, height):
        x = (win.winfo_screenwidth() // 2) - (width // 2)
        y = (win.winfo_screenheight() // 2) - (height // 2)
        win.geometry('{}x{}+{}+{}'.format(width, height, x, y))

    def create_prgbar_widgets(self, detail):
        # Progressbar
        self.prgbar = ttk.Progressbar(self.prgbar_win, orient='horizontal', length=630, mode='indeterminate')
        self.prgbar.place(x=30, y=30)
        self.prgbar.configure(value=0, maximum=100)
        self.prgbar_report = tk.StringVar()
        self.prgbar_report.set(detail)
        prgbar_label = ttk.Label(self.prgbar_win, textvariable=self.prgbar_report, font=("", 14))
        prgbar_label.place(x=30, y=55)
        self.prgbar.start(interval=10)

    def close_window(self):
        self.prgbar_win.destroy()

class TwoPathImporter(object):
    def __init__(self, master, title, type1, type2, fmt1, fmt2):
        self.master = master
        self.type1, self.type2 = type1, type2
        self.fmt1, self.fmt2 = fmt1, fmt2
        self.two_path_imp_win = tk.Toplevel(master)
        self.two_path_imp_win.attributes('-topmost', False)
        self.center_position(self.two_path_imp_win, width=690, height=145)
        self.two_path_imp_win.resizable(width=False, height=False)
        self.two_path_imp_win.title(title)
        self.title = title
        self.two_path_imp_win.grab_set()
        self.two_path_imp_win.iconphoto(False, tk.PhotoImage(file=IconPath))
        self.create_widgets()
        self.process_done = False
        self.two_path_imp_win.protocol('WM_DELETE_WINDOW', 
            lambda : self.close_window())

    def create_widgets(self):
        # path1
        label1 = ttk.Label(self.two_path_imp_win, text=self.type1+' :')
        label1.place(x=30, y=20)
        self.path1 = tk.StringVar()
        self.path1_entry = ttk.Entry(self.two_path_imp_win, textvariable=self.path1, width=72)
        self.path1_entry.place(x=120, y=20)
        self.path1_btn = ttk.Button(self.two_path_imp_win, text="Select", command=self.get_path1)
        self.path1_btn.place(x=580, y=18)
        # path2
        label2 = ttk.Label(self.two_path_imp_win, text=self.type2+' :')
        label2.place(x=30, y=60)
        self.path2 = tk.StringVar()
        self.path2_entry = ttk.Entry(self.two_path_imp_win, textvariable=self.path2, width=72)
        self.path2_entry.place(x=120, y=60)
        self.path2_btn = ttk.Button(self.two_path_imp_win, text="Select", command=self.get_path2)
        self.path2_btn.place(x=580, y=58)
        # Data format1
        dataformat_label = ttk.Label(self.two_path_imp_win, text="Data format :")
        dataformat_label.place(x=30, y=103)
        format_label1 = ttk.Label(self.two_path_imp_win, text=self.fmt1)
        format_label1.place(x=120, y=103)
        self.format1 = tk.StringVar()
        self.format1.set("")
        self.data1_rbtn1 = ttk.Radiobutton(self.two_path_imp_win, text="PeakList", value="PeakList", variable=self.format1)
        self.data1_rbtn2 = ttk.Radiobutton(self.two_path_imp_win, text="Alignment", value="Alignment", variable=self.format1)
        self.data1_rbtn1.place(x=145, y=102)
        self.data1_rbtn2.place(x=215, y=102)
        # Data format2
        format_label2 = ttk.Label(self.two_path_imp_win, text=self.fmt2)
        format_label2.place(x=330, y=103)
        self.format2 = tk.StringVar()
        self.format2.set("")
        self.data2_rbtn1 = ttk.Radiobutton(self.two_path_imp_win, text="PeakList", value="PeakList", variable=self.format2)
        self.data2_rbtn2 = ttk.Radiobutton(self.two_path_imp_win, text="Alignment", value="Alignment", variable=self.format2)
        self.data2_rbtn1.place(x=360, y=102)
        self.data2_rbtn2.place(x=430, y=102)
        # Start Button
        self.finish_btn = ttk.Button(self.two_path_imp_win, 
            text="Start", command=self.check_input_data_path)
        self.finish_btn.place(x=580, y=100) 
    
    def check_input_data_path(self):
        path1 = self.path1_entry.get()
        format1 = self.format1.get()
        path2 = self.path2_entry.get()
        format2 = self.format2.get()
        if path1 == "":
            messagebox.showwarning("Empty input data", "Please select {} file".format(self.type1))
        elif format1 == "":
            messagebox.showwarning("Select data format", "Please select {} file format".format(self.type1))
        elif path2 == "":
            messagebox.showwarning("Empty input data", "Please select {} file".format(self.type2))
        elif format2 == "":
            messagebox.showwarning("Select data format", "Please select {} file format".format(self.type2))
        else:
            thread_1 = threading.Thread(target=self.create_preprocess_prgbar_window)
            thread_2 = threading.Thread(target=self.generate_input_data)
            thread_1.start()
            thread_2.start()
            self.process_done = False

    def create_preprocess_prgbar_window(self):
        # logging.debug('create_preprocess_prgbar_window -> Start')
        self.temp_prgbar = ProgressBar(self.master,
            title="Data preprocess in progress", 
            detail="--- Data preprocessing ---")
        while not self.process_done:
            time.sleep(0.1)
        self.temp_prgbar.close_window()
        # logging.debug('create_preprocess_prgbar_window -> End')

    def generate_input_data(self):
        # logging.debug('generate_input_data -> Start')
        path1 = self.path1_entry.get()
        format1 = self.format1.get()
        path2 = self.path2_entry.get()
        format2 = self.format2.get()
        data_preprocessor =  DataPreprocessor(
            path1=path1, fmt1=format1, path2=path2, fmt2=format2)
        # print(f'Title={self.two_path_imp_win.title}')
        self.two_path_imp_win.destroy()
        try:
            if self.title == "Merge Neg&Pos CID-MS/MS data":
                data_preprocessor.merge_bipolarity_cid_data()
            elif self.title == "Merge CID&OAD MS/MS data":
                data_preprocessor.merge_cid_and_oad_data()
        except Exception as e:
            self.temp_prgbar.prgbar_report.set("--- Error occurred ---")
            messagebox.showwarning("Error message", "Error has occurred while the data preprocess, please confirm input data.")
            messagebox.showwarning("Error details", "{}".format(e))
            traceback.print_exc()
        self.process_done = True
        # logging.debug('generate_input_data -> End')

    def get_path1(self):
        old_path = self.path1_entry.get()
        filetype_list = [("text file", "*.txt")]
        new_path = filedialog.askopenfilename(filetypes=filetype_list, 
            title="Select {} file".format(self.type1))
        if new_path == "":
            self.path1.set(old_path)
        else:
            self.path1.set(new_path)

    def get_path2(self):
        old_path = self.path2_entry.get()
        filetype_list = [("text file", "*.txt")]
        new_path = filedialog.askopenfilename(filetypes=filetype_list, 
            title="Select {} file".format(self.type2))
        if new_path == "":
            self.path2.set(old_path)
        else:
            self.path2.set(new_path)

    def center_position(self, win, width, height):
        x = (win.winfo_screenwidth() // 2) - (width // 2)
        y = (win.winfo_screenheight() // 2) - (height // 2)
        win.geometry('{}x{}+{}+{}'.format(width, height, x, y))

    def close_window(self):
        self.two_path_imp_win.destroy()

class PopUpWindow(object):
    def __init__(self, master, title, message):
        self.master = master
        self.popup_win = tk.Toplevel(master)
        self.popup_win.attributes('-topmost', False)
        self.center_position(self.popup_win, width=300, height=80)
        self.popup_win.resizable(width=False, height=False)
        self.popup_win.title(title)
        self.popup_win.iconphoto(False, tk.PhotoImage(file=IconPath))
        label = ttk.Label(self.popup_win, text=message, font=("", 12))
        label.pack(anchor='center', expand=1)

    def center_position(self, win, width, height):
        x = (win.winfo_screenwidth() // 2) - (width // 2)
        y = (win.winfo_screenheight() // 2) - (height // 2)
        win.geometry('{}x{}+{}+{}'.format(width, height, x, y))

    def close_window(self):
        self.popup_win.destroy()

class ReAnalysisWindow(object):
    def __init__(self, master, name, info):
        self.master = master
        self.metabolite_name = name
        self.structure_info = info
        self.window = tk.Toplevel(master)
        self.window.attributes('-topmost', False)
        self.center_position(self.window, width=600, height=200)
        self.window.resizable(width=False, height=False)
        self.window.title("Re-analysis of C=C positions")
        self.window.grab_set()
        self.window.iconphoto(False, tk.PhotoImage(file=IconPath))
        self.create_widgets()

    def create_widgets(self):
        txt = f'Metabolite name : {self.metabolite_name}'
        chain1, chain2, chain3, lock1, lock2, lock3 \
            = self.check_unsaturated_moiety()
        name_label = ttk.Label(self.window, text=txt, font=("", 13))
        name_label.place(relx=0.5, rely=0.07, anchor='center')
        flame1 = tk.Frame(self.window)
        flame1.place(relx=0.01, rely=0.15, relwidth=0.32, relheight=0.6)
        flame2 = tk.Frame(self.window)
        flame2.place(relx=0.34, rely=0.15, relwidth=0.32, relheight=0.6)
        flame3 = tk.Frame(self.window)
        flame3.place(relx=0.67, rely=0.15, relwidth=0.32, relheight=0.6)
        title1 = ttk.Label(flame1, text=f'moiety-1: {chain1}', font=("", 13))
        title1.place(relx=0.5, rely=0.05, anchor='center')
        self.text1 = scrolledtext.ScrolledText(flame1, font=(None, 14))
        self.text1.place(
            relx=0.05, rely=0.2, relwidth=0.9, relheight=0.7)
        title2 = ttk.Label(flame2, text=f'moiety-2: {chain2}', font=("", 13))
        title2.place(relx=0.5, rely=0.05, anchor='center')
        self.text2 = scrolledtext.ScrolledText(flame2, font=(None, 14))
        self.text2.place(
            relx=0.05, rely=0.2, relwidth=0.9, relheight=0.7)
        title3 = ttk.Label(flame3, text=f'moiety-3: {chain3}', font=("", 13))
        title3.place(relx=0.5, rely=0.05, anchor='center')
        self.text3 = scrolledtext.ScrolledText(flame3, font=(None, 14))
        self.text3.place(
            relx=0.05, rely=0.2, relwidth=0.9, relheight=0.7)

        notice_txt = ("* Set C=C position as (6,9,12,15). "
                     +"Multiple candidates must be \n   "
                     +"written in each line.")
        notice = ttk.Label(self.window, text=notice_txt)
        notice.place(relx=0.01, rely=0.8, relwidth=0.65)
        self.start_btn = ttk.Button(self.window, text="Start")
        self.start_btn.place(relx=0.67, rely=0.82)
        self.close_btn = ttk.Button(self.window, text="Close")
        self.close_btn.place(relx=0.83, rely=0.82)
        self.lock_text(lock1, lock2, lock3)

    def check_unsaturated_moiety(self):
        def get_moiety_str(c_num, db):
            b = False if db > 0 else True
            return f'{c_num}:{db}', b
        chain_num = self.structure_info['Valid moiety num']
        unsaturated_moiety = self.structure_info['Unsaturated moiety']
        moiety_d = self.structure_info['Each moiety info']
        ontology = self.structure_info['Ontology']
        chain1, lock1 = 'none', True
        chain2, lock2 = 'none', True
        chain3, lock3 = 'none', True
        if chain_num >= 1:
            chain1, lock1 = get_moiety_str(moiety_d['chain-1'],moiety_d['db-1'])
        if chain_num >= 2:
            chain2, lock2 = get_moiety_str(moiety_d['chain-2'],moiety_d['db-2'])
        if chain_num >= 3:
            chain3, lock3 = get_moiety_str(moiety_d['chain-3'],moiety_d['db-3'])
        is_sph = check_lipid_category('Sphingolipids', ontology)
        if is_sph:
            chain1, chain2, chain3 = self.refine_sph_structure(
                chain1, chain2, chain3)
        return chain1, chain2, chain3, lock1, lock2, lock3

    def refine_sph_structure(self, chain1, chain2, chain3):
        ox_num = re.findall(r';\dO', self.metabolite_name)[0]
        chain1 = f'{chain1}{ox_num}'
        if 'O-' in self.metabolite_name:
            chain3 = re.findall(r'O\-\d+\:\d', self.metabolite_name)[0]
        return chain1, chain2, chain3

    def lock_text(self, lock1, lock2, lock3):
        lock_color = '#dcdcdc'
        if lock1:
            self.text1['state'] = tk.DISABLED
            self.text1['bg'] = lock_color
        if lock2:
            self.text2['state'] = tk.DISABLED
            self.text2['bg'] = lock_color
        if lock3:
            self.text3['state'] = tk.DISABLED
            self.text3['bg'] = lock_color

    def center_position(self, win, width, height):
        x = (win.winfo_screenwidth() // 2) - (width // 2)
        y = (win.winfo_screenheight() // 2) - (height // 2)
        win.geometry('{}x{}+{}+{}'.format(width, height, x, y))

    def close_window(self):
        self.window.destroy()

class VerticalProgressBars(object):
    def __init__(self, master, title, start_text):
        self.vert_prgbars_win = tk.Toplevel(master)
        self.vert_prgbars_win.attributes('-topmost', True)
        self.center_position(self.vert_prgbars_win, width=690, height=155)
        self.vert_prgbars_win.resizable(width=False, height=False)
        self.vert_prgbars_win.title(title)
        self.vert_prgbars_win.grab_set()
        self.vert_prgbars_win.iconphoto(False, tk.PhotoImage(file=IconPath))
        self.vert_prgbars_win.protocol(
            'WM_DELETE_WINDOW', lambda : self.close_window())
        self.start_text = start_text
        self.create_widgets()
        self.running = True
        # self.process_timer()

    def center_position(self, win, width, height):
        x = (win.winfo_screenwidth() // 2) - (width // 2)
        y = (win.winfo_screenheight() // 2) - (height // 2)
        win.geometry('{}x{}+{}+{}'.format(width, height, x, y))

    def create_widgets(self):
        # Section Bar
        self.section_bar = ttk.Progressbar(self.vert_prgbars_win, 
            orient='horizontal', length=534, mode='determinate')
        self.section_bar.place(x=30, y=30)
        self.section_report = tk.StringVar()
        self.section_report.set(self.start_text)
        section_label = ttk.Label(self.vert_prgbars_win, 
            textvariable=self.section_report, font=("", 14))
        section_label.place(x=30, y=55)
        # Section Timer
        self.section_timer = tk.StringVar()
        self.section_timer_label = ttk.Label(
            self.vert_prgbars_win, 
            textvariable=self.section_timer, font=("", 14))
        self.section_timer_label.place(x=580, y=28)
        self.section_timer.set("00:00:00")
        # Each progress
        self.each_prgbar = ttk.Progressbar(self.vert_prgbars_win, 
            orient='horizontal', length=620, mode='determinate')
        self.each_prgbar.place(x=30, y=90)
        self.each_report = tk.StringVar()
        self.each_report.set("")
        each_label = ttk.Label(self.vert_prgbars_win, 
            textvariable=self.each_report, font=("", 14))
        each_label.place(x=30, y=115)
        # Each progress Timer
        self.small_timer = tk.StringVar()

    def process_timer(self):
        hour = 0
        minute = 0
        second = 0
        while self.running:
            second += 1
            time.sleep(1)
            if second == 60:
                minute += 1
                second = 0
                if minute == 60:
                    minute = 0
                    hour += 1
            if minute >= 10 and second >= 10:
                self.section_timer.set(f'0{hour}:{minute}:{second}')
            elif minute < 10 and second >= 10:
                self.section_timer.set(f'0{hour}:0{minute}:{second}')
            elif minute >= 10 and second < 10:
                self.section_timer.set(f'0{hour}:{minute}:0{second}')
            else:
                self.section_timer.set(f'0{hour}:0{minute}:0{second}')

    def close_window(self):
        self.vert_prgbars_win.destroy()

#region Utilities
def math_floor(num, digit):
    floored = math.floor(num*10**digit)/(10**digit)
    return floored

def get_timestamp():
    dt_now = datetime.datetime.now()
    idx = str(dt_now).find('.') +1
    stamp = str(dt_now)[:idx].replace('-', '').replace(' ', '_').replace('.', 's')
    timestamp = stamp.replace(':', 'h', 1).replace(':', 'm', 1)
    return timestamp

def check_lipid_category(category, target):
    lipidclass_dict =  {
        'Fatty acyls': ['FA', 'NAGly', 'NAGlySer', 'NAOrn', 'NAE', 'CAR', 
                        'FAHFA'],
        'Glycerolipids': ['DG', 'EtherDG', 'DGDG', 'EtherDGDG', 'MGDG', 
                          'EtherMGDG', 'SQDG', 'EtherSMGDG', 'MG', 'ADGGA', 
                          'DGCC', 'DGGA', 'DGTS/A', 'LDGCC', 'LDGTS/A', 
                          'EtherTG', 'TG'],
        'Glycerophospholipids': ['LPA', 'PA', 'EtherLPC', 'EtherPC', 'LPC', 
                                 'PC', 'EtherLPE', 'EtherPE', 'EtherPE(P)', 
                                 'PlasmPE', 'LNAPE', 'LPE', 'PE', 'BMP', 
                                 'EtherLPG', 'EtherPG', 'HBMP', 'LPG', 'PG', 
                                 'CL', 'DLCL', 'MLCL', 'Ac2PIM1', 'Ac2PIM2', 
                                 'Ac3PIM2', 'Ac4PIM2', 'EtherPI', 'LPI', 'PI', 
                                 'EtherPS', 'LNAPS', 'LPS', 'PS', 'PEtOH', 
                                 'PMeOH', 'EtherOxPE', 'OxPC', 'OxPE', 'OxPG', 
                                 'OxPI', 'OxPS'],
        'Prenol lipids': ['VAE', 'CoQ', 'Vitamine E'],
        'Saccharolipids': ['LipidA'], 
        'Sphingolipids': ['GM3', 'SHexCer', 'SHexCer+O', 
                          'Cer_ADS', 'Cer_AP', 'Cer_AS', 'Cer_BDS', 'Cer_BS', 
                          'Cer_HDS', 'Cer_HS', 'Cer_EBDS', 'Cer_EODS', 
                          'Cer_EOS', 'Cer_NDS', 'Cer_NP', 'Cer_NS', 'CerP', 
                          'AHexCer', 'HexCer_ADS', 'HexCer_AP', 'HexCer_AS', 
                          'HexCer_BDS', 'HexCer_BS', 'HexCer_HDS', 'HexCer_HS', 
                          'HexCer_EOS', 'HexCer_NDS', 'HexCer_NP', 'HexCer_NS', 
                          'Hex2Cer', 'Hex3Cer', 'ASM', 'PE-Cer', 'PE-Cer+O', 
                          'PI-Cer', 'SM', 'SM+O', 'PhytoSph', 'SL', 'SL+O', 
                          'DHSph', 'Sph'],
        'Sterol lipids': ['CASulfate', 'CA', 'DCAE', 'GDCAE', 'GLCAE', 'TDCAE', 
                          'TLCAE', 'AHexCAS', 'AHexCS', 'AHexSIS', 'AHexBRS', 
                          'AHexSTS', 'Vitamine D', 'SSulfate', 'BRSE', 'CASE', 
                          'CE', 'Cholesterol', 'SHex', 'SISE', 'STSE', 'SPE', 
                          'BAHex', 'BASulfate', 'SPEHex', 'SPGHex', 'BRSLPHex', 
                          'BRSPHex', 'CASLPHex', 'CASPHex', 'SISLPHex', 
                          'SISPHex', 'STSLPHex', 'STSPHex']
                    }
    if category == 'Fatty acyls':
        return target in lipidclass_dict['Fatty acyls']
    elif category == 'Glycerolipids':
        return target in lipidclass_dict['Glycerolipids']
    elif category == 'Glycerophospholipids':
        return target in lipidclass_dict['Glycerophospholipids']
    elif category == 'Prenol lipids':
        return target in lipidclass_dict['Prenol lipids']
    elif category == 'Saccharolipids':
        return target in lipidclass_dict['Saccharolipids']
    elif category == 'Sphingolipids':
        return target in lipidclass_dict['Sphingolipids']
    elif category == 'Sterol lipids':
        return target in lipidclass_dict['Sterol lipids']

#endregion

#region Text format exporter
def text_sheet_exporter(path, target_table, cid, oad, structure_info, 
    normalized, stamp):
    ion_v_col = 'pmol/mg tissue' if normalized else 'Height'
    idxs = target_table['ID'].values
    alignment_ids = [-1 for _ in range(len(idxs))]
    if 'Alignment ID' in target_table.columns:
        alignment_ids = target_table['Alignment ID'].values
    moiety_info = [fill_in_oad_result_on_summary_df(i, oad[idx]) 
                   for i, idx in enumerate(idxs)]
    data_d = {'ID': target_table['ID'].values, 
              'RT': target_table['RT(min)'].values, 
              'MS1 m/z': target_table['Precursor m/z'].values, 
              'Ref m/z': [structure_info[i]['Ref precursor Mz'] 
                          for i in idxs], 
              'Precise m/z': target_table['Precise m/z'].values, 
              'm/z type': target_table['Precise m/z type'].values, 
              'Metabolite name': target_table['Metabolite name'].values, 
              'OAD result name': target_table['OAD result name'].values, 
              'Comment': target_table['User comment'].values, 
              'Data from': target_table['Data from'].values, 
              ion_v_col: target_table[ion_v_col].values, 
              'Manually changed': [v[3] for v in moiety_info], 
              'Moiety1 Result': [v[0] for v in moiety_info], 
              'Moiety2 Result': [v[1] for v in moiety_info], 
              'Moiety3 Result': [v[2] for v in moiety_info], 
              'Moiety1 resolved': [v[4] for v in moiety_info], 
              'Moiety2 resolved': [v[5] for v in moiety_info], 
              'Moiety3 resolved': [v[6] for v in moiety_info],
              'C=C in Moiety1': [v[7] for v in moiety_info], 
              'C=C in Moiety2': [v[8] for v in moiety_info], 
              'C=C in Moiety3': [v[9] for v in moiety_info], 
              'Ontology': target_table['Ontology'].values, 
              'Heads': [cid[i]['Lipid subclass']['Presence'] 
                        if cid[i]['Lipid subclass'] else 0 for i in idxs], 
              'Moiety': [cid[i]['Moiety']['Presence'] 
                         if cid[i]['Moiety'] else 0 for i in idxs], 
              'Alignment ID': alignment_ids}
    summary_df = pd.DataFrame(data_d, index=range(len(target_table)))
    summary_df = summary_df.dropna(how='all')
    excel_path = path + '/' + stamp + '_MSRIDD_Analysis_Result.txt'
    summary_df.to_csv(excel_path, sep='\t', index=False)
#endregion

#region Excel exporter
def find_null_cell(excel_sheet):
    """ Find null cell in excel sheet to add new value. 
        Return null cell position """
    column_count = 1
    while True:
        if excel_sheet.cell(row=8, column=column_count).value == None:
            return column_count
        else:
            column_count += 1

def write_dataframe_into_excel_sheet(sheet, df, msms_df, cid, oad, lipid_info):
    #region lipid_info structure
    # lipid_structural_dict = {'Status': '', 'Adduct': '', 'Precursor Mz': 0, 'MS2 Mz': 0,
    #                          'Precise precursor Mz': '', 'Ref precursor Mz': 0, 
    #                          'RT(min)': 0, 'Reference RT': 0, 'Ontology': '',
    #                          'Brutto': '', 'Valid moiety num': 0, 'Each moiety info': {}, 
    #                          'Unsaturated moiety': 0, 'Unsaturated sphingobase': False, 
    #                          'SMILES': '', 'Formula': '', 'Atom dict': '', 'Oxidized': 0,
    #                          'NL type': []}
    #endregion
    #region CID result dict structure
    # cid_dict = {'Lipid subclass': {'Glycine': [ref_mz, measured_mz, intensity, ratio, ppm],
    #                                'Presence': 50.00},
    #             'Moiety':         {'acyl-2': [ref_mz, measured_mz, ppm],
    #                                'Presence': 50.00}
    #            }
    #endregion
    #region OAD result dict structure
    # result_dict = {'Resolved level': 'All' or 'Partial' or 'None', 
    #                'Validated num': 0 ~ 3,
    #                'Each bools': [True, False, ...], 
    #                'Moiety-1': {0: {'Positions': '',
    #                                 'N-description': '',
    #                                 'Score': float,
    #                                 'Ratio sum': float,
    #                                 'Presence': float,
    #                                 'Notice': '',
    #                                 'Peaks dict': {'n-9/dis@n-8/+O/': [Ref m/z, Ref delta, Measured m/z, Measured ratio, ppm]}
    #                                },
    #                             1: {....},
    #                             'Determined db: {'Positions': '',
    #                                              'N-description': '',
    #                                              'Score': float,
    #                                              'Ratio sum': float,
    #                                              'Presence': float,
    #                                              'Notice': '',
    #                                              'Measured peaks': [[Measured m/z, Measured ratio, ppm], [...]],
    #                                              'Ref peaks': [[OAD type, Ref m/z, Ref NL, Ref ratio], [...]]
    #                                             }
    #                             }
    #               }
    #endregion
    #region Rearrange CID Result
    cid_result_txt = ''
    cid_head, cid_moiety = cid['Lipid subclass'], cid['Moiety']
    for key, li in cid_head.items():
        if key != 'Presence':
            cid_result_txt += key + '= m/z '+str(li[1])+'(ppm='+str(li[4])+')\n'
    for key, li in cid_moiety.items():
        if key != 'Presence':
            cid_result_txt += key + '= m/z '+str(li[1])+'(ppm='+str(li[4])+')\n'
    head_percent = cid_head['Presence'] if cid_head else 0
    moiety_percent = cid_moiety['Presence'] if cid_moiety else 0
    cid_result_txt += 'Head= ' + str(head_percent) + '%\n'
    cid_result_txt += 'Moiety= ' + str(moiety_percent) + '%\n'
    #endregion
    def get_rearragned_txt(moiety):
        txt = ''
        for rank, d in moiety.items():
            if rank != 'Determined db':
                txt += 'Position= ' + d['N-description'] + ', ' + 'Score= ' \
                + str(d['Score']) + ', ' + 'Presence= ' + str(d['Presence']) + ', ' \
                + 'Ratio sum= ' + str(d['Ratio sum']) + '\n'
        return txt
    def get_result_info(moiety):
        no1_pos, no2_pos, no3_pos = 'N/A', 'N/A', 'N/A'
        no1_score, no2_score, no3_score = 'N/A', 'N/A', 'N/A'
        no1_presence, no2_presence, no3_presence = 'N/A', 'N/A', 'N/A'
        no1_ratio, no2_ratio, no3_ratio = 'N/A', 'N/A', 'N/A'
        key1, key2, key3, key4 = 'N-description','Score','Presence','Ratio sum'
        for rank, d in moiety.items():
            if rank == 0:
                no1_pos, no1_score, no1_presence, no1_ratio \
                = d[key1], d[key2], d[key3], d[key4]
            if rank == 1:
                if d[key1] == 'Unresolved':
                    pass
                else:
                    no2_pos, no2_score, no2_presence, no2_ratio \
                    = d[key1], d[key2], d[key3], d[key4]
            if rank == 2:
                if d[key1] == 'Unresolved':
                    pass
                else:
                    no3_pos, no3_score, no3_presence, no3_ratio \
                    = d[key1], d[key2], d[key3], d[key4]
        return no1_pos, no2_pos, no3_pos, no1_score, no2_score, no3_score, \
        no1_presence, no2_presence, no3_presence, no1_ratio, no2_ratio, no3_ratio
    #region fill in basic information
    s_col = find_null_cell(sheet)
    sheet.cell(row=1, column=s_col, value='Metabolite name')
    sheet.cell(row=1, column=s_col+1, value=df['Metabolite name'])
    if all(oad['Each bools']):
        sheet.cell(row=1, column=s_col+1).font = Font(bold=True)
    sheet.cell(row=1, column=s_col+2, value='OAD result name')
    sheet.cell(row=1, column=s_col+3, value=df['OAD result name'])
    if all(oad['Each bools']):
        sheet.cell(row=1, column=s_col+3).font = Font(bold=True, color='FF0000')
    sheet.cell(row=2, column=s_col, value='MS1 m/z')
    sheet.cell(row=2, column=s_col+1, value=df['Precursor m/z'])
    sheet.cell(row=2, column=s_col+2, value=df['Precise m/z'])
    mz_type = '<-' + df['Precise m/z type']
    sheet.cell(row=2, column=s_col+3, value=mz_type)
    sheet.cell(row=3, column=s_col, value='Ref m/z')
    sheet.cell(row=3, column=s_col+1, value=lipid_info['Ref precursor Mz'])
    sheet.cell(row=4, column=s_col, value='RT')
    sheet.cell(row=4, column=s_col+1, value=df['RT(min)'])
    sheet.cell(row=5, column=s_col, value='ID')
    sheet.cell(row=5, column=s_col+1, value=df['ID'])
    sheet.cell(row=6, column=s_col, value='Smiles')
    sheet.cell(row=6, column=s_col+1, value=lipid_info['SMILES'])
    sheet.cell(row=7, column=s_col, value='Subclass info')
    sheet.cell(row=7, column=s_col+1, value=cid_result_txt)
    #endregion
    #region fill in each OAD result
    if len(oad) >= 4:
        all_result_1 = get_rearragned_txt(oad['Moiety-1'])
        no1_pos, no2_pos, no3_pos, no1_score, no2_score, no3_score, no1_presence, \
        no2_presence, no3_presence, no1_ratio, no2_ratio, no3_ratio \
        = get_result_info(oad['Moiety-1'])
        #region OAD1
        sheet.cell(row=1, column=s_col+5, value='No.1')
        sheet.cell(row=1, column=s_col+6, value='No.2')
        sheet.cell(row=1, column=s_col+7, value='No.3')
        sheet.cell(row=2, column=s_col+4, value='Positions1')
        sheet.cell(row=2, column=s_col+5, value=no1_pos)
        sheet.cell(row=2, column=s_col+6, value=no2_pos)
        sheet.cell(row=2, column=s_col+7, value=no3_pos)
        sheet.cell(row=3, column=s_col+4, value='Score1')
        sheet.cell(row=3, column=s_col+5, value=no1_score)
        sheet.cell(row=3, column=s_col+6, value=no2_score)
        sheet.cell(row=3, column=s_col+7, value=no3_score)
        sheet.cell(row=4, column=s_col+4, value='Presence1')
        sheet.cell(row=4, column=s_col+5, value=no1_presence)
        sheet.cell(row=4, column=s_col+6, value=no2_presence)
        sheet.cell(row=4, column=s_col+7, value=no3_presence)
        sheet.cell(row=5, column=s_col+4, value='Ratio sum1')
        sheet.cell(row=5, column=s_col+5, value=no1_ratio)
        sheet.cell(row=5, column=s_col+6, value=no2_ratio)
        sheet.cell(row=5, column=s_col+7, value=no3_ratio)
        sheet.cell(row=6, column=s_col+4, value='All results1')
        sheet.cell(row=6, column=s_col+5, value=all_result_1)
        #endregion
    if len(oad) >= 5:
        all_result_2 = get_rearragned_txt(oad['Moiety-2'])
        no1_pos, no2_pos, no3_pos, no1_score, no2_score, no3_score, no1_presence, \
        no2_presence, no3_presence, no1_ratio, no2_ratio, no3_ratio \
        = get_result_info(oad['Moiety-2'])
        #region OAD2
        sheet.cell(row=1, column=s_col+9, value='No.1')
        sheet.cell(row=1, column=s_col+10, value='No.2')
        sheet.cell(row=1, column=s_col+11, value='No.3')
        sheet.cell(row=2, column=s_col+8, value='Positions2')
        sheet.cell(row=2, column=s_col+9, value=no1_pos)
        sheet.cell(row=2, column=s_col+10, value=no2_pos)
        sheet.cell(row=2, column=s_col+11, value=no3_pos)
        sheet.cell(row=3, column=s_col+8, value='Score2')
        sheet.cell(row=3, column=s_col+9, value=no1_score)
        sheet.cell(row=3, column=s_col+10, value=no2_score)
        sheet.cell(row=3, column=s_col+11, value=no3_score)
        sheet.cell(row=4, column=s_col+8, value='Presence2')
        sheet.cell(row=4, column=s_col+9, value=no1_presence)
        sheet.cell(row=4, column=s_col+10, value=no2_presence)
        sheet.cell(row=4, column=s_col+11, value=no3_presence)
        sheet.cell(row=5, column=s_col+8, value='Ratio sum2')
        sheet.cell(row=5, column=s_col+9, value=no1_ratio)
        sheet.cell(row=5, column=s_col+10, value=no2_ratio)
        sheet.cell(row=5, column=s_col+11, value=no3_ratio)
        sheet.cell(row=6, column=s_col+8, value='All results2')
        sheet.cell(row=6, column=s_col+9, value=all_result_2)
        #endregion
    if len(oad) >= 6:
        all_result_3 = get_rearragned_txt(oad['Moiety-3'])
        no1_pos, no2_pos, no3_pos, no1_score, no2_score, no3_score, no1_presence, \
        no2_presence, no3_presence, no1_ratio, no2_ratio, no3_ratio \
        = get_result_info(oad['Moiety-3'])
        #region OAD3
        sheet.cell(row=1, column=s_col+13, value='No.1')
        sheet.cell(row=1, column=s_col+14, value='No.2')
        sheet.cell(row=1, column=s_col+15, value='No.3')
        sheet.cell(row=2, column=s_col+12, value='Positions3')
        sheet.cell(row=2, column=s_col+13, value=no1_pos)
        sheet.cell(row=2, column=s_col+14, value=no2_pos)
        sheet.cell(row=2, column=s_col+15, value=no3_pos)
        sheet.cell(row=3, column=s_col+12, value='Score3')
        sheet.cell(row=3, column=s_col+13, value=no1_score)
        sheet.cell(row=3, column=s_col+14, value=no2_score)
        sheet.cell(row=3, column=s_col+15, value=no3_score)
        sheet.cell(row=4, column=s_col+12, value='Presence3')
        sheet.cell(row=4, column=s_col+13, value=no1_presence)
        sheet.cell(row=4, column=s_col+14, value=no2_presence)
        sheet.cell(row=4, column=s_col+15, value=no3_presence)
        sheet.cell(row=5, column=s_col+12, value='Ratio sum3')
        sheet.cell(row=5, column=s_col+13, value=no1_ratio)
        sheet.cell(row=5, column=s_col+14, value=no2_ratio)
        sheet.cell(row=5, column=s_col+15, value=no3_ratio)
        sheet.cell(row=6, column=s_col+12, value='All results3')
        sheet.cell(row=6, column=s_col+13, value=all_result_3)
        #endregion
    #endregion
    #region fill in Mass Table columns
    sheet.cell(row=8, column=s_col, value='intensity')
    sheet.cell(row=8, column=s_col+1, value='frag m/z')
    sheet.cell(row=8, column=s_col+2, value='Ratio(%)')
    sheet.cell(row=8, column=s_col+3, value='Delta')
    if len(oad) == 4:
        sheet.cell(row=8, column=s_col+4, value='Ref NL1')
        sheet.cell(row=8, column=s_col+5, value='Ref Mz1')
        sheet.cell(row=8, column=s_col+6, value='ppm1')
        sheet.cell(row=8, column=s_col+7, value='Type1')
    elif len(oad) == 5:
        sheet.cell(row=8, column=s_col+4, value='Ref NL1')
        sheet.cell(row=8, column=s_col+5, value='Ref Mz1')
        sheet.cell(row=8, column=s_col+6, value='ppm1')
        sheet.cell(row=8, column=s_col+7, value='Type1')
        sheet.cell(row=8, column=s_col+8, value='Ref NL2')
        sheet.cell(row=8, column=s_col+9, value='Ref Mz2')
        sheet.cell(row=8, column=s_col+10, value='ppm2')
        sheet.cell(row=8, column=s_col+11, value='Type2')
    elif len(oad) == 6:
        sheet.cell(row=8, column=s_col+4, value='Ref NL1')
        sheet.cell(row=8, column=s_col+5, value='Ref Mz1')
        sheet.cell(row=8, column=s_col+6, value='ppm1')
        sheet.cell(row=8, column=s_col+7, value='Type1')
        sheet.cell(row=8, column=s_col+8, value='Ref NL2')
        sheet.cell(row=8, column=s_col+9, value='Ref Mz2')
        sheet.cell(row=8, column=s_col+10, value='ppm2')
        sheet.cell(row=8, column=s_col+11, value='Type2')
        sheet.cell(row=8, column=s_col+12, value='Ref NL3')
        sheet.cell(row=8, column=s_col+13, value='Ref Mz3')
        sheet.cell(row=8, column=s_col+14, value='ppm3')
        sheet.cell(row=8, column=s_col+15, value='Type3')
    #endregion
    #region fill in Mass Table
    mz_list = msms_df['frag m/z'].values.tolist()
    intensity_list = msms_df['intensity'].values.tolist()
    ratio_list = msms_df['Ratio(%)'].values.tolist()
    delta_list = msms_df['Delta'].values.tolist()
    cid_ions_dict = {}
    for d in cid.values():
        for key, li in d.items():
            if key != 'Presence':
                cid_ions_dict[key] = li
    m1, m2, m3, peaks = 'Moiety-1', 'Moiety-2', 'Moiety-3', 'Peaks dict'
    moiety_1 = {'None': [0, 0, 0]} 
    moiety_2 = {'None': [0, 0, 0]}
    moiety_3 = {'None': [0, 0, 0]}
    if (len(oad) >= 4 and oad[m1]): moiety_1 = oad[m1][0][peaks] 
    if (len(oad) >= 5 and oad[m2]): moiety_2 = oad[m2][0][peaks] 
    if (len(oad) >= 6 and oad[m3]): moiety_3 = oad[m3][0][peaks]
    for i, (mz, intensity, ratio, delta) in enumerate(zip(mz_list, 
        intensity_list, ratio_list, delta_list), start=9):
        sheet.cell(row=i, column=s_col, value=intensity)
        sheet.cell(row=i, column=s_col+1, value=mz)
        sheet.cell(row=i, column=s_col+2, value=ratio)
        sheet.cell(row=i, column=s_col+3, value=delta)
        for key, li in cid_ions_dict.items():
            #{'Glycine': [ref_mz, measured_mz, intensity, ratio, ppm], ...}
            if mz == li[1]:
                sheet.cell(row=i, column=s_col+4, value=li[1])
                sheet.cell(row=i, column=s_col+5, value=li[0])
                sheet.cell(row=i, column=s_col+6, value=li[4])
                sheet.cell(row=i, column=s_col+7, value=key)
        for key, li in moiety_1.items():
            #{'n-9/dis@n-8/+O/': 
            # [Ref m/z, Ref delta, Measured m/z, Measured ratio, ppm]}
            if mz == li[2]:
                sheet.cell(row=i, column=s_col+4, value=li[1])
                sheet.cell(row=i, column=s_col+5, value=li[0])
                sheet.cell(row=i, column=s_col+6, value=li[4])
                sheet.cell(row=i, column=s_col+7, value=key)
        for key, li in moiety_2.items():
            if mz == li[2]:
                sheet.cell(row=i, column=s_col+8, value=li[1])
                sheet.cell(row=i, column=s_col+9, value=li[0])
                sheet.cell(row=i, column=s_col+10, value=li[4])
                sheet.cell(row=i, column=s_col+11, value=key)
        for key, li in moiety_3.items():
            if mz == li[2]:
                sheet.cell(row=i, column=s_col+12, value=li[1])
                sheet.cell(row=i, column=s_col+13, value=li[0])
                sheet.cell(row=i, column=s_col+14, value=li[4])
                sheet.cell(row=i, column=s_col+15, value=key)
        #region Coloring
        if ratio < 0.005:
            # background_color = lightcyan
            sheet.cell(row=i, column=s_col+2).fill = PatternFill(patternType='solid', fgColor="E0FFFF")
        elif (0.005 <= ratio < 0.1):
            # background_color = mistyrose
            sheet.cell(row=i, column=s_col+2).fill = PatternFill(patternType='solid', fgColor="FBEFEF")
        elif  (0.1 <= ratio < 1.0):
            # background_color = generated by color picker
            sheet.cell(row=i, column=s_col+2).fill = PatternFill(patternType='solid', fgColor="F6CECE")
            # sheet.cell(row=i, column=s_col+3).font = Font(b=True)
        elif (1.0 <= ratio < 10.0):
            sheet.cell(row=i, column=s_col+2).fill = PatternFill(patternType='solid', fgColor="F78181")
            sheet.cell(row=i, column=s_col+3).font = Font(b=True)
        elif (10.0 <= ratio < 50.0):
            # background_color = red
            sheet.cell(row=i, column=s_col+2).fill = PatternFill(patternType='solid', fgColor="FF0000")
            sheet.cell(row=i, column=s_col+3).font = Font(b=True)
        elif (50.0 <= ratio):
            # background_color = red
            sheet.cell(row=i, column=s_col+2).fill = PatternFill(patternType='solid', fgColor="000000")
            sheet.cell(row=i, column=s_col+2).font = Font(color='FFFFFF')
            sheet.cell(row=i, column=s_col+3).font = Font(b=True)
        #endregion
    #endregion
    return sheet

def excel_sheet_exporter2(path, target_table, msms, cid, oad, structure_info, 
    normalized, stamp, each_bar, each_rep):
    #region CID result dict structure
    # cid_dict = {'Lipid subclass': {'Glycine': [ref_mz, measured_mz, ppm],
    #                                'Presence': 50.00},
    #             'Moiety':         {'acyl-2': [ref_mz, measured_mz, ppm],
    #                                'Presence': 50.00}
    #            }
    #endregion
    #region OAD result dict structure
    # result_dict = {'Resolved level': 'All' or 'Partial' or 'None', 
    #                'Validated num': 0 ~ 3,
    #                'Each bools': [True, False, ...], 
    #                'Moiety-1': {0: {'Positions': '',
    #                                 'N-description': '',
    #                                 'Score': float,
    #                                 'Ratio sum': float,
    #                                 'Presence': float,
    #                                 'Notice': '',
    #                                 'Peaks dict': {'n-9/dis@n-8/+O/': [Ref m/z, Ref delta, Measured m/z, Measured ratio, ppm]}
    #                                },
    #                             1: {....},
    #                             'Determined db: {'Positions': '',
    #                                              'N-description': '',
    #                                              'Score': float,
    #                                              'Ratio sum': float,
    #                                              'Presence': float,
    #                                              'Notice': '',
    #                                              'Measured peaks': [[Measured m/z, Measured ratio, ppm], [...]],
    #                                              'Ref peaks': [[OAD type, Ref m/z, Ref NL, Ref ratio], [...]]
    #                                             }
    #                             }
    #               }
    #endregion
    wb = openpyxl.Workbook()
    sheet_names = []
    ion_v_col = 'pmol/mg tissue' if normalized else 'Height'
    idxs = target_table['ID'].values
    alignment_ids = [-1 for _ in range(len(idxs))]
    if 'Alignment ID' in target_table.columns:
        alignment_ids = target_table['Alignment ID'].values
    moiety_info = [fill_in_oad_result_on_summary_df(i, oad[idx]) 
                   for i, idx in enumerate(idxs)]
    data_d = {'ID': target_table['ID'].values, 
              'RT': target_table['RT(min)'].values, 
              'MS1 m/z': target_table['Precursor m/z'].values, 
              'Ref m/z': [structure_info[i]['Ref precursor Mz'] 
                          for i in idxs], 
              'Precise m/z': target_table['Precise m/z'].values, 
              'm/z type': target_table['Precise m/z type'].values, 
              'Metabolite name': target_table['Metabolite name'].values, 
              'OAD result name': target_table['OAD result name'].values, 
              'Comment': target_table['User comment'].values, 
              'Data from': target_table['Data from'].values, 
              ion_v_col: target_table[ion_v_col].values, 
              'Manually changed': [v[3] for v in moiety_info], 
              'Moiety1 Result': [v[0] for v in moiety_info], 
              'Moiety2 Result': [v[1] for v in moiety_info], 
              'Moiety3 Result': [v[2] for v in moiety_info], 
              'Moiety1 resolved': [v[4] for v in moiety_info], 
              'Moiety2 resolved': [v[5] for v in moiety_info], 
              'Moiety3 resolved': [v[6] for v in moiety_info],
              'C=C in Moiety1': [v[7] for v in moiety_info], 
              'C=C in Moiety2': [v[8] for v in moiety_info], 
              'C=C in Moiety3': [v[9] for v in moiety_info], 
              'Ontology': target_table['Ontology'].values, 
              'Heads': [cid[i]['Lipid subclass']['Presence'] 
                        if cid[i]['Lipid subclass'] else 0 for i in idxs], 
              'Moiety': [cid[i]['Moiety']['Presence'] 
                         if cid[i]['Moiety'] else 0 for i in idxs], 
              'Alignment ID': alignment_ids}
    # table_range = range(len(target_table))
    summary_df = pd.DataFrame(data_d, index=range(len(target_table)))
    total = len(target_table)
    for i, (table_id, row) in enumerate(zip(idxs, target_table.index)):
        start = time.time()
        df = target_table.loc[row]
        msms_df = msms[table_id]
        cid_result, oad_result = cid[table_id], oad[table_id]
        lipid_info = structure_info[table_id]
        new_ontology_candidate = lipid_info['Ontology']
        sheet_names = wb.sheetnames
        if new_ontology_candidate in sheet_names:
            sheetname = check_max_col_in_sheet(wb, new_ontology_candidate, i)
            sheet = wb[sheetname]
        else:
            sheet = wb.create_sheet(title=new_ontology_candidate)
        sheet.freeze_panes = 'A9'
        sheet = write_dataframe_into_excel_sheet(sheet=sheet, 
            df=df, msms_df=msms_df, cid=cid_result, oad=oad_result, 
            lipid_info=lipid_info)
        each_rep.set(
            f'Exporting {i+1}/{total} >>> {time.time()-start:.3f} [sec]')
        each_bar.step(1)
    each_bar.step(0.99)
    each_rep.set(f'Saving files...')
    summary_df = summary_df.dropna(how='all')
    wb._sheets.sort(key=lambda ws: ws.title)
    sheet = wb.create_sheet(index=0, title='Summary')
    sheet.freeze_panes = 'A2'
    sheet.column_dimensions['G'].width = 30
    sheet.column_dimensions['H'].width = 50
    for row in dataframe_to_rows(summary_df, index=False, header=True):
        sheet.append(row)
    del wb['Sheet']
    excel_path = path + '/' + stamp + '_MSRIDD_Analysis_Report.xlsx'
    wb.save(filename=excel_path)

def excel_sheet_exporter(path, target_table, msms, cid, oad, structure_info, 
    normalized, stamp):
    #region CID result dict structure
    # cid_dict = {'Lipid subclass': {'Glycine': [ref_mz, measured_mz, ppm],
    #                                'Presence': 50.00},
    #             'Moiety':         {'acyl-2': [ref_mz, measured_mz, ppm],
    #                                'Presence': 50.00}
    #            }
    #endregion
    #region OAD result dict structure
    # result_dict = {'Resolved level': 'All' or 'Partial' or 'None', 
    #                'Validated num': 0 ~ 3,
    #                'Each bools': [True, False, ...], 
    #                'Moiety-1': {0: {'Positions': '',
    #                                 'N-description': '',
    #                                 'Score': float,
    #                                 'Ratio sum': float,
    #                                 'Presence': float,
    #                                 'Notice': '',
    #                                 'Peaks dict': {'n-9/dis@n-8/+O/': [Ref m/z, Ref delta, Measured m/z, Measured ratio, ppm]}
    #                                },
    #                             1: {....},
    #                             'Determined db: {'Positions': '',
    #                                              'N-description': '',
    #                                              'Score': float,
    #                                              'Ratio sum': float,
    #                                              'Presence': float,
    #                                              'Notice': '',
    #                                              'Measured peaks': [[Measured m/z, Measured ratio, ppm], [...]],
    #                                              'Ref peaks': [[OAD type, Ref m/z, Ref NL, Ref ratio], [...]]
    #                                             }
    #                             }
    #               }
    #endregion
    wb = openpyxl.Workbook()
    sheet_names = []
    ion_v_col = 'pmol/mg tissue' if normalized else 'Height'
    columns_list = ['ID', 'RT', 'MS1 m/z', 'Ref m/z', 'Precise m/z', 'm/z type', 
                    'Metabolite name', 'OAD result name', 'Comment', 
                    'Data from', ion_v_col, 'Manually changed', 
                    'Moiety1 Result', 'Moiety2 Result', 'Moiety3 Result', 
                    'Moiety1 resolved', 'Moiety2 resolved', 'Moiety3 resolved',
                    'C=C in Moiety1', 'C=C in Moiety2', 'C=C in Moiety3', 
                    'Ontology', 'Heads', 'Moiety', 'Alignment ID']
    table_range = range(len(target_table))
    summary_df = pd.DataFrame(index=table_range, columns=columns_list)
    is_batch = 'Alignment ID' in target_table.columns
    total = len(target_table)
    for i, (row, df) in enumerate(target_table.iterrows()):
        start = time.time()
        table_id = df['ID']
        alignment_id = df['Alignment ID'] if is_batch else 0
        msms_df = msms[table_id]
        cid_result = cid[table_id]
        oad_result = oad[table_id]
        lipid_info = structure_info[table_id]
        new_ontology_candidate = lipid_info['Ontology']
        summary_df.loc[i:i, ['ID', 'RT', 'MS1 m/z', 'Ref m/z', 'Precise m/z', 
            'm/z type', 'Metabolite name', 'OAD result name', 'Comment', 
            'Data from', ion_v_col, 'Ontology', 'Alignment ID']] \
        = df['ID'], df['RT(min)'], df['Precursor m/z'], \
        lipid_info['Ref precursor Mz'], df['Precise m/z'], \
        df['Precise m/z type'], df['Metabolite name'], df['OAD result name'], \
        df['User comment'], df['Data from'], df[ion_v_col], df['Ontology'], alignment_id
        summary_df.loc[i:i, ['Heads']] = cid_result['Lipid subclass']['Presence'] if cid_result['Lipid subclass'] else 0
        summary_df.loc[i:i, ['Moiety']] = cid_result['Moiety']['Presence'] if cid_result['Moiety'] else 0
        summary_df.loc[i:i, ['Moiety1 Result', 'Moiety2 Result', 'Moiety3 Result', 
        'Manually changed', 'Moiety1 resolved', 'Moiety2 resolved', 'Moiety3 resolved', 
        'C=C in Moiety1', 'C=C in Moiety2', 'C=C in Moiety3']] \
        = fill_in_oad_result_on_summary_df(row=i, oad_result=oad_result)
        sheet_names = wb.sheetnames
        if new_ontology_candidate in sheet_names:
            sheetname = check_max_col_in_sheet(wb, new_ontology_candidate, i)
            sheet = wb.get_sheet_by_name(sheetname)
        else:
            sheet = wb.create_sheet(title=new_ontology_candidate)
        sheet.freeze_panes = 'A9'
        sheet = write_dataframe_into_excel_sheet(sheet=sheet, 
            df=df, msms_df=msms_df, cid=cid_result, oad=oad_result, 
            lipid_info=lipid_info)
        print(f'Required time for Export: {time.time() - start:.4f} [sec]\t{i+1}/{total}')
    summary_df = summary_df.dropna(how='all')
    wb._sheets.sort(key=lambda ws: ws.title)
    sheet = wb.create_sheet(index=0, title='Summary')
    sheet.freeze_panes = 'A2'
    sheet.column_dimensions['G'].width = 30
    sheet.column_dimensions['H'].width = 50
    for row in dataframe_to_rows(summary_df, index=False, header=True):
        sheet.append(row)
    del wb['Sheet']
    excel_path = path + '/' + stamp + '_MSRIDD_Analysis_Report.xlsx'
    wb.save(filename=excel_path)

def check_max_col_in_sheet(wb, sheetname, num, counter=1):
    sheet = wb.get_sheet_by_name(sheetname)
    max_col = sheet.max_column
    if max_col > 16300:
        counter += 1
        if re.search(r'\d+', sheetname) is None:
            candidate = sheetname + str(counter)
        else:
            candidate = re.sub(r'\d+', str(counter), sheetname)
        if candidate not in wb.sheetnames:
            sheet = wb.create_sheet(title=candidate)
        return check_max_col_in_sheet(wb, candidate, num, counter)
    return sheetname

def fill_in_oad_result_on_summary_df(row, oad_result):
    #region Return values
    # ['Moiety1 Result', 'Moiety2 Result', 'Moiety3 Result', 'Manually changed', 
    #  'Moiety1 resolved', 'Moiety2 resolved', 'Moiety3 resolved', 
    #  'C=C in Moiety1', 'C=C in Moiety2', 'C=C in Moiety3']
    #endregion
    moiety_1_result, moiety_2_result, moiety_3_result = '', '', ''
    manually_changed = 'False'
    resolved_1, resolved_2, resolved_3 = '', '', ''
    db_1, db_2, db_3 = '', '', ''
    def get_rearragned_txt(moiety):
        txt = ''
        for rank, d in moiety.items():
            if rank == 'Determined db' or d['N-description'] == 'Unresolved':
                continue
            txt += (f"Position= {d['N-description']}, Score= {d['Score']},"
                    +f"Presence= {d['Presence']}, Ratio sum= {d['Ratio sum']}\n")
        return txt
    if len(oad_result) == 4:
        if any(oad_result['Each bools']):
            moiety_1_result = get_rearragned_txt(oad_result['Moiety-1'])
            db_1 = oad_result['Moiety-1']['Determined db']['N-description']
            resolved_1 = 'True'
            if 'Modified' in oad_result['Moiety-1']['Determined db']['Notice']:
                manually_changed = 'True'
        else:
            db_1 = 'n-'
            resolved_1 = 'False'
    if len(oad_result) == 5:
        if all(oad_result['Each bools']):
            moiety_1_result = get_rearragned_txt(oad_result['Moiety-1'])
            moiety_2_result = get_rearragned_txt(oad_result['Moiety-2'])
            db_1 = oad_result['Moiety-1']['Determined db']['N-description']
            db_2 = oad_result['Moiety-2']['Determined db']['N-description']
            resolved_1, resolved_2 = 'True', 'True'
            if ('Modified' in oad_result['Moiety-1']['Determined db']['Notice']
                or 'Modified' in oad_result['Moiety-2']['Determined db']['Notice']):
                manually_changed = 'True'
        elif oad_result['Each bools'][0]:
            moiety_1_result = get_rearragned_txt(oad_result['Moiety-1'])
            db_1 = oad_result['Moiety-1']['Determined db']['N-description']
            db_2 = 'n-'
            resolved_1, resolved_2 = 'True', 'False'
            if 'Modified' in oad_result['Moiety-1']['Determined db']['Notice']:
                manually_changed = 'True'
        elif oad_result['Each bools'][1]:
            moiety_2_result = get_rearragned_txt(oad_result['Moiety-2'])
            db_1 = 'n-'
            db_2 = oad_result['Moiety-2']['Determined db']['N-description']
            resolved_1, resolved_2 = 'False', 'True'
            if 'Modified' in oad_result['Moiety-2']['Determined db']['Notice']:
                manually_changed = 'True'
        else:
            db_1, db_2 = 'n-', 'n-'
            resolved_1, resolved_2 = 'False', 'False'
    if len(oad_result) == 6:
        if all(oad_result['Each bools']):
            moiety_1_result = get_rearragned_txt(oad_result['Moiety-1'])
            moiety_2_result = get_rearragned_txt(oad_result['Moiety-2'])
            moiety_3_result = get_rearragned_txt(oad_result['Moiety-3'])
            db_1 = oad_result['Moiety-1']['Determined db']['N-description']
            db_2 = oad_result['Moiety-2']['Determined db']['N-description']
            db_3 = oad_result['Moiety-3']['Determined db']['N-description']
            resolved_1, resolved_2, resolved_3 = 'True', 'True', 'True'
            if ('Modified' in oad_result['Moiety-1']['Determined db']['Notice']
                or 'Modified' in oad_result['Moiety-2']['Determined db']['Notice']
                or 'Modified' in oad_result['Moiety-3']['Determined db']['Notice']):
                manually_changed = 'True'
        elif oad_result['Each bools'][2] and oad_result['Each bools'][1]:
            moiety_2_result = get_rearragned_txt(oad_result['Moiety-2'])
            moiety_3_result = get_rearragned_txt(oad_result['Moiety-3'])
            db_1 = 'n-'
            db_2 = oad_result['Moiety-2']['Determined db']['N-description']
            db_3 = oad_result['Moiety-3']['Determined db']['N-description']
            resolved_1, resolved_2, resolved_3 = 'False', 'True', 'True'
            if ('Modified' in oad_result['Moiety-2']['Determined db']['Notice']
                or 'Modified' in oad_result['Moiety-3']['Determined db']['Notice']):
                manually_changed = 'True'
        elif oad_result['Each bools'][2] and oad_result['Each bools'][0]:
            moiety_1_result = get_rearragned_txt(oad_result['Moiety-1'])
            moiety_3_result = get_rearragned_txt(oad_result['Moiety-3'])
            db_1 = oad_result['Moiety-1']['Determined db']['N-description']
            db_2 = 'n-'
            db_3 = oad_result['Moiety-3']['Determined db']['N-description']
            resolved_1, resolved_2, resolved_3 = 'True', 'False', 'True'
            if ('Modified' in oad_result['Moiety-1']['Determined db']['Notice']
                or 'Modified' in oad_result['Moiety-3']['Determined db']['Notice']):
                manually_changed = 'True'
        elif oad_result['Each bools'][1] and oad_result['Each bools'][0]:
            moiety_1_result = get_rearragned_txt(oad_result['Moiety-1'])
            moiety_2_result = get_rearragned_txt(oad_result['Moiety-2'])
            db_1 = oad_result['Moiety-1']['Determined db']['N-description']
            db_2 = oad_result['Moiety-2']['Determined db']['N-description']
            db_3 = 'n-'
            resolved_1, resolved_2, resolved_3 = 'True', 'True', 'False'
            if ('Modified' in oad_result['Moiety-1']['Determined db']['Notice']
                or 'Modified' in oad_result['Moiety-2']['Determined db']['Notice']):
                manually_changed = 'True'
        elif oad_result['Each bools'][2]:
            moiety_3_result = get_rearragned_txt(oad_result['Moiety-3'])
            db_1, db_2 = 'n-', 'n-'
            db_3 = oad_result['Moiety-3']['Determined db']['N-description']
            resolved_1, resolved_2, resolved_3 = 'False', 'False', 'True'
            if 'Modified' in oad_result['Moiety-3']['Determined db']['Notice']:
                manually_changed = 'True'
        elif oad_result['Each bools'][1]:
            moiety_2_result = get_rearragned_txt(oad_result['Moiety-2'])
            db_1, db_3 = 'n-', 'n-'
            db_2 = oad_result['Moiety-2']['Determined db']['N-description']
            resolved_1, resolved_2, resolved_3 = 'False', 'True', 'False'
            if 'Modified' in oad_result['Moiety-2']['Determined db']['Notice']:
                manually_changed = 'True'
        elif oad_result['Each bools'][0]:
            moiety_1_result = get_rearragned_txt(oad_result['Moiety-1'])
            db_2, db_3 = 'n-', 'n-'
            db_1 = oad_result['Moiety-1']['Determined db']['N-description']
            resolved_1, resolved_2, resolved_3 = 'True', 'False', 'False'
            if 'Modified' in oad_result['Moiety-1']['Determined db']['Notice']:
                manually_changed = 'True'
        else:
            db_1, db_2, db_3 = 'n-', 'n-', 'n-'
            resolved_1, resolved_2, resolved_3 = 'False', 'False', 'False'
    return moiety_1_result, moiety_2_result, moiety_3_result, \
    manually_changed, resolved_1, resolved_2, resolved_3, db_1, db_2, db_3
#endregion

#region PowerPoint exporter
def generate_msms_figures(path, msms_dict, cid, oad, structure_info, graph, 
    stamp, each_bar, each_rep):
    #region base settings
    dpi_setting = 300
    fig_x, fig_y = 4.8, 2.4
    precursor_mz_font_size = 10
    frag_mz_font_size = 6
    x_label_font_size, y_label_font_size = 10, 10
    small_font_size = 8
    middle_font_size = 10
    line_color = 'black'
    moiety_1_color, moiety_2_color, moiety_3_color = 'red', 'blue', 'green'
    xtick_num = 10
    axes = ['top', 'bottom', 'left', 'right']
    axes_width = [0, 1, 1, 0]
    tick_width = 1
    tick_pad = 1
    #endregion
    #region Data structure
    # oad_dict = {0:              {'Positions': '', 'N-description': '', 'Score': float, 
    #                              'Ratio sum': float, 'Presence': float, 'Notice': '',
    #                              'Peaks dict': {'n-9/dis@n-8/+O/': [Ref m/z, Ref delta, Measured m/z, Measured ratio, ppm]}
    #                             },
    #             1:              {....},
    #             Determined db:  {'Positions': '', 'N-description': '', 'Score': float,
    #                              'Ratio sum': float, 'Presence': float, 'Notice': '',
    #                              'Measured peaks': [[Measured m/z, Measured ratio, ppm], [...]],
    #                              'Ref peaks': [[OAD type, Ref m/z, Ref NL, Ref ratio], [...]]
    #                             }
    #             }
    # 'Peaks dict': {'n-9/dis@n-8/+O/': [Ref m/z, Ref delta, Measured m/z, Measured ratio, ppm]}
    #endregion
    total = len(msms_dict)
    for i, (idx, msms_df) in enumerate(msms_dict.items()):
        start = time.time()
        #region Data Pre-processing
        lipid_info = structure_info[idx]
        cid_dict = cid[idx]
        oad_dict = oad[idx]
        graph_dict = graph[idx]['OAD']
        len_oad_dict = len(oad_dict)
        dict_1, dict_2, dict_3 = {}, {}, {}
        if len_oad_dict == 4:
            if all(oad_dict['Each bools']):
                dict_1 = oad_dict['Moiety-1']
        elif len_oad_dict == 5:
            if all(oad_dict['Each bools']):
                dict_1 = oad_dict['Moiety-1']
                dict_2 = oad_dict['Moiety-2']
            elif oad_dict['Each bools'][1]:
                dict_2 = oad_dict['Moiety-2']
            elif oad_dict['Each bools'][0]:
                dict_1 = oad_dict['Moiety-1']
        elif len_oad_dict == 6:
            if all(oad_dict['Each bools']):
                dict_1 = oad_dict['Moiety-1']
                dict_2 = oad_dict['Moiety-2']
                dict_3 = oad_dict['Moiety-3']
            elif oad_dict['Each bools'][2] and oad_dict['Each bools'][1]:
                dict_2 = oad_dict['Moiety-2']
                dict_3 = oad_dict['Moiety-3']
            elif oad_dict['Each bools'][2] and oad_dict['Each bools'][0]:
                dict_1 = oad_dict['Moiety-1']
                dict_3 = oad_dict['Moiety-3']
            elif oad_dict['Each bools'][1] and oad_dict['Each bools'][0]:
                dict_1 = oad_dict['Moiety-1']
                dict_2 = oad_dict['Moiety-2']
            elif oad_dict['Each bools'][2]:
                dict_3 = oad_dict['Moiety-3']
            elif oad_dict['Each bools'][1]:
                dict_2 = oad_dict['Moiety-2']
            elif oad_dict['Each bools'][0]:
                dict_1 = oad_dict['Moiety-1']
        #endregion
        #region Detail setting
        positions_list = []
        # peaks_dict_1, peaks_dict_2, peaks_dict_3 = {}, {}, {}
        ref_mz_1, ref_ratio_1, ref_mz_2, ref_ratio_2, ref_mz_3, ref_ratio_3 \
        = [], [], [], [], [], []
        # measured_oad_ions_1, measured_oad_ions_2, measured_oad_ions_3 = [], [], []
        # oad04_ions_1, oad04_ions_2, oad04_ions_3 = [], [], []
        # oad10_ions_1, oad10_ions_2, oad10_ions_3 = [], [], []
        is_this_plasmalogen = 'Plasm' in lipid_info['Ontology']
        # ref_precursor_mz = lipid_info['Ref precursor Mz']
        if dict_3:
            positions_list.append(dict_3[0]['Positions'])
            # peaks_dict_3 = dict_3[0]['Peaks dict']
        elif dict_2:
            positions_list.append(dict_2[0]['Positions'])
            # peaks_dict_2 = dict_2[0]['Peaks dict']
        elif dict_1:
            positions_list.append(dict_1[0]['Positions'])
            # peaks_dict_1 = dict_1[0]['Peaks dict']
        if dict_3:
            ref_mz_3 = [li[1] for li in dict_3['Determined db']['Ref peaks']]
            ref_ratio_3 = [
                li[3]*(-1) for li in dict_3['Determined db']['Ref peaks']]
            # measured_oad_ions_3 = [li[1] for li in dict_3['Determined db']['Measured peaks']]
            # oad04_ions_3 = [[li[1], li[3]] for li in dict_3['Determined db']['Ref peaks'] if 'OAD04' in li[0]]
            # oad10_ions_3 = [[li[1], li[3]] for li in dict_3['Determined db']['Ref peaks'] if 'OAD10' in li[0]]
        if dict_2:
            ref_mz_2 = [li[1] for li in dict_2['Determined db']['Ref peaks']]
            ref_ratio_2 = [
                li[3]*(-1) for li in dict_2['Determined db']['Ref peaks']]
            # measured_oad_ions_2 = [li[1] for li in dict_2['Determined db']['Measured peaks']]
            # oad04_ions_2 = [[li[1], li[3]] for li in dict_2['Determined db']['Ref peaks'] if 'OAD04' in li[0]]
            # oad10_ions_2 = [[li[1], li[3]] for li in dict_2['Determined db']['Ref peaks'] if 'OAD10' in li[0]]
        if dict_1:
            ref_mz_1 = [li[1] for li in dict_1['Determined db']['Ref peaks']]
            ref_ratio_1 = [
                li[3]*(-1) for li in dict_1['Determined db']['Ref peaks']]
            # measured_oad_ions_1 = [li[1] for li in dict_1['Determined db']['Measured peaks']]
            # oad04_ions_1 = [[li[1], li[3]] for li in dict_1['Determined db']['Ref peaks'] if 'OAD04' in li[0]]
            # oad10_ions_1 = [[li[1], li[3]] for li in dict_1['Determined db']['Ref peaks'] if 'OAD10' in li[0]]
        #endregion
        #region Range setting
        precursor_mz = graph_dict['MS2 Mz']
        ref_precursor_mz = graph_dict['Ref precursor Mz']
        x_min, x_max = graph_dict['x-range'][0], graph_dict['x-range'][1]
        bar_width = graph_dict['Bar_width']
        y_min, y_max = -100, 100
        #endregion
        #region magnification setting
        final_magnification = graph_dict['Magnification']
        #endregion
        #region matplotlib base setting
        fig = Figure(figsize=(fig_x, fig_y), dpi=dpi_setting)
        fig.patch.set_facecolor('white')
        fig.patch.set_alpha(0)
        ax = fig.add_subplot(1,1,1)
        ax.patch.set_facecolor('white')
        ax.patch.set_alpha(0)
        ax.set_xlim([x_min, x_max])
        ax.set_ylim([y_min, y_max])
        for axis, width in zip(axes, axes_width):
            ax.spines[axis].set_linewidth(width)
        ax.tick_params(labelsize=small_font_size, width=tick_width, pad=tick_pad)
        ax.set_yticks([-100, -50, 0, 50, 100])
        ax.set_yticklabels(['100', '50', '0', '50', '100'])
        #endregion
        #region Plotting OAD-MS/MS
        ex_msms = msms_df[msms_df['frag m/z'] >= x_min]
        total_fig_x = ex_msms['frag m/z'].values.tolist()
        total_fig_y = ex_msms['Ratio(%)'].values.tolist()
        total_fig_y = [v*final_magnification for v in total_fig_y]
        ref_ratio_1 = [v*final_magnification for v in ref_ratio_1]
        ref_ratio_2 = [v*final_magnification for v in ref_ratio_2]
        ref_ratio_3 = [v*final_magnification for v in ref_ratio_3]
        ax.bar(total_fig_x, total_fig_y, color=line_color, width=bar_width, 
            linewidth=0)
        ax.text(precursor_mz, 101, str(precursor_mz), 
            fontsize=precursor_mz_font_size, ha='center', va='bottom')
        ax.bar(ref_precursor_mz, -99, color=moiety_1_color, width=bar_width, 
            linewidth=0)
        ax.text(ref_precursor_mz, -137, str(ref_precursor_mz), 
            fontsize=precursor_mz_font_size, ha='center', va='bottom', 
            color=moiety_1_color)
        if ref_mz_1:
            ax.bar(ref_mz_1, ref_ratio_1, color=moiety_1_color, 
                width=bar_width, linewidth=0)
        if ref_mz_2:
            ax.bar(ref_mz_2, ref_ratio_2, color=moiety_2_color, 
                width=bar_width, linewidth=0)
        if ref_mz_3:
            ax.bar(ref_mz_3, ref_ratio_3, color=moiety_3_color, 
                width=bar_width, linewidth=0)
        #endregion
        #region matplotlib regend setting
        magnified_mz = round(precursor_mz-10)
        ax.hlines(y=100, xmin=(x_min+2), xmax=magnified_mz, colors='gray', 
                  linestyle='dashed')
        ax.hlines(y=0, xmin=x_min, xmax=x_max, linewidths=0.5)
        hlines_label = '×{}'.format(final_magnification)
        label_position_x = statistics.mean([x_min+5, magnified_mz])
        ax.text(label_position_x, 102, hlines_label, fontsize=small_font_size)
        msms_label_measured = '$\it{Measurement}$'
        msms_label_ref = '$\it{Reference}$'
        ax.text(x_min+1, 102, msms_label_measured, fontsize=small_font_size)
        ax.text(x_min+1, -98, msms_label_ref, fontsize=small_font_size)
        y_label = '$\it{Relative}$' + ' ' + '$\it{abundance}$'
        ax.set_xlabel('$\it{m/z}$', fontsize=x_label_font_size, labelpad=tick_pad)
        ax.set_ylabel(y_label, fontsize=y_label_font_size, labelpad=tick_pad)
        fig.tight_layout()
        #endregion
        #region saving setting
        save_path = f'{path}/{stamp}_tempfiles/ID{idx}.png'
        if os.path.isdir(f'{path}/{stamp}_tempfiles'): pass
        else: os.mkdir(f'{path}/{stamp}_tempfiles')
        fig.savefig(save_path, dpi=dpi_setting)
        each_rep.set(
            f'Generating {i+1}/{total} >>> {time.time()-start:.3f} [sec]')
        each_bar.step(1)
        fig.clear()
        # plt.close()
        #endregion

def save_msms_fig_as_pptx(path, stamp, target_table, each_bar, each_rep):
    tempfiles_path = f'{path}/{stamp}_tempfiles'
    ppt = pptx.Presentation()
    ppt_width = ppt.slide_width
    ppt_height = ppt.slide_height
    all_images = glob(f'{tempfiles_path}/*.png')
    blank_slide_layout = ppt.slide_layouts[6]
    txt_left = txt_top = txt_width = txt_height = Inches(1)*0.2
    total = len(all_images)
    for i, each_images in enumerate(all_images):
        start = time.time()
        matched_id = re.search(r'ID\d+\.', each_images)
        extracted_id = int(matched_id.group().replace('ID','').replace('.',''))
        extracted_index = target_table[target_table['ID'] == extracted_id].index
        oad_result_name = target_table['OAD result name'][extracted_index[0]]
        if oad_result_name == '':
            oad_result_name = target_table['Metabolite name'][extracted_index[0]]
        slide = ppt.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(
            each_images, 0, 0, ppt_width, ppt_width*9/16)
        pic.left = int((ppt_width - pic.width)/2)
        pic.top = int((ppt_height - pic.height)/2)
        txt_box = slide.shapes.add_textbox(
            txt_left, txt_top, txt_width*30, txt_height*2)
        txt_frame = txt_box.text_frame
        txt_frame.text = 'ID' + str(extracted_id) + ' ' + oad_result_name
        each_rep.set(
            f'Exporting {i+1}/{total} >>> {time.time()-start:.3f} [sec]')
        each_bar.step(1)
    pptx_path = f'{path}/{stamp}_MSRIDD_spectrum.pptx'
    ppt.save(pptx_path)
    shutil.rmtree(tempfiles_path)
#endregion

#region HeatMap of Neutral Loss
def generate_heatmap_of_nl(path, target_table, msms, stamp):
    nl_range, digit = 300, 0.01
    idx = [math_floor(digit*i, 2) for i in range(nl_range*100+1)]
    acyl_nl_df, acyl_ratio_df = pd.DataFrame(index=idx), pd.DataFrame(index=idx)
    sph_nl_df, sph_ratio_df = pd.DataFrame(index=idx), pd.DataFrame(index=idx)
    spl_nl_df, spl_ratio_df = pd.DataFrame(index=idx), pd.DataFrame(index=idx)
    #region generate NL&Ratio list
    for i, (row, df) in enumerate(target_table.iterrows()):
        nls = [math_floor(digit*i, 2) for i in range(nl_range*100+1)]
        ratios = [-1 for _ in range(len(nls))]
        table_id = df['ID']
        msms_df = msms[table_id]
        name, oad_name = df['Metabolite name'], df['OAD result name']
        res = exclude_null(nls, ratios, msms_df, label=name)
        nls, ratios = res[0], res[1]
        if '|' in name: rep_name = name.split('|')[-1]
        else: rep_name = name
        # col_name = f'{i}_{rep_name}'
        if re.search(r"\(d\d\)", name):
            spl_nl_df[f'{i}_{rep_name}'] = nls
            spl_ratio_df[f'{i}_{rep_name}'] = ratios
        elif re.search(r"\d+\:\d\;\dO\/\d+\:\d", name):
            sph_nl_df[f'{i}_{rep_name}'] = nls
            sph_ratio_df[f'{i}_{rep_name}'] = ratios
        else:
            acyl_nl_df[f'{i}_{rep_name}'] = nls
            acyl_ratio_df[f'{i}_{rep_name}'] = ratios
    #endregion
    #region Dropna
    acyl_nl_df = acyl_nl_df[acyl_nl_df > 0].dropna(how='all')
    acyl_ratio_df = acyl_ratio_df[acyl_ratio_df > 0].dropna(how='all')
    sph_nl_df = sph_nl_df[sph_nl_df > 0].dropna(how='all')
    sph_ratio_df = sph_ratio_df[sph_ratio_df > 0].dropna(how='all')
    spl_nl_df = spl_nl_df[spl_nl_df > 0].dropna(how='all')
    spl_ratio_df = spl_ratio_df[spl_ratio_df > 0].dropna(how='all')
    #endregion
    #region Excel generating&coloring
    wb = openpyxl.Workbook()
    res = generate_excelsheet(
        wb, acyl_nl_df, title='Acyls_NL', label='Excel generating')
    wb, acyls_nl_sht = res[0], res[1]
    acyls_nl_sht = coloring_by_ratio(
        acyls_nl_sht, acyl_ratio_df, label='Coloring cells')
    
    res = generate_excelsheet(
        wb, acyl_ratio_df, title='Acyls_Ratio', label='Excel generating')
    wb, acyls_ratio_sht = res[0], res[1]
    acyls_ratio_sht = coloring_by_ratio(
        acyls_ratio_sht, acyl_ratio_df, label='Coloring cells')
    
    res = generate_excelsheet(
        wb, sph_nl_df, title='Sphingo_NL', label='Excel generating')
    wb, sphingo_nl_sht = res[0], res[1]
    sphingo_nl_sht = coloring_by_ratio(
        sphingo_nl_sht, sph_ratio_df, label='Coloring cells')

    res = generate_excelsheet(
        wb, sph_ratio_df, title='Sphingo_Ratio', label='Excel generating')
    wb, sphingo_ratio_sht = res[0], res[1]
    sphingo_ratio_sht = coloring_by_ratio(
        sphingo_ratio_sht, sph_ratio_df, label='Coloring cells')

    res = generate_excelsheet(
        wb, spl_nl_df, title='SPLASH_NL', label='Excel generating')
    wb, splash_nl_sht = res[0], res[1]
    splash_nl_sht = coloring_by_ratio(
        splash_nl_sht, spl_ratio_df, label='Coloring cells')

    res = generate_excelsheet(
        wb, spl_ratio_df, title='SPLASH_Ratio', label='Excel generating')
    wb, splash_ratio_sht = res[0], res[1]
    splash_ratio_sht = coloring_by_ratio(
        splash_ratio_sht, spl_ratio_df, label='Coloring cells')
    #endregion
    del wb['Sheet']
    save_path = f'{path}/{stamp}_MSRIDD_HeatMap_of_NLs.xlsx'
    wb.save(filename=save_path)

def process_timer(func):
    @wraps(func)
    def _process_timer(*args, **kwargs):
        start = time.time()
        label = kwargs['label']
        print(f'{label} -> start')
        res = func(*args, **kwargs)
        req = time.time() - start
        req = math_floor(req, 2)
        print(f'{label} -> end (Required time: {req} [sec]')
        return res
    return _process_timer

# @process_timer
def exclude_null(nls, ratios, msms_df, **kwargs):
    for j, nl in enumerate(nls):
        ex_df = msms_df[(msms_df['Delta']>=nl)&(msms_df['Delta']<nl+0.01)]
        if ex_df.empty:
            nls[j] = -1
        else:
            nls[j] = ex_df['Delta'].values[0]
            ratios[j] = ex_df['Ratio(%)'].values[0]
    return nls, ratios

# @process_timer
def generate_excelsheet(wb, df, title='', **kwargs):
    sheet = wb.create_sheet(title=title)
    for row in dataframe_to_rows(df, index=True, header=True):
        sheet.append(row)
    sheet.delete_rows(2)
    sheet.freeze_panes = 'A2'
    return wb, sheet

# @process_timer
def coloring_by_ratio(sheet, ratio_df, **kwargs):
    arr = ratio_df.values
    rows, cols = arr.shape
    for row in range(rows):
        for col in range(cols):
            ratio = arr[row][col]
            if ratio < 0.005:
                # background_color = lightcyan
                sheet.cell(row=row+2, column=col+2).fill \
                = PatternFill(patternType='solid', fgColor="E0FFFF")
            elif (0.005 <= ratio < 0.1):
                # background_color = mistyrose
                sheet.cell(row=row+2, column=col+2).fill \
                = PatternFill(patternType='solid', fgColor="FBEFEF")
            elif  (0.1 <= ratio < 1.0):
                # background_color = generated by color picker
                sheet.cell(row=row+2, column=col+2).fill \
                = PatternFill(patternType='solid', fgColor="F6CECE")
                # sheet.cell(row=row+2, column=col+2).font = Font(b=True)
            elif (1.0 <= ratio < 10.0):
                sheet.cell(row=row+2, column=col+2).fill \
                = PatternFill(patternType='solid', fgColor="F78181")
            elif (10.0 <= ratio < 50.0):
                # background_color = red
                sheet.cell(row=row+2, column=col+2).fill \
                = PatternFill(patternType='solid', fgColor="FF0000")
                sheet.cell(row=row+2, column=col+2).font \
                = Font(color='FFFFFF', b=True)
            elif (50.0 <= ratio):
                # background_color = black
                sheet.cell(row=row+2, column=col+2).fill \
                = PatternFill(patternType='solid', fgColor="000000")
                sheet.cell(row=row+2, column=col+2).font \
                = Font(color='FFFFFF', b=True)
    return sheet
#endregion

